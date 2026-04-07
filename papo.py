import streamlit as st
import streamlit.components.v1 as components
import plotly.graph_objects as go
import pandas as pd
import base64 as _b64

try:
    import os as _os
    _BASE = _os.path.dirname(_os.path.abspath(__file__))
    _LOGO_B64 = _b64.b64encode(open(_os.path.join(_BASE, "Yuno logo.png"),"rb").read()).decode()
except:
    _LOGO_B64 = ""

try:
    _DANIELA_B64 = _b64.b64encode(open(_os.path.join(_BASE, "daniela.png"),"rb").read()).decode()
except:
    _DANIELA_B64 = ""

try:
    import sys as _sys
    _sys.path.insert(0, _os.path.join(_BASE, "data"))
    from partner_contacts import PARTNER_CONTACTS as _EXT_CONTACTS
except:
    _EXT_CONTACTS = {}

# ── Page Config ────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="PAPO · Yuno Partner Portal",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── Session State ──────────────────────────────────────────────────────────────
for k, v in [("role", None), ("pending_role", None), ("pw_error", False), ("page", "Home"), ("cat_filter", "all"), ("insight_tab", "market"), ("selected_partner", None), ("partner_tab", "Overview")]:
    if k not in st.session_state:
        st.session_state[k] = v

# Sync page from URL query params on first load
_qp = st.query_params
if "role" in _qp and st.session_state.role is None:
    if _qp["role"] in ("internal", "partner"):
        st.session_state.role = _qp["role"]
if "page" in _qp:
    valid_pages = {"Home","Pipeline","Partners","Merchants","Performance","Benchmarks","Insights","PayRec","MerchSim","MissionCtrl"}
    if _qp["page"] in valid_pages:
        st.session_state.page = _qp["page"]

# ── Data (loaded from Excel) ──────────────────────────────────────────────────
_PARTNERS_EXCEL = _os.path.join(_os.path.dirname(_os.path.abspath(__file__)), "data", "strategic_accounts.xlsx")

_TYPE_COLOR = {"Acquirer": "#3b82f6", "PSP/Aggregator": "#8b5cf6", "APM": "#14b8a6", "Fraud Provider": "#ef4444", "BaaS": "#f59e0b", "Other": "#64748b", "Plug-In": "#06b6d4"}
_TYPE_SHORT = {"Acquirer": "Acquirer", "PSP/Aggregator": "PSP", "APM": "APM", "Fraud Provider": "Fraud", "BaaS": "BaaS", "Other": "Other", "Plug-In": "Plug-In"}
_STAGE_MAP  = {"Live Partner": "Live", "Agreement Signed": "Agreement Signed", "Agreement Review": "Agreement Review", "Initial Negotiation": "Negotiation", "Opportunity Identification": "Prospect", "Lost": "Lost", "Only to be integrated": "Integration Only", "Agreement Signed - Only referrals": "Referral Only", "Non-qualified Partner": "Non-Qualified"}

@st.cache_data
def load_partners_excel():
    try:
        df = pd.read_excel(_PARTNERS_EXCEL, sheet_name="All partners")
        seen, partners = set(), []
        for _, row in df.iterrows():
            name = str(row.get("Provider", "")).strip()
            if not name or name in seen or name == "nan":
                continue
            seen.add(name)
            offering = str(row.get("Partner offering (acquirer, APM, Fraud etc)", "Other")).strip()
            stage_raw = str(row.get("Deal Stage", "")).strip()
            region = str(row.get("Region", "")).strip()
            country = str(row.get("Country", "")).strip()
            tier = str(row.get("Type of Partner", "")).strip()
            manager = str(row.get("Partner Manager", "")).strip()
            strategic = bool(row.get("Strategic?", False))
            mgmt_type = str(row.get("Type of Management", "")).strip()
            initials = "".join(w[0] for w in name.replace("/", " ").replace("(", " ").split() if w)[:2].upper()
            partners.append({
                "name": name, "type": _TYPE_SHORT.get(offering, offering), "offering_raw": offering,
                "region": region if region != "nan" else "", "country": country if country != "nan" else "",
                "status": _STAGE_MAP.get(stage_raw, stage_raw), "stage_raw": stage_raw,
                "tier": tier if tier != "nan" else "", "manager": manager if manager != "nan" else "",
                "strategic": strategic, "mgmt_type": mgmt_type if mgmt_type != "nan" else "",
                "logo": initials, "color": _TYPE_COLOR.get(offering, "#64748b"),
                "cat": _TYPE_SHORT.get(offering, offering),
                "nda": bool(row.get("NDA Signed and in drive", False)),
                "revshare": bool(row.get("Revshare Contract", False)),
                "revshare_active": bool(row.get("Revshare active", False)),
                "integration_ready": bool(row.get("Integration Ready by Yuno", False)),
                "integration_used": bool(row.get("Integration Used by Merchants", False)),
                "comments": str(row.get("Comments", "")).strip() if str(row.get("Comments", "")).strip() != "nan" else "",
            })
        return sorted(partners, key=lambda x: x["name"].lower())
    except Exception:
        return []

PARTNERS_DATA = load_partners_excel()

# Revenue share data from PPTX (Feb 2026)
_REVSHARE_BY_PARTNER = {"DLocal": 25338, "Stripe": 21705, "Cielo": 17635, "Bamboo": 9491, "Mercado Pago": 7365, "Unlimint": 5949, "PagBank": 4516, "Nuvei/Paymentez": 3772, "Pagar.me": 3441, "PicPay": 3180}
_REVSHARE_MONTHLY = [("Nov 24",61000),("Dec 24",61000),("Jan 25",122000),("Feb 25",122000),("Mar 25",183000),("Apr 25",183000),("May 25",144000),("Jun 25",145000),("Jul 25",166000),("Aug 25",219000),("Sep 25",245000),("Oct 25",245000),("Nov 25",178000),("Dec 25",161000),("Jan 26",123000),("Feb 26",123000)]
_REVSHARE_BY_REGION = {"LATAM": 42, "Brazil": 29, "Global": 25, "APAC": 3, "North America": 0}
_REGION_STATS = {
    "Brazil":        {"total":73,"live":32,"strategic":9,"tier1":4,"revshare":"$32.9K/mo"},
    "LATAM":         {"total":105,"live":15,"strategic":2,"tier1":31,"revshare":"$47.6K/mo"},
    "EMEA":          {"total":81,"live":2,"strategic":3,"tier1":21,"revshare":"-"},
    "Global":        {"total":107,"live":18,"strategic":13,"tier1":12,"revshare":"$28.3K/mo"},
    "APAC":          {"total":53,"live":6,"strategic":4,"tier1":10,"revshare":"$32.9K/mo"},
    "North America": {"total":37,"live":3,"strategic":5,"tier1":4,"revshare":"$0.4K/mo"},
    "Africa":        {"total":5,"live":0,"strategic":0,"tier1":2,"revshare":"-"},
}
_MANAGER_STATS = {
    "Alessandra Rospigliosi": {"regions":"EMEA, APAC, Africa, Global","strategic":8,"tier1":17,"tier2":44,"tier3":16},
    "Francisco Quintana":     {"regions":"MENAT, APAC, Global","strategic":1,"tier1":12,"tier2":25,"tier3":15},
    "Lily":                   {"regions":"APAC, Global","strategic":3,"tier1":9,"tier2":19,"tier3":1},
    "Johanderson Guevara":    {"regions":"Brazil, LATAM, Global","strategic":8,"tier1":13,"tier2":12,"tier3":22},
    "Talita Diaz Gama":       {"regions":"LATAM, Brazil, Global","strategic":3,"tier1":20,"tier2":16,"tier3":24},
    "Partha":                 {"regions":"APAC","strategic":1,"tier1":0,"tier2":1,"tier3":1},
}

CAT_CLASS   = {"PSP":"cat-psp","Acquirer":"cat-acquirer","APM":"cat-apm","Fraud":"cat-fraud","Fraud / KYC":"cat-fraud","BaaS":"cat-baas"}
STATUS_CLASS = {"Live":"p-green","Integration":"p-blue","Prospect":"p-amber","In Dev":"p-purple"}

# ── Source of Truth data for Payment Recommendations ─────────────────────────
_SOT_PATH = _os.path.join(_os.path.dirname(_os.path.abspath(__file__)), "data", "source_of_truth.xlsx")

_VERTICAL_COLS = {
    "High Risk":     "ACCEPTS_HIGH_RISK",
    "Gambling":      "ACCEPTS_GAMBLING",
    "Gaming":        "ACCEPTS_GAMING",
    "Forex":         "ACCEPTS_FOREX",
    "Crypto":        "ACCEPTS_CRYPTO",
    "Adult":         "ACCEPTS_ADULT",
    "MLM":           "ACCEPTS_MULTI_LEVEL_MARKETING",
    "Airlines":      "ACCEPTS_AIRLINES",
}

_ISO_TO_COUNTRY = {}
try:
    import pycountry as _pyc
    for c in _pyc.countries:
        _ISO_TO_COUNTRY[c.alpha_2] = c.name
except Exception:
    pass

@st.cache_data
def load_sot_data():
    try:
        df = pd.read_excel(_SOT_PATH, sheet_name="Partners")
        df = df[df["PROVIDER_CATEGORY"].isin(["ACQUIRER","GATEWAY","AGREGATOR","AGREGATOR / GATEWAY","PAYMENT_METHOD"])].copy()
        return df
    except Exception:
        return pd.DataFrame()

_SOT_DF = load_sot_data()
_SOT_COUNTRIES = sorted(set(_ISO_TO_COUNTRY.get(c, c) for c in _SOT_DF["COUNTRY_ISO"].dropna().unique() if _ISO_TO_COUNTRY.get(c, c) != c)) if len(_SOT_DF) > 0 else []
_SOT_PROVIDERS = sorted(_SOT_DF["PROVIDER_NAME"].dropna().unique().tolist()) if len(_SOT_DF) > 0 else []

def find_partners(country_iso=None, verticals=None, live_only=True, processing_type=None):
    """Find matching partners from Source of Truth."""
    df = _SOT_DF.copy()
    if len(df) == 0:
        return []
    if live_only:
        df = df[df["Live/NonLive Partner/Contract signed"] == "Live"]
    if country_iso:
        df = df[df["COUNTRY_ISO"] == country_iso]
    if processing_type:
        df = df[df["PROCESSING_TYPE"] == processing_type]
    if verticals:
        for v in verticals:
            col = _VERTICAL_COLS.get(v)
            if col and col in df.columns:
                df = df[(df[col] == True) | (df[col] == 1) | (df[col] == 1.0)]
    # Group by provider
    results = []
    for provider, grp in df.groupby("PROVIDER_NAME"):
        cats = grp["PROVIDER_CATEGORY"].unique().tolist()
        countries = sorted(grp["COUNTRY_ISO"].dropna().unique().tolist())
        pm_types = sorted(grp["PAYMENT_METHOD_TYPE"].dropna().unique().tolist())
        proc_types = grp["PROCESSING_TYPE"].dropna().unique().tolist()
        status = "Live" if "Live" in grp["Live/NonLive Partner/Contract signed"].values else "Non Live"
        supports = {}
        for feat in ["SUPPORTS_TOKENIZATION","SUPPORTS_RECURRING_PAYMENTS","SUPPORTS_PAYOUTS","SUPPORTS_INSTALLMENTS","SUPPORTS_PAYFAC","SUPPORTS_SPLIT_PAYMENTS","3DS"]:
            vals = grp[feat].dropna().unique()
            supports[feat] = any(v == True or v == 1 or v == 1.0 for v in vals)
        results.append({
            "name": provider,
            "categories": cats,
            "countries_iso": countries,
            "countries_count": len(countries),
            "payment_methods": pm_types[:8],
            "processing_types": proc_types,
            "status": status,
            "row_count": len(grp),
            "supports": supports,
        })
    results.sort(key=lambda x: (-x["countries_count"], x["name"]))
    return results

PIPELINE_STAGES = {
    "Prospect":    {"color":"#86868b","count":3},
    "Qualified":   {"color":"#60a5fa","count":3},
    "Evaluation":  {"color":"#c084fc","count":4},
    "Negotiation": {"color":"#fbbf24","count":2},
    "Won":         {"color":"#4ade80","count":2},
}

CONTACTS = [
    {"init":"TK","bg":"rgba(59,130,246,.2)","color":"#60a5fa","name":"Tom Kuehn","role":"Head of LATAM Partnerships","company":"Adyen · PSP","badge":"Champion","badge_class":"p-green","last":"2d ago","rel":"★★★★★","deals":"3 active"},
    {"init":"RV","bg":"rgba(168,85,247,.2)","color":"#c084fc","name":"Ricardo Vega","role":"VP Business Development","company":"Nuvei · PSP","badge":"Warm","badge_class":"p-amber","last":"5h ago","rel":"★★★★☆","deals":"1 active"},
    {"init":"AP","bg":"rgba(20,184,166,.2)","color":"#2dd4bf","name":"Ana Pacheco","role":"Strategic Partnerships Director","company":"Kushki · Acquirer","badge":"Active","badge_class":"p-blue","last":"1w ago","rel":"★★★★☆","deals":"2 active"},
    {"init":"FM","bg":"rgba(245,158,11,.2)","color":"#fbbf24","name":"Felipe Morales","role":"Head of Digital Products","company":"Getnet · Acquirer (Santander)","badge":"Executive Sponsor","badge_class":"p-green","last":"Today","rel":"★★★★★","deals":"1 at Negotiation"},
    {"init":"DH","bg":"rgba(239,68,68,.2)","color":"#fca5a5","name":"Diana Herrera","role":"Fraud Partnerships Lead","company":"SEON · Fraud & Risk","badge":"Live Partner","badge_class":"p-green","last":"3d ago","rel":"★★★★☆","deals":"Integration live"},
    {"init":"MC","bg":"rgba(245,158,11,.2)","color":"#fde68a","name":"Martín Castillo","role":"CEO & Co-Founder","company":"Pomelo · BaaS","badge":"New Vertical","badge_class":"p-amber","last":"Yesterday","rel":"★★★☆☆","deals":"1 at Evaluation"},
]

# ── CSS ─────────────────────────────────────────────────────────────────────────
_STATIC_CSS = """
@import url('https://fonts.googleapis.com/css2?family=Titillium+Web:wght@300;400;600;700&display=swap');
:root {
  --bg:#f5f5f7; --bg2:#ffffff; --bg3:#f0f0f5; --bg4:#e8e8ed; --bg5:#d1d1d6;
  --border:rgba(0,0,0,0.07); --border2:rgba(0,0,0,0.13);
  --text:#1d1d1f; --text2:#6e6e73; --text3:#86868b;
  --indigo:#4F46E5; --indigo-light:rgba(79,70,229,0.08); --indigo-mid:rgba(79,70,229,0.15);
  --blue:#2563eb; --green:#16a34a; --red:#dc2626; --amber:#d97706;
  --purple:#7c3aed; --teal:#0d9488;
  --shadow:0 1px 3px rgba(0,0,0,0.07),0 0 0 1px rgba(0,0,0,0.04);
  --shadow-md:0 4px 16px rgba(0,0,0,0.08),0 0 0 1px rgba(0,0,0,0.04);
  --font:'Titillium Web',sans-serif; --mono:Menlo,'Courier New',monospace;
}
html,body,[class*="css"] {
  font-family:var(--font) !important;
  background-color:var(--bg) !important;
  color:var(--text) !important;
}
#MainMenu,footer,header { visibility:hidden; }
.block-container { padding-top:2.5rem !important; padding-bottom:2.5rem !important; padding-left:2.5rem !important; padding-right:2.5rem !important; max-width:100% !important; background:var(--bg) !important; }
[data-testid="stAppViewContainer"] { background:var(--bg) !important; }
[data-testid="stAppViewContainer"] > section:nth-child(2) { background:var(--bg) !important; }
[data-testid="stMain"] { background:var(--bg) !important; }
.main { background:var(--bg) !important; }
section[data-testid="stMain"] > div { background:var(--bg) !important; }
[data-testid="stSidebar"] { background:#0A0F1E !important; border-right:0.5px solid rgba(255,255,255,0.05) !important; }
[data-testid="stSidebar"] > div:first-child { padding:0 !important; }
/* Role toggle buttons (inside columns in sidebar) */
[data-testid="stSidebar"] [data-testid="column"] [data-testid="stButton"] > button {
  border-radius:20px !important; border:none !important; font-size:11.5px !important; font-weight:600 !important; padding:7px 12px !important;
}
[data-testid="stSidebar"] [data-testid="column"] [data-testid="stButton"] > button[data-testid="baseButton-primary"] {
  background:var(--indigo) !important; color:#fff !important;
}
[data-testid="stSidebar"] [data-testid="column"] [data-testid="stButton"] > button[data-testid="baseButton-secondary"] {
  background:transparent !important; color:var(--text3) !important;
}
/* Nav buttons (sidebar buttons NOT in columns) */
[data-testid="stSidebar"] [data-testid="stVerticalBlock"] > div > [data-testid="stButton"] > button {
  background:transparent !important; border:none !important; color:var(--text2) !important;
  text-align:left !important; justify-content:flex-start !important; border-radius:8px !important;
  font-size:14px !important; font-weight:400 !important; padding:10px 14px !important;
  display:flex !important; align-items:center !important; width:100% !important;
}
[data-testid="stSidebar"] [data-testid="stVerticalBlock"] > div > [data-testid="stButton"] > button:hover {
  background:var(--bg3) !important; color:var(--text) !important;
}
[data-testid="stSidebar"] button[data-nav-active="true"] {
  background:var(--indigo-light) !important; border-left:3px solid var(--indigo) !important;
  border-radius:0 8px 8px 0 !important; color:var(--indigo) !important; font-weight:600 !important; padding-left:11px !important;
}
[data-testid="stSidebar"] button[data-badge]::after {
  content:attr(data-badge); margin-left:auto; font-size:10px; font-family:var(--mono);
  background:var(--bg3); color:var(--text3); padding:2px 7px; border-radius:10px; flex-shrink:0;
}
[data-testid="stSidebar"] button[data-nav-active="true"][data-badge]::after {
  background:var(--indigo-mid); color:var(--indigo);
}
[data-testid="stTextInput"] input { background:#fff !important; border:1px solid var(--border2) !important; color:var(--text) !important; border-radius:8px !important; font-family:var(--font) !important; font-size:15px !important; }
[data-testid="stTextInput"] input[type="password"] { background:#fff !important; color:var(--text) !important; -webkit-text-fill-color:var(--text) !important; caret-color:var(--text) !important; }
[data-testid="stTextInput"] input::placeholder { color:var(--text3) !important; }
[data-testid="stTextInput"] input:focus { border-color:var(--indigo) !important; box-shadow:0 0 0 3px rgba(79,70,229,0.12) !important; }
.stButton > button { background:#fff !important; color:var(--text2) !important; border:1px solid var(--border2) !important; border-radius:8px !important; font-family:var(--font) !important; font-weight:600 !important; font-size:14px !important; }
.stButton > button:hover { border-color:var(--indigo) !important; color:var(--indigo) !important; }
.stButton > button[data-testid="baseButton-primary"] { background:var(--indigo) !important; color:#fff !important; border:none !important; }
.stButton > button[data-testid="baseButton-primary"]:hover { background:#4338ca !important; }
.pill{display:inline-flex;align-items:center;font-size:10px;font-weight:600;padding:3px 9px;border-radius:20px;font-family:var(--mono);letter-spacing:.2px;white-space:nowrap;}
.p-green{background:rgba(22,163,74,.1);color:#15803d;}
.p-blue{background:rgba(37,99,235,.1);color:#1d4ed8;}
.p-amber{background:rgba(217,119,6,.1);color:#b45309;}
.p-red{background:rgba(220,38,38,.1);color:#b91c1c;}
.p-purple{background:rgba(124,58,237,.1);color:#6d28d9;}
.p-teal{background:rgba(13,148,136,.1);color:#0f766e;}
.p-grey{background:rgba(107,114,128,.1);color:#4b5563;}
.cat-acquirer{background:rgba(0,0,0,.04);color:var(--text3);border:1px solid rgba(0,0,0,.07);}
.cat-psp{background:rgba(0,0,0,.04);color:var(--text3);border:1px solid rgba(0,0,0,.07);}
.cat-apm{background:rgba(0,0,0,.04);color:var(--text3);border:1px solid rgba(0,0,0,.07);}
.cat-fraud{background:rgba(0,0,0,.04);color:var(--text3);border:1px solid rgba(0,0,0,.07);}
.cat-baas{background:rgba(0,0,0,.04);color:var(--text3);border:1px solid rgba(0,0,0,.07);}
::-webkit-scrollbar{width:5px;height:5px;}
::-webkit-scrollbar-track{background:transparent;}
::-webkit-scrollbar-thumb{background:var(--bg4);border-radius:3px;}
.pipeline-wrap{overflow-x:auto;margin-bottom:20px;padding-bottom:8px;}
.pipeline-board{display:flex;gap:12px;min-width:1000px;align-items:flex-start;}
.stage-col{background:#fff;border:none;box-shadow:var(--shadow);border-radius:12px;width:200px;flex-shrink:0;overflow:hidden;}
.stage-head{padding:10px 14px;background:#fafafa;border-bottom:1px solid rgba(0,0,0,0.06);display:flex;align-items:center;justify-content:space-between;}
.stage-name{font-size:9.5px;font-weight:700;letter-spacing:.8px;text-transform:uppercase;}
.stage-num{font-size:10px;font-family:var(--mono);background:var(--bg3);color:var(--text3);padding:2px 7px;border-radius:10px;}
.pcard{margin:8px;padding:11px 12px;background:#fff;border:1px solid rgba(0,0,0,0.06);border-radius:10px;cursor:pointer;transition:box-shadow .12s,border-color .12s;}
.pcard:hover{box-shadow:var(--shadow-md);border-color:rgba(0,0,0,0.1);}
.pcard-name{font-size:12.5px;font-weight:600;margin-bottom:5px;color:var(--text);}
.pcard-meta{display:flex;align-items:center;gap:5px;flex-wrap:wrap;margin-bottom:5px;}
.pcard-val{font-family:var(--mono);font-size:10.5px;color:var(--text2);}
.pcard-owner{display:flex;align-items:center;gap:5px;font-size:10px;color:var(--text3);margin-top:7px;}
.mini-av{width:16px;height:16px;border-radius:50%;font-size:7px;font-weight:700;display:inline-flex;align-items:center;justify-content:center;}
.ma-y{background:rgba(79,70,229,0.15);color:#4F46E5;}
.ma-e{background:rgba(37,99,235,0.15);color:#2563eb;}
.pbar{height:3px;background:var(--bg3);border-radius:2px;margin-top:9px;overflow:hidden;}
.pfill{height:100%;border-radius:2px;}
"""

_LANDING_CSS = """
[data-testid="stSidebar"] { display: none !important; }
[data-testid="collapsedControl"] { display: none !important; }
.stApp, section[data-testid="stMain"], .main, .block-container { background: #f5f5f7 !important; }
.main .block-container { max-width: 100% !important; padding: 0 !important; }
div[data-testid="column"] { padding: 0 6px !important; }
div[data-testid="stVerticalBlock"] { gap: 0 !important; }
div[data-testid="stVerticalBlockBorderWrapper"] { gap: 0 !important; }
div[data-testid="stButton"] { margin-top: -1px !important; }
div[data-testid="stButton"] button {
  border-radius: 0 0 12px 12px !important; border: none !important;
  font-size: 11px !important; font-weight: 600 !important;
  padding: 10px 16px !important; width: 100% !important;
  background: #4F46E5 !important; color: #fff !important;
}
div[data-testid="stButton"] button[data-partner-btn="true"] {
  background: #fff !important; color: #4F46E5 !important;
  border: 1.5px solid rgba(79,70,229,.3) !important;
}
div[data-testid="stButton"] button:hover { opacity: 0.88 !important; }
"""

def inject_css(role):
    css = _STATIC_CSS + ("" if role else _LANDING_CSS)
    escaped = css.replace('\\', '\\\\').replace('`', '\\`')
    if not role:
        btn_js = """
setInterval(function(){
  var doc = window.parent.document;
  var allBtns = doc.querySelectorAll('button[data-testid="baseButton-secondary"]');
  // Check if password modal is open (pw_back button exists)
  var hasPwBack = Array.from(allBtns).some(function(b){ return b.innerText.trim().startsWith('←'); });
  if(hasPwBack){
    allBtns.forEach(function(b){
      var txt = b.innerText.trim();
      if(txt.startsWith('←')){
        b.style.setProperty('background','rgba(255,255,255,.08)','important');
        b.style.setProperty('color','#aaa','important');
        b.style.setProperty('border','1px solid rgba(255,255,255,.12)','important');
      }
    });
  } else {
    if(allBtns[0]){allBtns[0].style.setProperty('background','rgba(255,255,255,.2)','important');allBtns[0].style.setProperty('color','#fff','important');}
    if(allBtns[1]){
      allBtns[1].setAttribute('data-partner-btn','true');
      allBtns[1].style.setProperty('background','#ffffff','important');
      allBtns[1].style.setProperty('color','#4F46E5','important');
      allBtns[1].style.setProperty('border','1.5px solid rgba(79,70,229,.3)','important');
    }
  }
},100);"""
    else:
        current_page = st.session_state.get('page', 'Pipeline')
        _page_to_nav = {"Pipeline":"Partner Leads","Partners":"Partner Portfolio","Merchants":"Merchants","MerchSim":"Merchant Sim","MissionCtrl":"Partners In Flight","Performance":"Partner Health","Benchmarks":"Rev Share","Insights":"Market Intel","Home":"Home"}
        active_nav = _page_to_nav.get(current_page, "Home")
        _badge_map = {"Partner Portfolio":"47","Partners In Flight":"3","Partner Leads":"24","Merchants":"12","Rev Share":"Apr 1"}
        btn_js = f"""
setTimeout(function(){{
  var activeNav='{active_nav}';
  var badges={str(_badge_map).replace("'",'"')};
  var sidebar=window.parent.document.querySelector('[data-testid="stSidebar"]');
  if(!sidebar)return;
  sidebar.querySelectorAll('button').forEach(function(btn){{
    var t=btn.textContent.trim();
    var isActive=(t===activeNav);
    var badge=badges[t]||'';
    // Style as dark sidebar nav item
    btn.style.cssText='display:flex !important;align-items:center !important;gap:9px !important;padding:8px 14px 8px 16px !important;cursor:pointer !important;border-radius:8px !important;margin:1px 8px !important;border:none !important;box-shadow:none !important;font-size:13px !important;font-weight:'+(isActive?'600':'400')+' !important;color:'+(isActive?'#A5B4FC':'#64748B')+' !important;background:'+(isActive?'rgba(91,95,222,0.15)':'transparent')+' !important;border-left:2px solid '+(isActive?'#5B5FDE':'transparent')+' !important;text-align:left !important;width:calc(100% - 16px) !important;';
    // Add badge
    if(badge && !btn.querySelector('.nav-badge')){{
      var b=document.createElement('span');
      b.className='nav-badge';
      b.style.cssText='margin-left:auto;font-size:10px;font-weight:600;padding:1px 6px;border-radius:8px;font-family:monospace;'+(badge==='3'?'background:rgba(239,68,68,0.2);color:#FCA5A5;':'background:rgba(255,255,255,0.07);color:#64748B;');
      b.textContent=badge;
      btn.appendChild(b);
    }}
  }});
}},300);"""
    components.html(
        f"<script>var old=window.parent.document.getElementById('papo-css');if(old)old.remove();var s=window.parent.document.createElement('style');s.id='papo-css';s.textContent=`{escaped}`;window.parent.document.head.appendChild(s);{btn_js}</script>",
        height=0,
    )

# ── Landing Page ───────────────────────────────────────────────────────────────
def show_landing():
    LOGO = _LOGO_B64
    pending = st.session_state.get("pending_role")

    # ── SSO / Email verification screen ─────────────────────────────────────
    if pending:
        sso_method = st.session_state.get("sso_method", "code")  # "google", "microsoft", or "code"

        if sso_method in ("google", "microsoft"):
            provider_name = "Google" if sso_method == "google" else "Microsoft"
            provider_color = "#4285F4" if sso_method == "google" else "#0078D4"
            provider_icon_svg = '<svg width="24" height="24" viewBox="0 0 18 18"><path d="M17.64 9.2c0-.637-.057-1.251-.164-1.84H9v3.481h4.844a4.14 4.14 0 01-1.796 2.716v2.259h2.908c1.702-1.567 2.684-3.875 2.684-6.615z" fill="#4285F4"/><path d="M9 18c2.43 0 4.467-.806 5.956-2.18l-2.908-2.259c-.806.54-1.837.86-3.048.86-2.344 0-4.328-1.584-5.036-3.711H.957v2.332A8.997 8.997 0 009 18z" fill="#34A853"/><path d="M3.964 10.71A5.41 5.41 0 013.682 9c0-.593.102-1.17.282-1.71V4.958H.957A8.997 8.997 0 000 9c0 1.452.348 2.827.957 4.042l3.007-2.332z" fill="#FBBC05"/><path d="M9 3.58c1.321 0 2.508.454 3.44 1.345l2.582-2.58C13.463.891 11.426 0 9 0A8.997 8.997 0 00.957 4.958L3.964 7.29C4.672 5.163 6.656 3.58 9 3.58z" fill="#EA4335"/></svg>' if sso_method == "google" else '<svg width="24" height="24" viewBox="0 0 21 21"><rect x="1" y="1" width="9" height="9" fill="#F25022"/><rect x="11" y="1" width="9" height="9" fill="#7FBA00"/><rect x="1" y="11" width="9" height="9" fill="#00A4EF"/><rect x="11" y="11" width="9" height="9" fill="#FFB900"/></svg>'

            # Loading redirect simulation
            if not st.session_state.get("sso_loaded"):
                import time
                components.html(f"""
<style>
@keyframes spin{{0%{{transform:rotate(0deg)}}100%{{transform:rotate(360deg)}}}}
@keyframes fadeIn{{0%{{opacity:0;transform:translateY(10px)}}100%{{opacity:1;transform:translateY(0)}}}}
</style>
<div style="display:flex;flex-direction:column;align-items:center;justify-content:center;padding:80px 20px;font-family:'Outfit',system-ui,sans-serif;">
  <div style="width:48px;height:48px;border-radius:50%;background:#fff;box-shadow:0 4px 24px rgba(0,0,0,0.1);display:flex;align-items:center;justify-content:center;margin-bottom:20px;">
    {provider_icon_svg}
  </div>
  <div style="font-size:15px;font-weight:600;color:#0f172a;margin-bottom:6px;">Redirecting to {provider_name}...</div>
  <div style="font-size:12px;color:#94a3b8;margin-bottom:24px;">Securely connecting to your organization</div>
  <div style="width:32px;height:32px;border:3px solid #e2e8f0;border-top:3px solid {provider_color};border-radius:50%;animation:spin 0.8s linear infinite;"></div>
  <div style="margin-top:24px;display:flex;align-items:center;gap:6px;animation:fadeIn 0.5s ease 0.3s both;">
    <svg width="12" height="12" viewBox="0 0 11 11" fill="none" stroke="#94a3b8" stroke-width="1.5" stroke-linecap="round"><path d="M5.5 1L2 2.5V6c0 2 1.5 3.5 3.5 4 2-0.5 3.5-2 3.5-4V2.5L5.5 1z"/></svg>
    <span style="font-size:10px;color:#94a3b8;">Protected by 256-bit TLS encryption</span>
  </div>
</div>
""", height=340)
                time.sleep(2)
                st.session_state["sso_loaded"] = True
                st.rerun()

            _, mc, _ = st.columns([1.2, 1.6, 1.2])
            with mc:
                st.markdown(f'<div style="text-align:center;padding:30px 0 12px;">'
                            f'<div style="width:48px;height:48px;border-radius:50%;background:#fff;box-shadow:0 2px 12px rgba(0,0,0,0.08);display:inline-flex;align-items:center;justify-content:center;margin-bottom:16px;">{provider_icon_svg}</div>'
                            f'<div style="font-size:18px;font-weight:700;color:#0f172a;margin-bottom:4px;">Verify your identity</div>'
                            f'<div style="font-size:12px;color:#94a3b8;margin-bottom:2px;">{provider_name} connected — confirm your email to continue</div>'
                            f'</div>', unsafe_allow_html=True)
                email = st.text_input("Email", placeholder="name@company.com", key="sso_email", label_visibility="collapsed")
                if st.session_state.get("pw_error"):
                    st.markdown('<p style="color:#ef4444;font-size:11px;text-align:center;margin:4px 0 0;">Please enter a valid work email address.</p>', unsafe_allow_html=True)
                c1, c2 = st.columns(2)
                with c1:
                    if st.button("← Back", key="sso_back", use_container_width=True):
                        st.session_state.pending_role = None
                        st.session_state.sso_method = None
                        st.session_state.sso_loaded = False
                        st.session_state.pw_error = False
                        st.rerun()
                with c2:
                    if st.button("Sign in →", key="sso_submit", use_container_width=True, type="primary"):
                        if email and "@" in email:
                            domain = email.split("@")[-1].lower()
                            role = "internal" if domain in ("yuno.com","yuno.co","y.uno") else "partner"
                            st.session_state.role = role
                            st.session_state.pending_role = None
                            st.session_state.pw_error = False
                            st.session_state.sso_method = None
                            st.session_state.sso_loaded = False
                            st.session_state.page = "Home"
                            st.query_params["role"] = role
                            st.query_params["page"] = "Home"
                            st.rerun()
                        else:
                            st.session_state.pw_error = True
                            st.rerun()
                st.markdown(f'<div style="text-align:center;margin-top:14px;font-size:10px;color:#94a3b8;line-height:1.6;">@yuno.com / @yuno.co / @y.uno → Yuno A-Team access<br>All other domains → Yuno Partner access</div>', unsafe_allow_html=True)
            return

        # Access code flow
        label = "Yuno A-Team" if pending == "internal" else "Yuno Partner"
        correct_pw = "Yuno" if pending == "internal" else "Partners"

        _, mc, _ = st.columns([1.2, 1.6, 1.2])
        with mc:
            st.markdown(f'<div style="text-align:center;padding:40px 0 16px;">'
                        f'<img src="data:image/png;base64,{LOGO}" style="height:40px;object-fit:contain;margin-bottom:20px;opacity:.5;">'
                        f'<div style="font-size:20px;font-weight:700;color:#1d1d1f;margin-bottom:6px;">Access {label}</div>'
                        f'<div style="font-size:12px;color:#6e6e73;margin-bottom:4px;">Enter the access code from your Yuno Partner Manager</div>'
                        f'</div>', unsafe_allow_html=True)
            pw = st.text_input("Code", type="password", key="pw_input",
                               placeholder="Enter access code…", label_visibility="collapsed")
            if st.session_state.get("pw_error"):
                st.markdown('<p style="color:#ef4444;font-size:11px;text-align:center;margin:4px 0 0;">Invalid access code — contact your Yuno Partner Manager.</p>', unsafe_allow_html=True)
            c1, c2 = st.columns(2)
            with c1:
                if st.button("← Back", key="pw_back", use_container_width=True):
                    st.session_state.pending_role = None
                    st.session_state.pw_error = False
                    st.rerun()
            with c2:
                if st.button("Enter →", key="pw_submit", use_container_width=True, type="primary"):
                    if pw == correct_pw:
                        st.session_state.role = pending
                        st.session_state.pending_role = None
                        st.session_state.pw_error = False
                        st.session_state.page = "Home"
                        st.query_params["role"] = pending
                        st.query_params["page"] = "Home"
                        st.rerun()
                    else:
                        st.session_state.pw_error = True
                        st.rerun()
        return

    # ── Portal sign-in screen (from yuno_portal_v3.html) ────────────────────
    components.html("""
<link href="https://fonts.googleapis.com/css2?family=Outfit:wght@300;400;500;600;700;800&display=swap" rel="stylesheet">
<style>
*{box-sizing:border-box;margin:0;padding:0}
:root{--ink:#0A0F1E;--indigo:#5B5FDE;--indigo2:#4F46E5;--font:'Outfit',system-ui,sans-serif;}
.login-wrap{display:flex;flex-direction:column;align-items:center;justify-content:center;min-height:100%;font-family:var(--font);}
.lm-lg{display:grid;grid-template-columns:repeat(4,7px);gap:2.5px;}
.lm-lg span{width:7px;height:7px;border-radius:1.5px;background:var(--indigo);}
.lm-lg .off{opacity:0}.lm-lg .mid{opacity:0.4}
.login-card{width:420px;background:rgba(255,255,255,0.03);border:0.5px solid rgba(255,255,255,0.08);border-radius:24px;padding:44px 40px;backdrop-filter:blur(12px);position:relative;z-index:1;}
.sso-btn{width:100%;padding:14px;border-radius:10px;font-size:14px;font-weight:600;font-family:var(--font);cursor:pointer;display:flex;align-items:center;justify-content:center;gap:10px;transition:all 0.15s;margin-bottom:10px;border:0.5px solid rgba(255,255,255,0.12);background:rgba(255,255,255,0.05);color:#E2E8F0;}
.sso-btn:hover{background:rgba(255,255,255,0.1);transform:translateY(-1px);box-shadow:0 8px 24px rgba(0,0,0,0.2);}
.sso-btn:active{transform:translateY(0);}
.sso-google{border-color:rgba(66,133,244,0.3);}
.sso-google:hover{border-color:rgba(66,133,244,0.5);background:rgba(66,133,244,0.08);}
.sso-microsoft{border-color:rgba(0,120,212,0.3);}
.sso-microsoft:hover{border-color:rgba(0,120,212,0.5);background:rgba(0,120,212,0.08);}
@keyframes pulse{0%,100%{opacity:1;transform:scale(1);}50%{opacity:0.6;transform:scale(1.3);}}
</style>
<div style="background:#0A0F1E;border-radius:20px;padding:60px 20px;position:relative;overflow:hidden;min-height:560px;display:flex;align-items:center;justify-content:center;">
  <div style="position:absolute;inset:0;background-image:radial-gradient(circle at 1px 1px,rgba(255,255,255,0.03) 1px,transparent 0);background-size:28px 28px;pointer-events:none;"></div>
  <div style="position:absolute;width:600px;height:600px;border-radius:50%;background:radial-gradient(circle,rgba(91,95,222,0.14) 0%,transparent 65%);top:50%;left:50%;transform:translate(-50%,-50%);pointer-events:none;"></div>
  <div class="login-wrap">
    <div class="login-card">
      <div style="display:flex;align-items:center;gap:10px;justify-content:center;margin-bottom:6px;">
        <div class="lm-lg" style="transform:scale(1.4);margin:4px 6px;">
          <span></span><span class="mid"></span><span class="mid"></span><span class="off"></span>
          <span class="mid"></span><span></span><span></span><span class="mid"></span>
          <span class="mid"></span><span></span><span></span><span class="mid"></span>
          <span class="off"></span><span class="mid"></span><span class="mid"></span><span></span>
        </div>
        <span style="font-size:28px;font-weight:800;color:#E2E8F0;letter-spacing:-0.03em;">yuno</span>
      </div>
      <div style="text-align:center;font-size:22px;font-weight:700;color:#F1F5F9;letter-spacing:-0.02em;margin:20px 0 8px;">Welcome to the<br>Yuno Partner Portal</div>
      <div style="text-align:center;font-size:13px;color:#475569;margin-bottom:32px;">Internal access only</div>
      <button class="sso-btn sso-google" onclick="ssoLogin('google')">
        <svg width="18" height="18" viewBox="0 0 18 18"><path d="M17.64 9.2c0-.637-.057-1.251-.164-1.84H9v3.481h4.844a4.14 4.14 0 01-1.796 2.716v2.259h2.908c1.702-1.567 2.684-3.875 2.684-6.615z" fill="#4285F4"/><path d="M9 18c2.43 0 4.467-.806 5.956-2.18l-2.908-2.259c-.806.54-1.837.86-3.048.86-2.344 0-4.328-1.584-5.036-3.711H.957v2.332A8.997 8.997 0 009 18z" fill="#34A853"/><path d="M3.964 10.71A5.41 5.41 0 013.682 9c0-.593.102-1.17.282-1.71V4.958H.957A8.997 8.997 0 000 9c0 1.452.348 2.827.957 4.042l3.007-2.332z" fill="#FBBC05"/><path d="M9 3.58c1.321 0 2.508.454 3.44 1.345l2.582-2.58C13.463.891 11.426 0 9 0A8.997 8.997 0 00.957 4.958L3.964 7.29C4.672 5.163 6.656 3.58 9 3.58z" fill="#EA4335"/></svg>
        Continue with Google
      </button>
      <button class="sso-btn sso-microsoft" onclick="ssoLogin('microsoft')">
        <svg width="18" height="18" viewBox="0 0 21 21"><rect x="1" y="1" width="9" height="9" fill="#F25022"/><rect x="11" y="1" width="9" height="9" fill="#7FBA00"/><rect x="1" y="11" width="9" height="9" fill="#00A4EF"/><rect x="11" y="11" width="9" height="9" fill="#FFB900"/></svg>
        Continue with Microsoft
      </button>
      <div style="display:flex;align-items:center;gap:10px;margin:20px 0;color:#1E293B;font-size:12px;">
        <div style="flex:1;height:0.5px;background:rgba(255,255,255,0.06);"></div>
        <span style="color:#334155;">or enter access code</span>
        <div style="flex:1;height:0.5px;background:rgba(255,255,255,0.06);"></div>
      </div>
      <div style="display:flex;gap:8px;">
        <button class="sso-btn" style="flex:1;margin:0;background:rgba(91,95,222,0.1);border-color:rgba(91,95,222,0.3);color:#A5B4FC;" onclick="enterCode('internal')">
          <svg width="16" height="16" viewBox="0 0 22 22" fill="none" stroke="#818CF8" stroke-width="1.5" stroke-linecap="round"><path d="M11 2L3 8v12h16V8L11 2z"/><path d="M8 20v-7h6v7"/></svg>
          Enter Portal
        </button>
      </div>
      <div style="text-align:center;margin-top:24px;font-size:11px;color:#1E293B;display:flex;align-items:center;justify-content:center;gap:5px;">
        <svg width="11" height="11" viewBox="0 0 11 11" fill="none" stroke="#2D3D55" stroke-width="1.5" stroke-linecap="round"><path d="M5.5 1L2 2.5V6c0 2 1.5 3.5 3.5 4 2-0.5 3.5-2 3.5-4V2.5L5.5 1z"/></svg>
        256-bit encrypted · SOC 2 Type II · SSO ready
      </div>
    </div>
  </div>
</div>
<script>
function ssoLogin(provider){
  // Store SSO method then trigger internal flow (email will determine role)
  var doc = window.parent.document;
  // Set hidden input to pass provider
  doc.querySelectorAll('button').forEach(function(btn){
    if(btn.textContent.trim() === 'SSO_'+provider) btn.click();
  });
}
function enterCode(role){
  var doc = window.parent.document;
  var target = role === 'internal' ? 'SSO Internal' : 'SSO Partner';
  doc.querySelectorAll('button').forEach(function(btn){
    if(btn.textContent.trim() === target) btn.click();
  });
}
</script>
""", height=620)

    # Hidden Streamlit buttons for all entry methods
    col_h1, col_h2, col_h3, col_h4 = st.columns(4)
    with col_h1:
        if st.button("SSO Internal", key="btn_internal", use_container_width=True):
            st.session_state.pending_role = "internal"
            st.session_state.sso_method = "code"
            st.session_state.pw_error = False
            st.rerun()
    with col_h2:
        if st.button("SSO Partner", key="btn_partner", use_container_width=True):
            st.session_state.pending_role = "partner"
            st.session_state.sso_method = "code"
            st.session_state.pw_error = False
            st.rerun()
    with col_h3:
        if st.button("SSO_google", key="btn_google", use_container_width=True):
            st.session_state.pending_role = "internal"
            st.session_state.sso_method = "google"
            st.session_state.pw_error = False
            st.rerun()
    with col_h4:
        if st.button("SSO_microsoft", key="btn_microsoft", use_container_width=True):
            st.session_state.pending_role = "internal"
            st.session_state.sso_method = "microsoft"
            st.session_state.pw_error = False
            st.rerun()

    components.html("""<script>
setTimeout(function(){
  var doc = window.parent.document;
  doc.querySelectorAll('button').forEach(function(btn){
    var t = btn.textContent.trim();
    if(t==='SSO Internal'||t==='SSO Partner'||t==='SSO_google'||t==='SSO_microsoft'){
      btn.style.cssText = 'opacity:0 !important;height:1px !important;overflow:hidden !important;position:absolute !important;pointer-events:auto !important;';
    }
  });
},100);
</script>""", height=0)


# ── Sidebar ────────────────────────────────────────────────────────────────────
def show_sidebar():
    is_internal = st.session_state.role == "internal"
    role_label = "INTERNAL" if is_internal else "PARTNER"

    with st.sidebar:
        # Logo + role badge header
        st.markdown(f"""
<div style="padding:18px 16px 14px;border-bottom:0.5px solid rgba(255,255,255,0.05);display:flex;align-items:center;gap:9px;">
  <div style="display:grid;grid-template-columns:repeat(4,5px);gap:1.5px;">
    <span style="width:5px;height:5px;border-radius:1px;background:#5B5FDE;"></span><span style="width:5px;height:5px;border-radius:1px;background:#5B5FDE;opacity:0.4;"></span><span style="width:5px;height:5px;border-radius:1px;background:#5B5FDE;opacity:0.4;"></span><span style="width:5px;height:5px;border-radius:1px;background:transparent;"></span>
    <span style="width:5px;height:5px;border-radius:1px;background:#5B5FDE;opacity:0.4;"></span><span style="width:5px;height:5px;border-radius:1px;background:#5B5FDE;"></span><span style="width:5px;height:5px;border-radius:1px;background:#5B5FDE;"></span><span style="width:5px;height:5px;border-radius:1px;background:#5B5FDE;opacity:0.4;"></span>
    <span style="width:5px;height:5px;border-radius:1px;background:#5B5FDE;opacity:0.4;"></span><span style="width:5px;height:5px;border-radius:1px;background:#5B5FDE;"></span><span style="width:5px;height:5px;border-radius:1px;background:#5B5FDE;"></span><span style="width:5px;height:5px;border-radius:1px;background:#5B5FDE;opacity:0.4;"></span>
    <span style="width:5px;height:5px;border-radius:1px;background:transparent;"></span><span style="width:5px;height:5px;border-radius:1px;background:#5B5FDE;opacity:0.4;"></span><span style="width:5px;height:5px;border-radius:1px;background:#5B5FDE;opacity:0.4;"></span><span style="width:5px;height:5px;border-radius:1px;background:#5B5FDE;"></span>
  </div>
  <span style="font-size:18px;font-weight:800;color:#E2E8F0;letter-spacing:-0.02em;">yuno</span>
  <span style="margin-left:auto;font-size:9px;font-weight:700;letter-spacing:0.07em;text-transform:uppercase;background:rgba(91,95,222,0.2);color:#A5B4FC;padding:2px 7px;border-radius:4px;border:0.5px solid rgba(91,95,222,0.3);">{role_label}</span>
</div>""", unsafe_allow_html=True)

        # Nav sections
        if is_internal:
            NAV_SECTIONS = [
                ("OVERVIEW", [
                    ("Home",        "Home"),
                    ("Partners",    "Partner Portfolio"),
                ]),
                ("PIPELINE & FLOW", [
                    ("MissionCtrl", "Partners In Flight"),
                    ("Pipeline",    "Partner Leads"),
                ]),
                ("PERFORMANCE", [
                    ("Performance", "Partner Health"),
                    ("Benchmarks",  "Rev Share"),
                ]),
                ("INTELLIGENCE", [
                    ("Insights",    "Market Intel"),
                ]),
                ("TOOLS", [
                    ("MerchSim",    "Merchant Sim"),
                ]),
            ]
        else:
            NAV_SECTIONS = [
                ("OVERVIEW", [
                    ("Home",        "Home"),
                    ("Partners",    "Partner Portfolio"),
                ]),
                ("PIPELINE & FLOW", [
                    ("Pipeline",    "Partner Leads"),
                ]),
                ("PERFORMANCE", [
                    ("Performance", "Partner Health"),
                    ("Benchmarks",  "Benchmarks"),
                ]),
                ("INTELLIGENCE", [
                    ("Insights",    "Market Intel"),
                ]),
            ]

        _NAV_BADGES = {"Partner Portfolio":"47","Partners In Flight":"3","Partner Leads":"24","Merchants":"12","Rev Share":"Apr 1","Recommendations":"12"}

        for section_label, items in NAV_SECTIONS:
            st.markdown(f'<div style="padding:14px 16px 4px;font-size:9.5px;font-weight:700;letter-spacing:0.1em;text-transform:uppercase;color:#334155;">{section_label}</div>', unsafe_allow_html=True)
            for page_key, label in items:
                if st.button(label, key=f"nav_{page_key}", use_container_width=True, type="secondary"):
                    st.session_state.page = page_key
                    st.query_params["page"] = page_key
                    st.query_params["role"] = st.session_state.role
                    st.rerun()

        # User footer
        if is_internal:
            st.markdown(f"""
<div style="margin-top:auto;padding:12px 14px;border-top:0.5px solid rgba(255,255,255,0.05);display:flex;align-items:center;gap:9px;">
  <div style="width:32px;height:32px;border-radius:50%;background:linear-gradient(135deg,#5B5FDE,#8B5CF6);display:flex;align-items:center;justify-content:center;font-size:11px;font-weight:700;color:white;flex-shrink:0;">DR</div>
  <div><div style="font-size:12.5px;font-weight:600;color:#CBD5E1;">Daniela Reyes</div><div style="font-size:10.5px;color:#475569;">Head of Partnerships · Yuno</div></div>
</div>""", unsafe_allow_html=True)
        else:
            st.markdown("""
<div style="margin-top:auto;padding:12px 14px;border-top:0.5px solid rgba(255,255,255,0.05);display:flex;align-items:center;gap:9px;">
  <div style="width:32px;height:32px;border-radius:50%;background:linear-gradient(135deg,#3B82F6,#60A5FA);display:flex;align-items:center;justify-content:center;font-size:11px;font-weight:700;color:white;flex-shrink:0;">TK</div>
  <div><div style="font-size:12.5px;font-weight:600;color:#CBD5E1;">Tom Kuehn</div><div style="font-size:10.5px;color:#475569;">LATAM Partnerships · Adyen</div></div>
</div>""", unsafe_allow_html=True)


# ── Stat Row ───────────────────────────────────────────────────────────────────
def stat_row(stats):
    cols = st.columns(len(stats))
    for col, s in zip(cols, stats):
        val_color = f"color:{s.get('val_color','#4F46E5')};"
        delta_color = {"up":"#16a34a","down":"#dc2626","flat":"#86868b"}.get(s.get("delta_type","flat"), "#86868b")
        col.markdown(f"""
<div style="background:#fff;box-shadow:0 1px 3px rgba(0,0,0,0.07),0 0 0 1px rgba(0,0,0,0.04);border-radius:12px;padding:18px 20px;position:relative;overflow:hidden;">
  <div style="position:absolute;top:0;left:0;right:0;height:3px;background:linear-gradient(90deg,#4F46E5,#818cf8);border-radius:12px 12px 0 0;"></div>
  <div style="font-size:10px;color:#86868b;letter-spacing:0.5px;text-transform:uppercase;font-weight:600;margin-bottom:8px;">{s['label']}</div>
  <div style="font-size:22px;font-weight:700;font-family:'Menlo',monospace;letter-spacing:-1px;line-height:1;{val_color}">{s['value']}</div>
  <div style="font-size:11px;margin-top:6px;color:{delta_color};">{s['delta']}</div>
</div>""", unsafe_allow_html=True)

# ── Pipeline View ──────────────────────────────────────────────────────────────
def show_pipeline():
    is_internal = st.session_state.role == "internal"

    # ── Hero ──────────────────────────────────────────────────────────────────
    st.markdown("""
<div style="background:#0C1220;border-radius:16px;padding:24px 28px;margin-bottom:14px;position:relative;overflow:hidden;">
  <div style="position:absolute;inset:0;background-image:radial-gradient(circle at 1px 1px,rgba(255,255,255,0.03) 1px,transparent 0);background-size:22px 22px;pointer-events:none;"></div>
  <div style="position:relative;z-index:1;">
    <div style="font-size:22px;font-weight:800;color:#F1F5F9;letter-spacing:-0.02em;margin-bottom:4px;">Partner Leads</div>
    <div style="font-size:12px;color:#64748b;margin-bottom:14px;">Track merchant opportunities from prospect to won. Filter by region to see deal flow.</div>
    <div style="display:flex;gap:10px;flex-wrap:wrap;">
      <div style="background:rgba(255,255,255,0.04);border:0.5px solid rgba(255,255,255,0.08);border-radius:8px;padding:10px 14px;text-align:center;"><div style="font-size:18px;font-weight:800;color:#818cf8;">$4.2M</div><div style="font-size:9px;color:#475569;text-transform:uppercase;margin-top:2px;">Pipeline Value</div></div>
      <div style="background:rgba(255,255,255,0.04);border:0.5px solid rgba(255,255,255,0.08);border-radius:8px;padding:10px 14px;text-align:center;"><div style="font-size:18px;font-weight:800;color:#C7D2FE;">24</div><div style="font-size:9px;color:#475569;text-transform:uppercase;margin-top:2px;">Active Deals</div></div>
      <div style="background:rgba(255,255,255,0.04);border:0.5px solid rgba(255,255,255,0.08);border-radius:8px;padding:10px 14px;text-align:center;"><div style="font-size:18px;font-weight:800;color:#FCD34D;">$175K</div><div style="font-size:9px;color:#475569;text-transform:uppercase;margin-top:2px;">Avg Deal Size</div></div>
      <div style="background:rgba(255,255,255,0.04);border:0.5px solid rgba(255,255,255,0.08);border-radius:8px;padding:10px 14px;text-align:center;"><div style="font-size:18px;font-weight:800;color:#6EE7B7;">38%</div><div style="font-size:9px;color:#475569;text-transform:uppercase;margin-top:2px;">Win Rate</div></div>
    </div>
  </div>
</div>""", unsafe_allow_html=True)

    # ── Region toggles ────────────────────────────────────────────────────────
    _pipe_regions = {
        "All":      {"deals":24,"value":"$4.2M","prospect":6,"qualified":5,"eval":7,"nego":4,"won":2},
        "LATAM":    {"deals":10,"value":"$1.8M","prospect":3,"qualified":2,"eval":3,"nego":1,"won":1},
        "Brasil":   {"deals":5,"value":"$0.9M","prospect":1,"qualified":1,"eval":2,"nego":1,"won":0},
        "N. America":{"deals":4,"value":"$0.8M","prospect":1,"qualified":1,"eval":1,"nego":1,"won":0},
        "EMEA":     {"deals":3,"value":"$0.5M","prospect":1,"qualified":1,"eval":0,"nego":0,"won":1},
        "APAC":     {"deals":2,"value":"$0.2M","prospect":0,"qualified":0,"eval":1,"nego":1,"won":0},
    }
    sel_pipe_region = st.session_state.get("pipe_region", "All")
    pr_cols = st.columns(len(_pipe_regions))
    for ri, (rname, rdata) in enumerate(_pipe_regions.items()):
        if pr_cols[ri].button(f"{rname} ({rdata['deals']})", key=f"pipe_reg_{ri}", use_container_width=True):
            st.session_state["pipe_region"] = rname
            st.rerun()

    # Style region buttons
    components.html(f"""<script>
setTimeout(function(){{
  var doc = window.parent.document;
  var active = '{sel_pipe_region}';
  doc.querySelectorAll('button').forEach(function(btn){{
    var t = btn.textContent.trim();
    if(t.match(/^(All|LATAM|Brasil|N\\. America|EMEA|APAC)\s*\(/)){{
      var isOn = t.startsWith(active+' (') || (active==='All' && t.startsWith('All ('));
      btn.style.cssText='background:'+(isOn?'#0C1220':'#fff')+' !important;border:1px solid '+(isOn?'#5B5FDE':'#E2E8F0')+' !important;border-radius:8px !important;padding:8px 8px !important;font-size:11px !important;color:'+(isOn?'#A5B4FC':'#64748b')+' !important;font-weight:'+(isOn?'700':'500')+' !important;cursor:pointer !important;box-shadow:none !important;';
    }}
  }});
}},200);
</script>""", height=0)

    # ── Region overview table ─────────────────────────────────────────────────
    rd = _pipe_regions.get(sel_pipe_region, _pipe_regions["All"])
    stages = [("Prospect", rd["prospect"], "#94a3b8"), ("Qualified", rd["qualified"], "#3B82F6"), ("Evaluation", rd["eval"], "#8B5CF6"), ("Negotiation", rd["nego"], "#F59E0B"), ("Won", rd["won"], "#10B981")]
    total_d = rd["deals"] or 1
    bar_html = '<div style="display:flex;gap:2px;height:8px;margin:10px 0;border-radius:4px;overflow:hidden;">'
    for label, count, color in stages:
        w = max(int(count / total_d * 100), 2) if count > 0 else 0
        if w > 0:
            bar_html += f'<div style="width:{w}%;background:{color};" title="{label}: {count}"></div>'
    bar_html += '</div>'

    stage_pills = "".join(f'<div style="text-align:center;flex:1;padding:8px;background:#fff;border:0.5px solid #E2E8F0;border-radius:8px;"><div style="font-size:16px;font-weight:800;color:{color};">{count}</div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;margin-top:2px;">{label}</div></div>' for label, count, color in stages)

    st.markdown(f'<div style="background:#fff;border:0.5px solid #E2E8F0;border-radius:12px;padding:16px 20px;margin-bottom:14px;">'
                f'<div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:6px;">'
                f'<span style="font-size:14px;font-weight:700;color:#0f172a;">{sel_pipe_region}</span>'
                f'<span style="font-size:13px;font-weight:800;color:#5B5FDE;">{rd["value"]}</span></div>'
                f'{bar_html}'
                f'<div style="display:flex;gap:6px;margin-top:8px;">{stage_pills}</div>'
                f'</div>', unsafe_allow_html=True)

    st.markdown('<div style="font-size:10px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:1px;margin-bottom:8px;">Deal Board</div>', unsafe_allow_html=True)

    st.markdown("""
<div style="margin-top:20px;margin-bottom:10px;">
  <span style="font-size:16px;font-weight:700;color:#1d1d1f;letter-spacing:-0.3px;">Deal Board</span>
  <span style="font-size:11px;color:#86868b;margin-left:10px;">Drag to move · Click to open · Internal notes visible to Yuno BD only</span>
</div>
""", unsafe_allow_html=True)

    # Kanban board as HTML (no <style> tag — styles injected via JS already)
    kanban_html = """
<div class="pipeline-wrap">
<div class="pipeline-board">

<!-- Prospect -->
<div class="stage-col">
  <div class="stage-head"><span class="stage-name" style="color:#6e6e73;">Prospect</span><span class="stage-num">6</span></div>
  <div class="pcard"><div class="pcard-name">Qatar Airways</div><div class="pcard-meta"><span class="pill cat-psp">Travel</span><span class="pill p-amber">Warm</span></div><div class="pcard-val">$80K ARR est.</div><div class="pcard-owner"><span class="mini-av ma-y">DR</span>Daniela R. · via Adyen ref</div><div class="pbar"><div class="pfill" style="width:15%;background:#86868b;"></div></div></div>
  <div class="pcard"><div class="pcard-name">Bancolombia</div><div class="pcard-meta"><span class="pill cat-acquirer">Banking</span><span class="pill p-blue">Inbound</span></div><div class="pcard-val">$120K ARR est.</div><div class="pcard-owner"><span class="mini-av ma-e">JL</span>Jorge L. · Direct</div><div class="pbar"><div class="pfill" style="width:10%;background:#3b82f6;"></div></div></div>
  <div class="pcard"><div class="pcard-name">iFood</div><div class="pcard-meta"><span class="pill cat-baas">Food</span><span class="pill p-amber">Pilot</span></div><div class="pcard-val">$200K ARR est.</div><div class="pcard-owner"><span class="mini-av ma-y">MC</span>Maria C. · Direct</div><div class="pbar"><div class="pfill" style="width:12%;background:#f59e0b;"></div></div></div>
</div>

<!-- Qualified -->
<div class="stage-col">
  <div class="stage-head"><span class="stage-name" style="color:#60a5fa;">Qualified</span><span class="stage-num">5</span></div>
  <div class="pcard"><div class="pcard-name">Rappi</div><div class="pcard-meta"><span class="pill cat-apm">Super App</span><span class="pill p-green">Hot</span></div><div class="pcard-val">$350K ARR est.</div><div class="pcard-owner"><span class="mini-av ma-y">DR</span>Daniela R.</div><div class="pbar"><div class="pfill" style="width:30%;background:#22c55e;"></div></div></div>
  <div class="pcard"><div class="pcard-name">Falabella</div><div class="pcard-meta"><span class="pill cat-acquirer">Retail</span><span class="pill p-blue">Active</span></div><div class="pcard-val">$180K ARR est.</div><div class="pcard-owner"><span class="mini-av ma-e">AP</span>Ana P. · Partner ref</div><div class="pbar"><div class="pfill" style="width:28%;background:#3b82f6;"></div></div></div>
  <div class="pcard"><div class="pcard-name">Despegar</div><div class="pcard-meta"><span class="pill cat-fraud">Travel</span><span class="pill p-purple">Trial</span></div><div class="pcard-val">$95K ARR est.</div><div class="pcard-owner"><span class="mini-av ma-y">DG</span>Diego G.</div><div class="pbar"><div class="pfill" style="width:25%;background:#a855f7;"></div></div></div>
</div>

<!-- Evaluation -->
<div class="stage-col">
  <div class="stage-head"><span class="stage-name" style="color:#c084fc;">Evaluation</span><span class="stage-num">7</span></div>
  <div class="pcard"><div class="pcard-name">Spotify</div><div class="pcard-meta"><span class="pill cat-psp">Streaming</span><span class="pill p-green">Hot</span></div><div class="pcard-val">$620K ARR est.</div><div class="pcard-owner"><span class="mini-av ma-e">RV</span>Ricardo V. · Direct</div><div class="pbar"><div class="pfill" style="width:55%;background:#a855f7;"></div></div></div>
  <div class="pcard"><div class="pcard-name">PedidosYa</div><div class="pcard-meta"><span class="pill cat-apm">Food</span><span class="pill p-teal">Active</span></div><div class="pcard-val">$240K ARR est.</div><div class="pcard-owner"><span class="mini-av ma-y">DR</span>Daniela R.</div><div class="pbar"><div class="pfill" style="width:50%;background:#14b8a6;"></div></div></div>
  <div class="pcard"><div class="pcard-name">Claro</div><div class="pcard-meta"><span class="pill cat-acquirer">Telco</span><span class="pill p-amber">Stalled</span></div><div class="pcard-val">$310K ARR est.</div><div class="pcard-owner"><span class="mini-av ma-e">LM</span>Luisa M. · Partner</div><div class="pbar"><div class="pfill" style="width:48%;background:#f59e0b;"></div></div></div>
  <div class="pcard"><div class="pcard-name">MercadoLibre</div><div class="pcard-meta"><span class="pill cat-baas">E-commerce</span><span class="pill p-amber">New</span></div><div class="pcard-val">$440K ARR est.</div><div class="pcard-owner"><span class="mini-av ma-y">MC</span>Maria C.</div><div class="pbar"><div class="pfill" style="width:42%;background:#f59e0b;"></div></div></div>
</div>

<!-- Negotiation -->
<div class="stage-col">
  <div class="stage-head"><span class="stage-name" style="color:#fbbf24;">Negotiation</span><span class="stage-num">2</span></div>
  <div class="pcard"><div class="pcard-name">Uber</div><div class="pcard-meta"><span class="pill cat-acquirer">Mobility</span><span class="pill p-green">High Priority</span></div><div class="pcard-val">$850K ARR est.</div><div class="pcard-owner"><span class="mini-av ma-e">FM</span>Felipe M. · Direct</div><div class="pbar"><div class="pfill" style="width:75%;background:#4F46E5;"></div></div></div>
  <div class="pcard"><div class="pcard-name">Cinépolis</div><div class="pcard-meta"><span class="pill cat-apm">Entertainment</span><span class="pill p-blue">Regional</span></div><div class="pcard-val">$110K ARR est.</div><div class="pcard-owner"><span class="mini-av ma-y">DG</span>Diego G.</div><div class="pbar"><div class="pfill" style="width:70%;background:#14b8a6;"></div></div></div>
</div>

<!-- Won -->
<div class="stage-col">
  <div class="stage-head"><span class="stage-name" style="color:#4ade80;">Won</span><span class="stage-num">2</span></div>
  <div class="pcard" style="border-color:rgba(34,197,94,.25);"><div class="pcard-name">Netflix</div><div class="pcard-meta"><span class="pill cat-psp">Streaming</span><span class="pill p-green">Live</span></div><div class="pcard-val">$1.2M ARR signed</div><div class="pcard-owner"><span class="mini-av ma-e">TK</span>Tom K. · Direct</div><div class="pbar"><div class="pfill" style="width:100%;background:#22c55e;"></div></div></div>
  <div class="pcard" style="border-color:rgba(34,197,94,.25);"><div class="pcard-name">Linio</div><div class="pcard-meta"><span class="pill cat-fraud">E-commerce</span><span class="pill p-green">Live</span></div><div class="pcard-val">$290K ARR signed</div><div class="pcard-owner"><span class="mini-av ma-y">DR</span>Daniela R.</div><div class="pbar"><div class="pfill" style="width:100%;background:#22c55e;"></div></div></div>
</div>

</div></div>"""
    st.markdown(kanban_html, unsafe_allow_html=True)

    # Bottom row: funnel + activity
    col_funnel, col_activity = st.columns([3, 2])

    with col_funnel:
        st.markdown('<div style="background:#fff;box-shadow:0 1px 3px rgba(0,0,0,0.07),0 0 0 1px rgba(0,0,0,0.04);border:none;border-radius:10px;padding:18px;">', unsafe_allow_html=True)
        st.markdown('<div style="font-size:13px;font-weight:700;color:#1d1d1f;margin-bottom:2px;">Pipeline Funnel</div><div style="font-size:11px;color:#86868b;margin-bottom:14px;">Current quarter · all categories</div>', unsafe_allow_html=True)
        funnel_data = [("Prospect",6,100,"#86868b"),("Qualified",5,83,"#6e6e73"),("Evaluation",7,100,"#a855f7"),("Negotiation",4,50,"#f59e0b"),("Won",2,33,"#22c55e")]
        funnel_html = ""
        for label, count, pct, color in funnel_data:
            funnel_html += f"""<div style="display:flex;align-items:center;gap:10px;margin-bottom:8px;">
  <span style="font-size:11px;color:#6e6e73;width:90px;text-align:right;flex-shrink:0;">{label}</span>
  <div style="flex:1;height:24px;background:#f9f9fb;border-radius:4px;overflow:hidden;">
    <div style="width:{pct}%;height:100%;background:{color};display:flex;align-items:center;padding-left:10px;font-size:10px;font-family:monospace;font-weight:600;color:rgba(0,0,0,0.7);">{count} deals</div>
  </div>
  <span style="font-size:11px;font-family:monospace;color:#6e6e73;width:20px;">{count}</span>
</div>"""
        st.markdown(funnel_html, unsafe_allow_html=True)

        if is_internal:
            st.markdown("""
<div style="border:1px solid rgba(245,158,11,.22);background:rgba(245,158,11,.04);border-radius:10px;padding:13px 16px;margin-top:16px;">
  <div style="font-size:9.5px;font-weight:700;color:#d97706;letter-spacing:.9px;text-transform:uppercase;margin-bottom:10px;">🔒 Internal Yuno Only</div>
  <div style="display:flex;justify-content:space-between;align-items:center;padding:5px 0;border-bottom:1px solid rgba(0,0,0,0.06);font-size:12px;"><span style="color:#6e6e73;">Blended margin (est.)</span><span style="font-family:monospace;color:#d97706;">61%</span></div>
  <div style="display:flex;justify-content:space-between;align-items:center;padding:5px 0;border-bottom:1px solid rgba(0,0,0,0.06);font-size:12px;"><span style="color:#6e6e73;">Committed quota coverage</span><span style="font-family:monospace;color:#d97706;">2.4×</span></div>
  <div style="display:flex;justify-content:space-between;align-items:center;padding:5px 0;font-size:12px;"><span style="color:#6e6e73;">At-risk deals (&gt;90d stalled)</span><span style="font-family:monospace;color:#d97706;">3</span></div>
</div>""", unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)

    with col_activity:
        st.markdown("""
<div style="background:#fff;box-shadow:0 1px 3px rgba(0,0,0,0.07),0 0 0 1px rgba(0,0,0,0.04);border:none;border-radius:10px;padding:18px;">
  <div style="font-size:13px;font-weight:700;color:#1d1d1f;margin-bottom:2px;">Recent Activity</div>
  <div style="font-size:11px;color:#86868b;margin-bottom:14px;">Last 48 hours</div>
  <div style="display:flex;gap:10px;padding:9px 0;border-bottom:1px solid rgba(0,0,0,0.07);">
    <div style="width:7px;height:7px;border-radius:50%;background:#22c55e;margin-top:4px;flex-shrink:0;"></div>
    <div><div style="font-size:12px;color:#6e6e73;line-height:1.55;"><strong style="color:#1d1d1f;">Uber</strong> moved to Negotiation. MSA draft sent.</div><div style="font-size:10px;color:#86868b;font-family:monospace;margin-top:2px;">2h ago · Daniela R.</div></div>
  </div>
  <div style="display:flex;gap:10px;padding:9px 0;border-bottom:1px solid rgba(0,0,0,0.07);">
    <div style="width:7px;height:7px;border-radius:50%;background:#3b82f6;margin-top:4px;flex-shrink:0;"></div>
    <div><div style="font-size:12px;color:#6e6e73;line-height:1.55;"><strong style="color:#1d1d1f;">Spotify</strong> technical review call scheduled for Dec 18.</div><div style="font-size:10px;color:#86868b;font-family:monospace;margin-top:2px;">5h ago · Ricardo V.</div></div>
  </div>
  <div style="display:flex;gap:10px;padding:9px 0;border-bottom:1px solid rgba(0,0,0,0.07);">
    <div style="width:7px;height:7px;border-radius:50%;background:#f59e0b;margin-top:4px;flex-shrink:0;"></div>
    <div><div style="font-size:12px;color:#6e6e73;line-height:1.55;"><strong style="color:#1d1d1f;">MercadoLibre</strong> added as new e-commerce opportunity.</div><div style="font-size:10px;color:#86868b;font-family:monospace;margin-top:2px;">Yesterday · Maria C.</div></div>
  </div>
  <div style="display:flex;gap:10px;padding:9px 0;border-bottom:1px solid rgba(0,0,0,0.07);">
    <div style="width:7px;height:7px;border-radius:50%;background:#4F46E5;margin-top:4px;flex-shrink:0;"></div>
    <div><div style="font-size:12px;color:#6e6e73;line-height:1.55;"><strong style="color:#1d1d1f;">Netflix</strong> integration live. First transaction processed.</div><div style="font-size:10px;color:#86868b;font-family:monospace;margin-top:2px;">Yesterday · Tom K.</div></div>
  </div>
  <div style="display:flex;gap:10px;padding:9px 0;">
    <div style="width:7px;height:7px;border-radius:50%;background:#a855f7;margin-top:4px;flex-shrink:0;"></div>
    <div><div style="font-size:12px;color:#6e6e73;line-height:1.55;"><strong style="color:#1d1d1f;">Despegar</strong> trial extended to Jan 15. Positive signal.</div><div style="font-size:10px;color:#86868b;font-family:monospace;margin-top:2px;">2d ago · Diego G.</div></div>
  </div>
</div>""", unsafe_allow_html=True)

# ── Partners View ──────────────────────────────────────────────────────────────
def _partner_detail(p):
    """Render partner detail page with tabs (Salesforce-style)."""
    if st.button("← Back to Directory", key="back_to_dir"):
        st.session_state.selected_partner = None
        st.session_state.partner_tab = "Overview"
        st.rerun()

    STATUS_BADGE = {"Live": ("background:#065f46;color:#fff;", "ACTIVE"), "Agreement Signed": ("background:#1e40af;color:#fff;", "AGREEMENT SIGNED"), "Agreement Review": ("background:#7c3aed;color:#fff;", "IN REVIEW"), "Negotiation": ("background:#d97706;color:#fff;", "NEGOTIATION"), "Prospect": ("background:#92400e;color:#fff;", "PROSPECT"), "Lost": ("background:#dc2626;color:#fff;", "LOST")}
    badge_style, badge_label = STATUS_BADGE.get(p["status"], ("background:#64748b;color:#fff;", p["status"].upper()))
    scope = "Global" if p.get("region","") in ("Global","APAC/EMEA") or p.get("country","") == "Global" else "Regional"
    strategic_tag = '<span style="background:#fef3c7;color:#92400e;font-size:10px;font-weight:700;padding:3px 10px;border-radius:4px;margin-left:10px;letter-spacing:0.5px;">STRATEGIC</span>' if p.get("strategic") else ""

    st.markdown(f"""
<div style="display:flex;justify-content:space-between;align-items:flex-start;margin-bottom:8px;">
  <div>
    <div style="display:flex;align-items:center;gap:0;"><span style="font-size:28px;font-weight:800;color:#0f172a;letter-spacing:-0.5px;">{p['name']}</span>{strategic_tag}</div>
    <div style="font-size:13px;color:#94a3b8;margin-top:4px;">{p['type']} &middot; {p.get('country','')} &middot; {p.get('region','')} &middot; {p.get('manager','Unassigned')}</div>
  </div>
  <span style="{badge_style}font-size:11px;font-weight:700;padding:6px 18px;border-radius:20px;letter-spacing:0.5px;">{badge_label}</span>
</div>""", unsafe_allow_html=True)

    # Completeness
    checks = [p.get("nda"), p.get("revshare"), p.get("integration_ready"), p.get("integration_used"), bool(p.get("comments"))]
    check_labels = ["NDA Signed", "Rev Share Contract", "Integration Ready", "Merchant Using", "Notes Added"]
    done = sum(1 for c in checks if c)
    pct = int(done / len(checks) * 100)
    bar_color = "#22c55e" if pct >= 70 else "#f59e0b" if pct >= 40 else "#ef4444"
    tags_html = "".join(f'<span style="font-size:10px;color:{"#065f46" if checks[i] else "#94a3b8"};padding:3px 8px;border-radius:4px;background:{"#ecfdf5" if checks[i] else "#f8fafc"};border:1px solid {"#a7f3d0" if checks[i] else "#e2e8f0"};">{check_labels[i]}</span>' for i in range(len(checks)))
    st.markdown(f"""
<div style="margin:12px 0 16px;">
  <div style="display:flex;align-items:center;gap:10px;margin-bottom:6px;">
    <div style="flex:1;height:5px;background:#e5e7eb;border-radius:3px;overflow:hidden;"><div style="width:{pct}%;height:100%;background:{bar_color};border-radius:3px;transition:width 0.3s;"></div></div>
    <span style="font-size:12px;font-weight:700;color:{bar_color};">{pct}%</span>
  </div>
  <div style="display:flex;flex-wrap:wrap;gap:5px;">{tags_html}</div>
</div>""", unsafe_allow_html=True)

    TABS = ["Overview", "Contacts", "Contracts", "SLA", "Revenue", "Questionnaires", "Audit Log"]
    tab_cols = st.columns(len(TABS))
    for i, (col, tab) in enumerate(zip(tab_cols, TABS)):
        if col.button(tab, key=f"ptab_{i}", use_container_width=True):
            st.session_state.partner_tab = tab
            st.rerun()

    active_tab = st.session_state.partner_tab

    # ── Overview ──────────────────────────────────────────────────────────────
    if active_tab == "Overview":
        st.markdown(f"""
<div style="border:1px solid #e2e8f0;border-radius:10px;padding:28px;margin-top:16px;background:#fff;">
  <div style="font-size:18px;font-weight:700;color:#0f172a;margin-bottom:24px;">Basic Information</div>
  <div style="display:grid;grid-template-columns:1fr 1fr;gap:0 48px;">
    <div style="padding:14px 0;border-bottom:1px solid #f1f5f9;"><div style="font-size:11px;font-weight:600;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-bottom:4px;">Provider Name</div><div style="font-size:15px;font-weight:600;color:#0f172a;">{p['name']}</div></div>
    <div style="padding:14px 0;border-bottom:1px solid #f1f5f9;"><div style="font-size:11px;font-weight:600;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-bottom:4px;">Type</div><div style="font-size:15px;font-weight:600;color:#0f172a;">{p['type']}</div></div>
    <div style="padding:14px 0;border-bottom:1px solid #f1f5f9;"><div style="font-size:11px;font-weight:600;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-bottom:4px;">Scope</div><div style="font-size:15px;font-weight:600;color:#0f172a;">{scope}</div></div>
    <div style="padding:14px 0;border-bottom:1px solid #f1f5f9;"><div style="font-size:11px;font-weight:600;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-bottom:4px;">Deal Stage</div><div style="font-size:15px;font-weight:600;color:#0f172a;">{p['status']}</div></div>
    <div style="padding:14px 0;border-bottom:1px solid #f1f5f9;"><div style="font-size:11px;font-weight:600;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-bottom:4px;">Country</div><div style="font-size:15px;font-weight:600;color:#0f172a;">{p.get('country','-')}</div></div>
    <div style="padding:14px 0;border-bottom:1px solid #f1f5f9;"><div style="font-size:11px;font-weight:600;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-bottom:4px;">Region</div><div style="font-size:15px;font-weight:600;color:#0f172a;">{p.get('region','-')}</div></div>
    <div style="padding:14px 0;border-bottom:1px solid #f1f5f9;"><div style="font-size:11px;font-weight:600;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-bottom:4px;">Partner Tier</div><div style="font-size:15px;font-weight:600;color:#0f172a;">{p.get('tier','-')}</div></div>
    <div style="padding:14px 0;border-bottom:1px solid #f1f5f9;"><div style="font-size:11px;font-weight:600;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-bottom:4px;">Management</div><div style="font-size:15px;font-weight:600;color:#0f172a;">{p.get('mgmt_type','-')}</div></div>
    <div style="padding:14px 0;border-bottom:1px solid #f1f5f9;"><div style="font-size:11px;font-weight:600;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-bottom:4px;">Partner Manager</div><div style="font-size:15px;font-weight:600;color:#0f172a;">{p.get('manager','-')}</div></div>
    <div style="padding:14px 0;border-bottom:1px solid #f1f5f9;"><div style="font-size:11px;font-weight:600;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-bottom:4px;">Strategic Partner</div><div style="font-size:15px;font-weight:600;color:#0f172a;">{'Yes' if p.get('strategic') else 'No'}</div></div>
    <div style="padding:14px 0;"><div style="font-size:11px;font-weight:600;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-bottom:4px;">Notes</div><div style="font-size:14px;color:#64748b;">{p.get('comments') or '-'}</div></div>
    <div style="padding:14px 0;"><div style="font-size:11px;font-weight:600;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-bottom:4px;">NDA Signed</div><div style="font-size:15px;font-weight:600;color:#0f172a;">{'Yes' if p.get('nda') else 'No'}</div></div>
  </div>
</div>""", unsafe_allow_html=True)

        # Region context card
        r_stats = _REGION_STATS.get(p.get("region",""), None)
        if r_stats:
            st.markdown(f"""
<div style="border:1px solid #e2e8f0;border-radius:10px;padding:24px;margin-top:12px;background:#fff;">
  <div style="font-size:16px;font-weight:700;color:#0f172a;margin-bottom:16px;">Region Context — {p.get('region','')}</div>
  <div style="display:grid;grid-template-columns:repeat(5,1fr);gap:12px;">
    <div style="text-align:center;padding:12px;border:1px solid #f1f5f9;border-radius:8px;"><div style="font-size:22px;font-weight:800;color:#0f172a;">{r_stats['total']}</div><div style="font-size:10px;color:#94a3b8;margin-top:2px;">Total</div></div>
    <div style="text-align:center;padding:12px;border:1px solid #f1f5f9;border-radius:8px;"><div style="font-size:22px;font-weight:800;color:#065f46;">{r_stats['live']}</div><div style="font-size:10px;color:#94a3b8;margin-top:2px;">Live</div></div>
    <div style="text-align:center;padding:12px;border:1px solid #f1f5f9;border-radius:8px;"><div style="font-size:22px;font-weight:800;color:#7c3aed;">{r_stats['strategic']}</div><div style="font-size:10px;color:#94a3b8;margin-top:2px;">Strategic</div></div>
    <div style="text-align:center;padding:12px;border:1px solid #f1f5f9;border-radius:8px;"><div style="font-size:22px;font-weight:800;color:#1e40af;">{r_stats['tier1']}</div><div style="font-size:10px;color:#94a3b8;margin-top:2px;">Tier 1</div></div>
    <div style="text-align:center;padding:12px;border:1px solid #f1f5f9;border-radius:8px;"><div style="font-size:22px;font-weight:800;color:#065f46;">{r_stats['revshare']}</div><div style="font-size:10px;color:#94a3b8;margin-top:2px;">Rev Share</div></div>
  </div>
</div>""", unsafe_allow_html=True)

        # Manager stats card
        m_stats = _MANAGER_STATS.get(p.get("manager",""), None)
        if m_stats:
            st.markdown(f"""
<div style="border:1px solid #e2e8f0;border-radius:10px;padding:24px;margin-top:12px;background:#fff;">
  <div style="font-size:16px;font-weight:700;color:#0f172a;margin-bottom:4px;">Account Manager — {p.get('manager','')}</div>
  <div style="font-size:12px;color:#94a3b8;margin-bottom:16px;">{m_stats['regions']}</div>
  <div style="display:grid;grid-template-columns:repeat(4,1fr);gap:12px;">
    <div style="text-align:center;padding:12px;border:1px solid #f1f5f9;border-radius:8px;"><div style="font-size:22px;font-weight:800;color:#7c3aed;">{m_stats['strategic']}</div><div style="font-size:10px;color:#94a3b8;margin-top:2px;">Strategic</div></div>
    <div style="text-align:center;padding:12px;border:1px solid #f1f5f9;border-radius:8px;"><div style="font-size:22px;font-weight:800;color:#1e40af;">{m_stats['tier1']}</div><div style="font-size:10px;color:#94a3b8;margin-top:2px;">Tier 1</div></div>
    <div style="text-align:center;padding:12px;border:1px solid #f1f5f9;border-radius:8px;"><div style="font-size:22px;font-weight:800;color:#0f172a;">{m_stats['tier2']}</div><div style="font-size:10px;color:#94a3b8;margin-top:2px;">Tier 2</div></div>
    <div style="text-align:center;padding:12px;border:1px solid #f1f5f9;border-radius:8px;"><div style="font-size:22px;font-weight:800;color:#64748b;">{m_stats['tier3']}</div><div style="font-size:10px;color:#94a3b8;margin-top:2px;">Tier 3</div></div>
  </div>
</div>""", unsafe_allow_html=True)

    # ── Contacts ──────────────────────────────────────────────────────────────
    elif active_tab == "Contacts":
        partner_contacts = [c for c in CONTACTS if p["name"].split("/")[0].split("(")[0].strip().lower() in c["company"].lower()] or CONTACTS[:2]
        rows_html = ""
        for c in partner_contacts:
            rows_html += f"""<tr style="border-bottom:1px solid #f1f5f9;">
  <td style="padding:14px 16px;font-size:13px;font-weight:600;color:#0f172a;">{c['name']}</td>
  <td style="padding:14px 16px;font-size:13px;color:#64748b;">{c['role']}</td>
  <td style="padding:14px 16px;font-size:13px;color:#64748b;">{c.get('company','')}</td>
  <td style="padding:14px 16px;font-size:13px;color:#64748b;">{c['last']}</td>
  <td style="padding:14px 16px;font-size:13px;color:#64748b;">{c['rel']}</td>
</tr>"""
        st.markdown(f"""
<div style="border:1px solid #e2e8f0;border-radius:10px;overflow:hidden;margin-top:16px;background:#fff;">
  <div style="padding:20px 20px 12px;font-size:18px;font-weight:700;color:#0f172a;">Partner Contacts</div>
  <table style="width:100%;border-collapse:collapse;">
    <thead><tr><th style="padding:10px 16px;font-size:11px;font-weight:600;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;text-align:left;border-bottom:2px solid #e2e8f0;">Name</th><th style="padding:10px 16px;font-size:11px;font-weight:600;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;text-align:left;border-bottom:2px solid #e2e8f0;">Role</th><th style="padding:10px 16px;font-size:11px;font-weight:600;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;text-align:left;border-bottom:2px solid #e2e8f0;">Company</th><th style="padding:10px 16px;font-size:11px;font-weight:600;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;text-align:left;border-bottom:2px solid #e2e8f0;">Last Contact</th><th style="padding:10px 16px;font-size:11px;font-weight:600;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;text-align:left;border-bottom:2px solid #e2e8f0;">Relationship</th></tr></thead>
    <tbody>{rows_html}</tbody>
  </table>
</div>""", unsafe_allow_html=True)

    # ── Contracts ─────────────────────────────────────────────────────────────
    elif active_tab == "Contracts":
        has_contract = p.get("revshare") or p["status"] == "Live"
        if has_contract:
            rev_status = "ACTIVE" if p.get("revshare_active") else "SIGNED"
            rev_bg = "#d1fae5" if p.get("revshare_active") else "#dbeafe"
            rev_color = "#065f46" if p.get("revshare_active") else "#1e40af"
            st.markdown(f"""
<div style="border:1px solid #e2e8f0;border-radius:10px;padding:24px;margin-top:16px;background:#fff;">
  <div style="font-size:18px;font-weight:700;color:#0f172a;margin-bottom:20px;">Contracts</div>
  <div style="border:1px solid #e2e8f0;border-radius:8px;padding:18px;">
    <div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:14px;">
      <div style="font-size:14px;font-weight:700;color:#0f172a;">Revenue Share Agreement</div>
      <span style="background:{rev_bg};color:{rev_color};font-size:11px;font-weight:600;padding:3px 12px;border-radius:20px;">{rev_status}</span>
    </div>
    <div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:16px;">
      <div><div style="font-size:11px;color:#94a3b8;">NDA</div><div style="font-size:13px;font-weight:600;color:#0f172a;">{'Signed' if p.get('nda') else 'Pending'}</div></div>
      <div><div style="font-size:11px;color:#94a3b8;">Rev Share Active</div><div style="font-size:13px;font-weight:600;color:#0f172a;">{'Yes' if p.get('revshare_active') else 'No'}</div></div>
      <div><div style="font-size:11px;color:#94a3b8;">Management Type</div><div style="font-size:13px;font-weight:600;color:#0f172a;">{p.get('mgmt_type','-')}</div></div>
    </div>
  </div>
</div>""", unsafe_allow_html=True)
        else:
            st.markdown('<div style="border:1px solid #e2e8f0;border-radius:10px;padding:48px;margin-top:16px;background:#fff;text-align:center;"><div style="font-size:15px;color:#94a3b8;">No contracts on file</div><div style="font-size:12px;color:#cbd5e1;margin-top:4px;">Contracts will appear here once the partnership is formalized.</div></div>', unsafe_allow_html=True)

    # ── SLA ───────────────────────────────────────────────────────────────────
    elif active_tab == "SLA":
        if p["status"] == "Live":
            st.markdown("""
<div style="border:1px solid #e2e8f0;border-radius:10px;padding:24px;margin-top:16px;background:#fff;">
  <div style="font-size:18px;font-weight:700;color:#0f172a;margin-bottom:20px;">Service Level Agreements</div>
  <div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:14px;">
    <div style="border:1px solid #e2e8f0;border-radius:8px;padding:18px;"><div style="font-size:11px;color:#94a3b8;margin-bottom:6px;">Uptime SLA</div><div style="font-size:26px;font-weight:700;color:#065f46;">99.95%</div><div style="font-size:11px;color:#94a3b8;margin-top:6px;">Target: 99.9%</div></div>
    <div style="border:1px solid #e2e8f0;border-radius:8px;padding:18px;"><div style="font-size:11px;color:#94a3b8;margin-bottom:6px;">Response Time</div><div style="font-size:26px;font-weight:700;color:#0f172a;">240ms</div><div style="font-size:11px;color:#94a3b8;margin-top:6px;">Target: &lt;500ms</div></div>
    <div style="border:1px solid #e2e8f0;border-radius:8px;padding:18px;"><div style="font-size:11px;color:#94a3b8;margin-bottom:6px;">Support Response</div><div style="font-size:26px;font-weight:700;color:#0f172a;">2h</div><div style="font-size:11px;color:#94a3b8;margin-top:6px;">Target: &lt;4h (P1)</div></div>
  </div>
</div>""", unsafe_allow_html=True)
        else:
            st.markdown('<div style="border:1px solid #e2e8f0;border-radius:10px;padding:48px;margin-top:16px;background:#fff;text-align:center;"><div style="font-size:15px;color:#94a3b8;">No SLA defined yet</div><div style="font-size:12px;color:#cbd5e1;margin-top:4px;">SLAs will be configured once the integration goes live.</div></div>', unsafe_allow_html=True)

    # ── Revenue ───────────────────────────────────────────────────────────────
    elif active_tab == "Revenue":
        # Match partner name to rev share data (fuzzy)
        rev_match = None
        pname_lower = p["name"].lower()
        for rk, rv in _REVSHARE_BY_PARTNER.items():
            if rk.lower() in pname_lower or pname_lower.split("/")[0].split("(")[0].strip() in rk.lower() or rk.lower().split("/")[0] in pname_lower:
                rev_match = (rk, rv)
                break
        region = p.get("region", "")
        region_pct = _REVSHARE_BY_REGION.get(region, 0)
        region_stats = _REGION_STATS.get(region, {})

        if rev_match or p.get("revshare_active"):
            rev_name, rev_val = rev_match if rev_match else (p["name"], 0)
            rev_display = f"${rev_val:,.0f}" if rev_val else "Active"
            arr_est = f"${rev_val * 12:,.0f}" if rev_val else "-"
            st.markdown(f"""
<div style="border:1px solid #e2e8f0;border-radius:10px;padding:24px;margin-top:16px;background:#fff;">
  <div style="font-size:18px;font-weight:700;color:#0f172a;margin-bottom:20px;">Revenue Summary</div>
  <div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:14px;margin-bottom:20px;">
    <div style="border:1px solid #e2e8f0;border-radius:8px;padding:18px;"><div style="font-size:11px;color:#94a3b8;margin-bottom:6px;">Feb 2026 Rev Share</div><div style="font-size:26px;font-weight:700;color:#065f46;">{rev_display}</div></div>
    <div style="border:1px solid #e2e8f0;border-radius:8px;padding:18px;"><div style="font-size:11px;color:#94a3b8;margin-bottom:6px;">Est. ARR</div><div style="font-size:26px;font-weight:700;color:#0f172a;">{arr_est}</div></div>
    <div style="border:1px solid #e2e8f0;border-radius:8px;padding:18px;"><div style="font-size:11px;color:#94a3b8;margin-bottom:6px;">Region Rev Share</div><div style="font-size:26px;font-weight:700;color:#0f172a;">{region_stats.get('revshare','-')}</div></div>
  </div>
  <div style="font-size:13px;font-weight:600;color:#0f172a;margin-bottom:10px;">Portfolio Revenue Context</div>
  <div style="display:grid;grid-template-columns:1fr 1fr 1fr 1fr;gap:12px;">
    <div style="padding:10px 0;border-bottom:1px solid #f1f5f9;"><div style="font-size:11px;color:#94a3b8;">Total Feb 2026</div><div style="font-size:14px;font-weight:600;color:#0f172a;">$123K</div></div>
    <div style="padding:10px 0;border-bottom:1px solid #f1f5f9;"><div style="font-size:11px;color:#94a3b8;">Portfolio ARR</div><div style="font-size:14px;font-weight:600;color:#0f172a;">$1.48M</div></div>
    <div style="padding:10px 0;border-bottom:1px solid #f1f5f9;"><div style="font-size:11px;color:#94a3b8;">2025 Total</div><div style="font-size:14px;font-weight:600;color:#0f172a;">$1.78M</div></div>
    <div style="padding:10px 0;border-bottom:1px solid #f1f5f9;"><div style="font-size:11px;color:#94a3b8;">Peak Month</div><div style="font-size:14px;font-weight:600;color:#0f172a;">$245K (Oct 25)</div></div>
  </div>
</div>""", unsafe_allow_html=True)
        else:
            st.markdown('<div style="border:1px solid #e2e8f0;border-radius:10px;padding:48px;margin-top:16px;background:#fff;text-align:center;"><div style="font-size:15px;color:#94a3b8;">No revenue data yet</div><div style="font-size:12px;color:#cbd5e1;margin-top:4px;">Revenue data will populate once rev share is active.</div></div>', unsafe_allow_html=True)

    # ── Questionnaires ────────────────────────────────────────────────────────
    elif active_tab == "Questionnaires":
        kyp_done = bool(p.get("integration_ready"))
        st.markdown(f"""
<div style="border:1px solid #e2e8f0;border-radius:10px;padding:24px;margin-top:16px;background:#fff;">
  <div style="font-size:18px;font-weight:700;color:#0f172a;margin-bottom:20px;">Questionnaires</div>
  <div style="border:1px solid #e2e8f0;border-radius:8px;padding:16px;margin-bottom:10px;"><div style="display:flex;justify-content:space-between;align-items:center;"><div><div style="font-size:14px;font-weight:600;color:#0f172a;">Know Your Partner (KYP)</div><div style="font-size:12px;color:#94a3b8;margin-top:2px;">Due diligence, compliance, and operational assessment</div></div><span style="background:{'#d1fae5' if kyp_done else '#fef3c7'};color:{'#065f46' if kyp_done else '#92400e'};font-size:11px;font-weight:600;padding:3px 12px;border-radius:20px;">{'COMPLETED' if kyp_done else 'PENDING'}</span></div></div>
  <div style="border:1px solid #e2e8f0;border-radius:8px;padding:16px;margin-bottom:10px;"><div style="display:flex;justify-content:space-between;align-items:center;"><div><div style="font-size:14px;font-weight:600;color:#0f172a;">Technical Integration Assessment</div><div style="font-size:12px;color:#94a3b8;margin-top:2px;">API capabilities, payment methods, certification</div></div><span style="background:{'#d1fae5' if p.get('integration_ready') else '#f8fafc'};color:{'#065f46' if p.get('integration_ready') else '#64748b'};font-size:11px;font-weight:600;padding:3px 12px;border-radius:20px;">{'COMPLETED' if p.get('integration_ready') else 'NOT STARTED'}</span></div></div>
  <div style="border:1px solid #e2e8f0;border-radius:8px;padding:16px;"><div style="display:flex;justify-content:space-between;align-items:center;"><div><div style="font-size:14px;font-weight:600;color:#0f172a;">Commercial Terms Review</div><div style="font-size:12px;color:#94a3b8;margin-top:2px;">Pricing, SLAs, revenue share, exclusivity</div></div><span style="background:{'#d1fae5' if p.get('revshare') else '#f8fafc'};color:{'#065f46' if p.get('revshare') else '#64748b'};font-size:11px;font-weight:600;padding:3px 12px;border-radius:20px;">{'COMPLETED' if p.get('revshare') else 'NOT STARTED'}</span></div></div>
</div>""", unsafe_allow_html=True)

    # ── Audit Log ─────────────────────────────────────────────────────────────
    elif active_tab == "Audit Log":
        manager = p.get("manager", "Team")
        st.markdown(f"""
<div style="border:1px solid #e2e8f0;border-radius:10px;padding:24px;margin-top:16px;background:#fff;">
  <div style="font-size:18px;font-weight:700;color:#0f172a;margin-bottom:24px;">Audit Log</div>
  <div style="border-left:2px solid #e2e8f0;margin-left:8px;padding-left:24px;">
    <div style="position:relative;padding-bottom:24px;"><div style="position:absolute;left:-31px;top:2px;width:12px;height:12px;border-radius:50%;background:#4F46E5;border:2px solid #fff;"></div><div style="font-size:13px;font-weight:600;color:#0f172a;">Status updated to {p['status']}</div><div style="font-size:12px;color:#94a3b8;">{manager} &middot; March 2026</div></div>
    <div style="position:relative;padding-bottom:24px;"><div style="position:absolute;left:-31px;top:2px;width:12px;height:12px;border-radius:50%;background:#e2e8f0;border:2px solid #fff;"></div><div style="font-size:13px;font-weight:600;color:#0f172a;">{'NDA signed and uploaded' if p.get('nda') else 'NDA requested'}</div><div style="font-size:12px;color:#94a3b8;">{manager} &middot; February 2026</div></div>
    <div style="position:relative;padding-bottom:24px;"><div style="position:absolute;left:-31px;top:2px;width:12px;height:12px;border-radius:50%;background:#e2e8f0;border:2px solid #fff;"></div><div style="font-size:13px;font-weight:600;color:#0f172a;">Partner account created</div><div style="font-size:12px;color:#94a3b8;">{manager} &middot; January 2026</div></div>
  </div>
</div>""", unsafe_allow_html=True)


def _build_partners_df():
    """Build a DataFrame from PARTNERS_DATA for export."""
    rows = []
    for p in PARTNERS_DATA:
        manager = p.get("manager", "—")
        if manager == "nan": manager = "—"
        country = p.get("country", "")
        if country == "nan": country = "—"
        rows.append({
            "Partner Name": p["name"],
            "Category": p.get("cat", ""),
            "Region": p.get("region", "—"),
            "Country": country,
            "Status": p["status"],
            "Tier": p.get("tier", "—"),
            "Manager": manager,
        })
    return pd.DataFrame(rows)


def _export_excel(df):
    import io
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as XlImage
    from openpyxl.styles import Font, PatternFill, Alignment
    wb = Workbook()
    ws = wb.active
    ws.title = "Partner Portfolio"
    # Logo
    logo_path = _os.path.join(_BASE, "Yuno logo.png")
    if _os.path.exists(logo_path):
        img = XlImage(logo_path)
        img.width = 150
        img.height = 150
        ws.add_image(img, "A1")
    # Title below logo
    ws["A9"] = "Yuno — Partner Portfolio Export"
    ws["A9"].font = Font(name="Helvetica", size=14, bold=True, color="4F46E5")
    ws["A10"] = f"Generated on {pd.Timestamp.now().strftime('%B %d, %Y')}"
    ws["A10"].font = Font(name="Helvetica", size=10, color="6B7280")
    # Header row
    header_row = 12
    header_fill = PatternFill(start_color="4F46E5", end_color="4F46E5", fill_type="solid")
    header_font = Font(name="Helvetica", size=10, bold=True, color="FFFFFF")
    for c_idx, col in enumerate(df.columns, 1):
        cell = ws.cell(row=header_row, column=c_idx, value=col)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center")
    # Data rows
    data_font = Font(name="Helvetica", size=9, color="1C1433")
    for r_idx, row in enumerate(df.itertuples(index=False), header_row + 1):
        for c_idx, val in enumerate(row, 1):
            cell = ws.cell(row=r_idx, column=c_idx, value=val)
            cell.font = data_font
    # Auto-width columns
    for c_idx, col in enumerate(df.columns, 1):
        max_len = max(len(str(col)), df[col].astype(str).str.len().max())
        ws.column_dimensions[ws.cell(row=header_row, column=c_idx).column_letter].width = min(max_len + 4, 30)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _export_pdf(df):
    from fpdf import FPDF
    pdf = FPDF(orientation="L", format="A4")
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_fill_color(255, 255, 255)
    # Yuno logo
    logo_path = _os.path.join(_BASE, "Yuno logo.png")
    if _os.path.exists(logo_path):
        pdf.image(logo_path, x=10, y=8, w=35, h=35)
    # Title next to logo
    pdf.set_xy(50, 14)
    pdf.set_font("Helvetica", "B", 20)
    pdf.set_text_color(62, 79, 224)
    pdf.cell(0, 10, "Yuno", ln=False)
    pdf.set_xy(50, 24)
    pdf.set_font("Helvetica", "B", 14)
    pdf.set_text_color(28, 20, 51)
    pdf.cell(0, 8, "Partner Portfolio Export", ln=True)
    pdf.set_xy(50, 33)
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(100, 116, 139)
    pdf.cell(0, 6, f"Generated on {pd.Timestamp.now().strftime('%B %d, %Y')}", ln=True)
    pdf.ln(8)
    # Table
    cols = list(df.columns)
    n_cols = len(cols)
    avail_w = 277
    col_w = avail_w / n_cols
    col_widths = [col_w] * n_cols
    # Header
    pdf.set_font("Helvetica", "B", 8)
    pdf.set_fill_color(62, 79, 224)
    pdf.set_text_color(255, 255, 255)
    for i, col in enumerate(cols):
        pdf.cell(col_widths[i], 7, col, border=0, fill=True)
    pdf.ln()
    # Rows
    pdf.set_font("Helvetica", "", 7)
    for r_idx, (_, row) in enumerate(df.iterrows()):
        if r_idx % 2 == 1:
            pdf.set_fill_color(248, 250, 252)
        else:
            pdf.set_fill_color(255, 255, 255)
        pdf.set_text_color(28, 20, 51)
        for i, col in enumerate(cols):
            val = str(row[col]) if pd.notna(row[col]) else ""
            val = val.replace("\u2014", "-").replace("\u2013", "-").replace("\u2018","'").replace("\u2019","'").replace("\u201c",'"').replace("\u201d",'"')
            pdf.cell(col_widths[i], 6, val[:35], border=0, fill=True)
        pdf.ln()
    # Footer line
    pdf.set_y(-20)
    pdf.set_font("Helvetica", "", 7)
    pdf.set_text_color(148, 163, 184)
    pdf.cell(0, 5, "Confidential - Yuno Partner Portal", align="C")
    return bytes(pdf.output())


def _export_pptx(df):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    import io
    # Use Yuno template if available
    tpl_path = _os.path.join(_BASE, "data", "template.pptx")
    if _os.path.exists(tpl_path):
        prs = Presentation(tpl_path)
        # Remove all existing slides (keep template layouts/master)
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    # Find title layout (first with TITLE placeholder) or fallback to first
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[0]
    for sl in prs.slide_layouts:
        if "title" in sl.name.lower():
            title_layout = sl
            break
    for sl in prs.slide_layouts:
        if "content" in sl.name.lower():
            content_layout = sl
            break
    # Brand colors from template theme
    brand_blue = RGBColor(0x3E, 0x4F, 0xE0)
    dark = RGBColor(0x28, 0x2A, 0x30)
    white = RGBColor(0xFF, 0xFF, 0xFF)
    gray = RGBColor(0x6B, 0x72, 0x80)
    # Title slide
    slide = prs.slides.add_slide(title_layout)
    for ph in slide.placeholders:
        if ph.placeholder_format.type is not None:
            if ph.placeholder_format.idx == 0:
                ph.text = "Partner Portfolio"
                for p in ph.text_frame.paragraphs:
                    p.font.size = Pt(36)
                    p.font.bold = True
                    p.font.color.rgb = dark
                    p.font.name = "Arial"
            elif ph.placeholder_format.idx == 1:
                ph.text = f"Generated on {pd.Timestamp.now().strftime('%B %d, %Y')}"
                for p in ph.text_frame.paragraphs:
                    p.font.size = Pt(16)
                    p.font.color.rgb = gray
                    p.font.name = "Arial"
    # Data slides
    cols = list(df.columns)
    chunk_size = 22
    for start in range(0, len(df), chunk_size):
        chunk = df.iloc[start:start + chunk_size]
        slide_n = prs.slides.add_slide(content_layout)
        # Clear placeholders
        for ph in slide_n.placeholders:
            ph.text = ""
        # Add table
        n_rows = len(chunk) + 1
        tbl_shape = slide_n.shapes.add_table(n_rows, len(cols), Inches(0.4), Inches(0.9), Inches(12.5), Inches(6.0))
        tbl = tbl_shape.table
        # Header
        for i, col in enumerate(cols):
            cell = tbl.cell(0, i)
            cell.text = col
            cell.fill.solid()
            cell.fill.fore_color.rgb = brand_blue
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.bold = True
                p.font.color.rgb = white
                p.font.name = "Arial"
        # Rows
        for r_idx in range(len(chunk)):
            for c_idx, col in enumerate(cols):
                cell = tbl.cell(r_idx + 1, c_idx)
                val = str(chunk.iloc[r_idx][col]) if pd.notna(chunk.iloc[r_idx][col]) else ""
                cell.text = val
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC) if r_idx % 2 == 1 else white
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(8)
                    p.font.color.rgb = dark
                    p.font.name = "Arial"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def show_partners():
    if st.session_state.selected_partner is not None:
        sel = st.session_state.selected_partner
        partner = next((p for p in PARTNERS_DATA if p["name"] == sel), None)
        if partner:
            _partner_detail(partner)
            return
        else:
            st.session_state.selected_partner = None

    # ── Build connector data for HTML component (cached) ────────────────────
    _connector_data, total, sot_live_ct, countries_ct = _build_connector_data()

    import json as _json
    _data_json = _json.dumps(_connector_data).replace("'","\\'").replace("</","<\\/")

    # Read the HTML template and inject data
    _html_file = _os.path.join(_BASE, "yuno_connectors.html")
    with open(_html_file, "r") as _f:
        _html_content = _f.read()

    # Replace the hardcoded DATA array with our dynamic data
    _html_content = _html_content.replace(
        "var DATA = [",
        f"var DATA = {_data_json}; var _ORIGINAL_DATA = ["
    )
    # Update stat numbers
    _html_content = _html_content.replace(">460<", f">{total}<")
    _html_content = _html_content.replace(">97<", f">{sot_live_ct}<")
    _html_content = _html_content.replace(">189<", f">{countries_ct}<")
    _html_content = _html_content.replace("var YUNO_LOGO_B64='';", f"var YUNO_LOGO_B64='{_LOGO_B64}';")

    components.html(_html_content, height=2400, scrolling=True)
    return


@st.cache_data
def _build_connector_data():
    total = len(PARTNERS_DATA)
    live_all = sum(1 for p in PARTNERS_DATA if p["status"] == "Live")
    countries_ct = len(set(p.get("country","") for p in PARTNERS_DATA if p.get("country","") and p.get("country","") != "nan"))
    sot_live_ct = len(_SOT_DF[_SOT_DF["Live/NonLive Partner/Contract signed"] == "Live"]["PROVIDER_NAME"].unique()) if len(_SOT_DF) > 0 else live_all

    import json as _json
    _connector_data = []
    for p in PARTNERS_DATA:
        enrich = _PARTNER_ENRICHMENT.get(p["name"].upper().replace("(","").replace(")","").strip())
        manager = p.get("manager","—")
        if manager == "nan": manager = "—"
        # Payment methods from SOT
        sot_pms = []
        if len(_SOT_DF) > 0:
            pm_rows = _SOT_DF[(_SOT_DF["PROVIDER_NAME"].str.upper() == p["name"].upper()) & (_SOT_DF["Live/NonLive Partner/Contract signed"] == "Live")]
            sot_pms = sorted(pm_rows["PAYMENT_METHOD_TYPE"].dropna().unique().tolist())[:8]
        country = p.get("country","")
        if country == "nan": country = "—"
        # SOT capabilities lookup
        _sot_caps = {}
        if len(_SOT_DF) > 0:
            _p_sot = _SOT_DF[_SOT_DF["PROVIDER_NAME"].str.upper() == p["name"].upper()]
            _sot_caps = {
                "gaming": bool(any(_p_sot["ACCEPTS_GAMING"].isin([True, 1, 1.0]))),
                "gambling": bool(any(_p_sot["ACCEPTS_GAMBLING"].isin([True, 1, 1.0]))),
                "crypto": bool(any(_p_sot["ACCEPTS_CRYPTO"].isin([True, 1, 1.0]))),
                "recurring": bool(any(_p_sot["SUPPORTS_RECURRING_PAYMENTS"].isin([True, 1, 1.0]))),
                "tokenization": bool(any(_p_sot["SUPPORTS_TOKENIZATION"].isin([True, 1, 1.0]))),
                "payouts": bool(any(_p_sot["SUPPORTS_PAYOUTS"].isin([True, 1, 1.0]))),
                "installments": bool(any(_p_sot["SUPPORTS_INSTALLMENTS"].isin([True, 1, 1.0]))),
                "airlines": bool(any(_p_sot["ACCEPTS_AIRLINES"].isin([1, 1.0]))),
                "regions": sorted(_p_sot["REGION"].dropna().unique().tolist()),
                "countries": sorted(_p_sot["COUNTRY_ISO"].dropna().unique().tolist()),
            }

        # External contacts from partner files
        _ext_key = p["name"].upper().replace("(","").replace(")","").strip()
        _ext = _EXT_CONTACTS.get(_ext_key, {})
        _comm_name = _ext.get("commercial_contact","") or (enrich["commercial_contact"] if enrich else "") or "TBD"
        _comm_email = _ext.get("commercial_email","") or "TBD"
        _comm_phone = _ext.get("commercial_phone","") or "TBD"
        _tech_name = _ext.get("technical_contact","") or (enrich["technical_contact"] if enrich else "") or "TBD"
        _tech_email = _ext.get("technical_email","") or "TBD"
        _tech_phone = _ext.get("technical_phone","") or "TBD"
        # Verticals from external data (Source of Truth features)
        _ext_verts = _ext.get("verticals","")
        _verticals_list = [v.strip() for v in _ext_verts.split(",") if v.strip()] if _ext_verts else (enrich["verticals"].split(", ") if enrich and enrich.get("verticals") else ["TBD"])

        _connector_data.append({
            "name": p["name"], "type": p["cat"], "region": p.get("region","—"), "country": country, "status": p["status"],
            "caps": _sot_caps,
            "contacts": [
                {"name": _comm_name, "role": "Commercial", "email": _comm_email, "phone": _comm_phone, "color": "#3B82F6"},
                {"name": _tech_name, "role": "Technical", "email": _tech_email, "phone": _tech_phone, "color": "#10B981"},
                {"name": manager if manager != "—" else "Partnerships Team", "role": "Escalation", "email": "", "phone": "", "color": "#EF4444"},
            ],
            "verticals": _verticals_list,
            "currencies": ["USD","Local"],
            "auth": {"type": "Full Auth", "settlement": "T+1", "settleCurrency": "Local / USD", "chargebackWindow": "90 days"},
            "onboarding": {"kyb": bool(p.get("nda")), "integration": enrich["onboarding"].split(",")[0] if enrich else "API", "timeline": enrich["onboarding"].split(",")[-1].strip() if enrich else "2-4 weeks", "localEntity": False, "jurisdiction": "—", "entityType": "Contact partnerships"},
            "pricing": {"standard": enrich["avg_pricing"] if enrich else "TBD", "discount": enrich["discount_rate"] if enrich else "TBD", "discountWhen": enrich["discount_condition"] if enrich else "TBD", "mdr": "Blended", "fx": "TBD", "minimums": "TBD"},
            "paymentMethods": sot_pms if sot_pms else ["Card"],
            "health": {"score": 85 if p["status"]=="Live" else 55 if "Agreement" in p["status"] else 30, "label": enrich["health"] if enrich else ("Live" if p["status"]=="Live" else "Pipeline")},
            "tier": p.get("tier","—").replace("Partners","").replace("(Global)","").replace("Strategic Partners: Very Important","Strategic").strip().rstrip(":").strip() or "—",
            "strategic": bool(p.get("strategic")),
            "intReady": bool(p.get("int_ready")),
            "intUsed": bool(p.get("int_used")),
            "nda": bool(p.get("nda")),
            "revshare": enrich["revshare"] if enrich else ("Yes" if p.get("revshare") else "No"),
        })
    return _connector_data, total, sot_live_ct, countries_ct

    # ── OLD CODE (kept for reference, never executes) ─────────────────────────
    st.markdown(f"""
<div style="margin-bottom:20px;">
  <div style="font-size:28px;font-weight:800;color:#0f172a;letter-spacing:-0.8px;margin-bottom:4px;">Yuno Connectors</div>
  <div style="font-size:13px;color:#94a3b8;">Search, filter and explore Yuno's payment partner ecosystem. Find the right connector for your merchant.</div>
</div>""", unsafe_allow_html=True)

    # ── Stat cards ────────────────────────────────────────────────────────────
    st.markdown(f"""
<div style="display:grid;grid-template-columns:1fr 1fr 1fr 1fr;gap:12px;margin-bottom:20px;">
  <div style="background:#fff;border:1px solid #E2E8F0;border-radius:14px;padding:20px;text-align:center;"><div style="font-size:28px;font-weight:800;color:#5B5FDE;letter-spacing:-1px;">{total}</div><div style="font-size:11px;color:#94a3b8;font-weight:600;text-transform:uppercase;letter-spacing:0.05em;margin-top:4px;">Total Connectors</div></div>
  <div style="background:#fff;border:1px solid #E2E8F0;border-radius:14px;padding:20px;text-align:center;"><div style="font-size:28px;font-weight:800;color:#10B981;letter-spacing:-1px;">{sot_live_ct}</div><div style="font-size:11px;color:#94a3b8;font-weight:600;text-transform:uppercase;letter-spacing:0.05em;margin-top:4px;">Live & Active</div></div>
  <div style="background:#fff;border:1px solid #E2E8F0;border-radius:14px;padding:20px;text-align:center;"><div style="font-size:28px;font-weight:800;color:#3B82F6;letter-spacing:-1px;">{countries_ct}</div><div style="font-size:11px;color:#94a3b8;font-weight:600;text-transform:uppercase;letter-spacing:0.05em;margin-top:4px;">Countries Covered</div></div>
  <div style="background:#fff;border:1px solid #E2E8F0;border-radius:14px;padding:20px;text-align:center;"><div style="font-size:28px;font-weight:800;color:#F59E0B;letter-spacing:-1px;">89.4%</div><div style="font-size:11px;color:#94a3b8;font-weight:600;text-transform:uppercase;letter-spacing:0.05em;margin-top:4px;">Avg Approval Rate</div></div>
</div>""", unsafe_allow_html=True)

    # ── AI Search box ─────────────────────────────────────────────────────────
    st.markdown("""
<div style="background:linear-gradient(160deg,#0a0e1a,#111827);border-radius:14px;padding:24px 28px 18px;margin-bottom:0;position:relative;overflow:hidden;">
  <div style="position:absolute;top:-40px;right:-40px;width:200px;height:200px;border-radius:50%;background:radial-gradient(circle,rgba(91,95,222,0.1),transparent 65%);pointer-events:none;"></div>
  <div style="position:relative;z-index:1;">
    <div style="display:flex;align-items:center;gap:8px;margin-bottom:8px;">
      <span style="font-size:16px;">✦</span>
      <span style="font-size:10px;font-weight:700;letter-spacing:1.2px;text-transform:uppercase;color:#818cf8;">Smart Search</span>
    </div>
    <div style="font-size:12px;color:#64748b;line-height:1.6;margin-bottom:6px;">Search by name, type, country, or describe what you need — even with typos. We'll find the closest match.</div>
  </div>
</div>""", unsafe_allow_html=True)

    st.markdown('<div style="background:#fff;border:1px solid #E2E8F0;border-radius:0 0 14px 14px;padding:14px 22px 8px;margin-bottom:4px;"></div>', unsafe_allow_html=True)
    search = st.text_input("Search", placeholder="Describe what you're looking for...", key="partner_search_input", label_visibility="collapsed")

    # Quick search prompts
    _QS_PROMPTS = [
        "Copa Airlines expanding into the U.S.",
        "Uber needs local processing in Brazil",
        "iGaming operator entering Turkey",
        "Crypto exchange launching in UAE",
        "SaaS recurring payments in Saudi Arabia",
    ]
    st.markdown('<div style="font-size:10px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:1px;margin:4px 0 6px;">Quick Search</div>', unsafe_allow_html=True)
    qs_cols = st.columns(len(_QS_PROMPTS))
    for qi, prompt in enumerate(_QS_PROMPTS):
        if qs_cols[qi].button(prompt, key=f"qs_conn_{qi}", use_container_width=True):
            st.session_state["qs_search"] = prompt
            st.rerun()

    # Apply quick search if set
    if st.session_state.get("qs_search"):
        if not search:
            search = st.session_state["qs_search"]
        else:
            st.session_state["qs_search"] = ""

    # Style quick search buttons
    components.html("""<script>
setTimeout(function(){
  var doc = window.parent.document;
  var prompts = ['Copa Airlines','Uber needs','iGaming','Crypto exchange','SaaS recurring'];
  doc.querySelectorAll('button').forEach(function(btn){
    if(prompts.some(function(p){return btn.textContent.trim().startsWith(p);})){
      btn.style.cssText='background:#EEF2FF !important;border:1px solid #C7D2FE !important;border-radius:10px !important;padding:10px 12px !important;font-size:12px !important;color:#5B5FDE !important;font-weight:500 !important;cursor:pointer !important;box-shadow:none !important;text-align:left !important;line-height:1.3 !important;min-height:0 !important;';
    }
  });
},200);
</script>""", height=0)

    # ── Filters (collapsible) ─────────────────────────────────────────────────
    with st.expander("Filters & Export", expanded=False):
        fc = st.columns([1, 1, 1, 1])
        with fc[0]:
            type_filter = st.selectbox("Type", ["Type","PSP","Acquirer","APM","Fraud","BaaS","Other"], key="pd_type_filter", label_visibility="collapsed")
        with fc[1]:
            _pd_regions = sorted(set(p.get("region","") for p in PARTNERS_DATA if p.get("region","")))
            region_filter = st.selectbox("Region", ["Region"] + _pd_regions, key="pd_region_filter", label_visibility="collapsed")
        with fc[2]:
            _pd_countries = sorted(set(p.get("country","") for p in PARTNERS_DATA if p.get("country","") and p.get("country","") != "nan"))
            country_filter = st.selectbox("Country", ["Country"] + _pd_countries[:40], key="pd_country_filter", label_visibility="collapsed")
        with fc[3]:
            stage_filter = st.selectbox("Status", ["Status","Live","Agreement Signed","Agreement Review","Negotiation","Prospect"], key="pd_stage_filter", label_visibility="collapsed")
        fc2 = st.columns([1, 1, 1, 1])
        with fc2[0]:
            _pm_options = ["Payment Method"]
            if len(_SOT_DF) > 0:
                _pm_options += sorted(_SOT_DF["PAYMENT_METHOD_TYPE"].dropna().unique().tolist())[:30]
            pm_filter = st.selectbox("PM", _pm_options, key="pd_pm_filter", label_visibility="collapsed")
        with fc2[1]:
            vert_filter = st.selectbox("Vertical", ["Vertical","High Risk","Gambling","Gaming","Forex","Crypto","Airlines"], key="pd_vert_filter", label_visibility="collapsed")
        with fc2[2]:
            sort_by = st.selectbox("Sort", ["Sort by","Popularity","Cost effective","High performance"], key="pd_sort", label_visibility="collapsed")
        with fc2[3]:
            export_fmt = st.selectbox("Export", ["Export","CSV","PDF","Excel","PPTX"], key="pd_export", label_visibility="collapsed")
        if export_fmt != "Export":
            st.info(f"Generating {export_fmt}... (production: file download with merchant logo)")
            st.session_state["pd_export"] = "Export"

    # ── Filter logic ──────────────────────────────────────────────────────────
    filt = list(PARTNERS_DATA)
    if type_filter != "Type":
        filt = [p for p in filt if p["cat"] == type_filter]
    if stage_filter != "Status":
        filt = [p for p in filt if p["status"] == stage_filter]
    if region_filter != "Region":
        filt = [p for p in filt if p.get("region") == region_filter]
    if country_filter != "Country":
        filt = [p for p in filt if p.get("country") == country_filter]
    # ── AI Smart Search ────────────────────────────────────────────────────────
    _ai_rec = None  # Will hold recommendation if complex query detected
    if search:
        sl = search.lower()
        _search_aliases = {"u.s.":"united states","u.s":"united states","usa":"united states","us":"united states",
                           "uk":"united kingdom","uae":"united arab emirates","emirates":"united arab emirates",
                           "brasil":"brazil","ksa":"saudi arabia","saudi":"saudi arabia","dubai":"united arab emirates",
                           "igaming":"gambling","betting":"gambling","casino":"gambling",
                           "crypto":"crypto","airline":"airlines","aviation":"airlines","flight":"airlines",
                           "forex":"forex","trading":"forex","saas":"psp","recurring":"psp"}
        _stop_words = {"the","a","an","in","into","for","to","and","or","of","with","is","by","on","at","from",
                        "needs","need","wants","want","expanding","entering","launching","looking","operator",
                        "exchange","company","processing","local","payments","payment","top","best","find",
                        "can","you","provide","short","list","partners","partner","high","low","rates","rate",
                        "approvals","approval","please","help","me","what","are","which","should","recommend",
                        "we","our","their","them","this","that","these","those"}
        # Detect intent qualifiers
        wants_low_cost = any(w in sl for w in ["low rate","low cost","cheap","cost effective","affordable","discount"])
        wants_high_perf = any(w in sl for w in ["high approval","high performance","best approval","top perform","reliable"])
        wants_live = any(w in sl for w in ["live","active","ready","available"])
        wants_acquirer = "acquir" in sl
        wants_psp = "psp" in sl or "gateway" in sl or "aggregat" in sl
        wants_apm = "apm" in sl or "wallet" in sl or "local method" in sl
        # Detect country from aliases
        detected_country = ""
        for alias, full in sorted(_search_aliases.items(), key=lambda x: -len(x[0])):
            if alias in sl:
                detected_country = full
                break
        # Also check raw country names in partner data
        if not detected_country:
            for p in PARTNERS_DATA:
                cname = p.get("country","").lower()
                if cname and cname != "nan" and cname in sl:
                    detected_country = p.get("country","")
                    break
        # Expand aliases
        expanded = sl
        for alias, replacement in _search_aliases.items():
            if alias in sl:
                expanded = expanded + " " + replacement
        tokens = [t.strip(".,!?") for t in expanded.split() if t.strip(".,!?").lower() not in _stop_words and len(t.strip(".,!?")) > 1]
        if not tokens:
            tokens = [sl]
        def _smart_match(p, tokens):
            searchable = f"{p['name']} {p['cat']} {p.get('region','')} {p.get('country','')} {p['status']} {p.get('offering_raw','')} {p.get('tier','')} {p.get('manager','')}".lower()
            for t in tokens:
                if t in searchable:
                    return True
                if len(t) >= 3 and any(w.startswith(t[:3]) for w in searchable.split()):
                    return True
            return False
        filt = [p for p in filt if _smart_match(p, tokens)]
        # Apply intent filters
        if wants_acquirer:
            filt = [p for p in filt if p["cat"] == "Acquirer"]
        elif wants_psp:
            filt = [p for p in filt if p["cat"] == "PSP"]
        elif wants_apm:
            filt = [p for p in filt if p["cat"] == "APM"]
        if wants_live:
            filt = [p for p in filt if p["status"] == "Live"]
        if detected_country:
            country_filt = [p for p in filt if detected_country.lower() in p.get("country","").lower() or detected_country.lower() in p.get("region","").lower()]
            if country_filt:
                filt = country_filt
        # Sort by intent
        if wants_low_cost:
            filt = sorted(filt, key=lambda x: 0 if x.get("strategic") else 1)
        if wants_high_perf:
            filt = sorted(filt, key=lambda x: (0 if x.get("int_used") else 1, 0 if x["status"]=="Live" else 1))
        # Build AI recommendation if this looks like a complex query (>4 words)
        if len(search.split()) > 3 and filt:
            top = filt[:5]
            top_names = [p["name"] for p in top]
            rec_type = "acquirers" if wants_acquirer else "PSPs" if wants_psp else "APMs" if wants_apm else "connectors"
            rec_country = detected_country if detected_country else "your target market"
            rec_why = []
            if wants_live: rec_why.append("live and processing")
            if wants_high_perf: rec_why.append("high approval rates")
            if wants_low_cost: rec_why.append("competitive pricing")
            if not rec_why: rec_why.append("best overall fit")
            _ai_rec = {"names": top_names, "type": rec_type, "country": rec_country, "why": rec_why, "count": len(filt)}
    if vert_filter != "Vertical" and len(_SOT_DF) > 0:
        vert_col = _VERTICAL_COLS.get(vert_filter)
        if vert_col and vert_col in _SOT_DF.columns:
            vert_providers = set(_SOT_DF[(_SOT_DF[vert_col] == True) | (_SOT_DF[vert_col] == 1)]["PROVIDER_NAME"].str.upper().unique())
            filt = [p for p in filt if p["name"].upper() in vert_providers]
    if pm_filter != "Payment Method" and len(_SOT_DF) > 0:
        pm_providers = set(_SOT_DF[_SOT_DF["PAYMENT_METHOD_TYPE"] == pm_filter]["PROVIDER_NAME"].str.upper().unique())
        filt = [p for p in filt if p["name"].upper() in pm_providers]
    if sort_by == "Popularity":
        filt = sorted(filt, key=lambda x: 0 if x["status"] == "Live" else 1)
    elif sort_by == "High performance":
        filt = sorted(filt, key=lambda x: (0 if x.get("int_used") else 1, 0 if x.get("revshare_active") else 1))

    # ── AI Recommendation card ──────────────────────────────────────────────
    _STATUS_TAG = {"Live":("#DCFCE7","#065F46"),"Agreement Signed":("#DBEAFE","#1D4ED8"),"Agreement Review":("#FEF3C7","#B45309"),"Negotiation":("#FEF3C7","#B45309"),"Prospect":("#F1F5F9","#475569"),"Lost":("#FEE2E2","#DC2626")}
    _TYPE_TAG = {"PSP":("#DBEAFE","#1D4ED8"),"Acquirer":("#DBEAFE","#1D4ED8"),"APM":("#F3E8FF","#6D28D9"),"Fraud":("#FEE2E2","#DC2626"),"BaaS":("#FEF3C7","#B45309"),"Other":("#F1F5F9","#475569")}

    if _ai_rec:
        r = _ai_rec
        top_pills = "".join(f'<span style="font-size:12px;font-weight:700;padding:4px 14px;border-radius:8px;background:#EEF2FF;color:#4F46E5;border:0.5px solid #C7D2FE;">{n}</span>' for n in r["names"])
        why_text = ", ".join(r["why"])
        # Get enrichment data for top picks
        detail_rows = ""
        for pname in r["names"]:
            enrich = _PARTNER_ENRICHMENT.get(pname.upper().replace("(","").replace(")","").strip())
            pricing = enrich["avg_pricing"] if enrich else "—"
            discount = enrich["discount_rate"] if enrich else "—"
            health = enrich["health"] if enrich else "—"
            detail_rows += (f'<tr style="border-bottom:0.5px solid #EEF2FF;">'
                            f'<td style="padding:8px 12px;font-size:13px;font-weight:700;color:#0f172a;">{pname}</td>'
                            f'<td style="padding:8px 12px;font-size:12px;color:#64748b;">{pricing}</td>'
                            f'<td style="padding:8px 12px;font-size:12px;font-weight:700;color:#065F46;">{discount}</td>'
                            f'<td style="padding:8px 12px;font-size:12px;color:#64748b;">{health}</td>'
                            f'</tr>')
        th_r = "padding:6px 12px;font-size:9px;font-weight:700;letter-spacing:0.08em;text-transform:uppercase;color:#94a3b8;border-bottom:1px solid #EEF2FF;"
        st.markdown(f"""
<div style="background:linear-gradient(135deg,#EEF2FF,#F8FAFF);border:1px solid #C7D2FE;border-radius:14px;padding:22px 24px;margin-bottom:16px;">
  <div style="display:flex;align-items:center;gap:8px;margin-bottom:10px;">
    <span style="font-size:14px;">✦</span>
    <span style="font-size:11px;font-weight:700;color:#4F46E5;text-transform:uppercase;letter-spacing:0.8px;">AI Recommendation</span>
  </div>
  <div style="font-size:14px;font-weight:600;color:#0f172a;line-height:1.6;margin-bottom:12px;">
    Based on your query, here are the <strong>top {len(r["names"])} {r["type"]}</strong> for <strong>{r["country"]}</strong> — selected for {why_text}.
  </div>
  <div style="display:flex;flex-wrap:wrap;gap:6px;margin-bottom:16px;">{top_pills}</div>
  <table style="width:100%;border-collapse:collapse;background:#fff;border-radius:8px;overflow:hidden;">
    <thead><tr><th style="{th_r}">Partner</th><th style="{th_r}">Std Rate</th><th style="{th_r}">Best Rate</th><th style="{th_r}">Health</th></tr></thead>
    <tbody>{detail_rows}</tbody>
  </table>
  <div style="margin-top:12px;font-size:11px;color:#64748b;line-height:1.5;">Click any partner below to see full details — contacts, payment methods, integration status, and onboarding requirements.</div>
</div>""", unsafe_allow_html=True)

    # ── Results table + expandable details ────────────────────────────────────
    th = "padding:10px 14px;text-align:left;font-size:10px;font-weight:700;letter-spacing:0.08em;text-transform:uppercase;color:#94A3B8;border-bottom:1px solid #E2E8F0;background:#FAFAFA;"
    rows_html = ""
    for i, p in enumerate(filt[:80]):
        s_bg, s_color = _STATUS_TAG.get(p["status"], ("#F1F5F9","#475569"))
        t_bg, t_color = _TYPE_TAG.get(p["cat"], ("#F1F5F9","#475569"))
        country = p.get("country","")
        if country == "nan": country = "—"
        bb = "border-bottom:0.5px solid #F1F5F9;" if i < min(len(filt),80)-1 else ""
        td = f"padding:11px 14px;{bb}font-size:13px;vertical-align:middle;"
        rows_html += (f'<tr onmouseover="this.style.background=\'#F8FAFC\'" onmouseout="this.style.background=\'#fff\'">'
                      f'<td style="{td}font-weight:700;color:#0f172a;">{p["name"]}</td>'
                      f'<td style="{td}"><span style="font-size:10px;font-weight:700;padding:2px 8px;border-radius:5px;background:{t_bg};color:{t_color};">{p["cat"]}</span></td>'
                      f'<td style="{td}color:#64748b;">{p.get("region","—")}</td>'
                      f'<td style="{td}color:#64748b;">{country}</td>'
                      f'<td style="{td}"><span style="font-size:10px;font-weight:700;padding:2px 8px;border-radius:5px;background:{s_bg};color:{s_color};">{p["status"]}</span></td>'
                      f'</tr>')

    st.markdown(f"""
<div style="display:flex;align-items:center;gap:8px;margin:8px 0 10px;"><span style="font-size:14px;font-weight:700;color:#0f172a;">{len(filt)}</span><span style="font-size:11px;color:#94a3b8;">connectors found</span></div>
<div style="background:#fff;border:0.5px solid #E2E8F0;border-radius:12px;overflow:hidden;margin-bottom:14px;">
  <table style="width:100%;border-collapse:collapse;">
    <thead><tr>
      <th style="{th}">Connector Name</th>
      <th style="{th}">Type</th>
      <th style="{th}">Region</th>
      <th style="{th}">Country</th>
      <th style="{th}">Status</th>
    </tr></thead>
    <tbody>{rows_html}</tbody>
  </table>
</div>""", unsafe_allow_html=True)

    # Expandable detail per partner
    st.markdown('<div style="font-size:11px;color:#94a3b8;margin-bottom:6px;">Click any connector below to view full details</div>', unsafe_allow_html=True)

    for i, p in enumerate(filt[:80]):
        s_bg, s_color = _STATUS_TAG.get(p["status"], ("#F1F5F9","#475569"))
        t_bg, t_color = _TYPE_TAG.get(p["cat"], ("#F1F5F9","#475569"))
        country = p.get("country","")
        if country == "nan": country = "—"

        with st.expander(p["name"], expanded=False):
            manager = p.get("manager","—")
            if manager == "nan": manager = "—"
            tier_short = p.get("tier","").replace("Partners","").replace("(Global)","").replace("Strategic Partners: Very Important","Strategic").strip().rstrip(":").strip()
            if not tier_short or tier_short == "nan": tier_short = "—"

            health_map = {"Live":85,"Agreement Signed":55,"Agreement Review":45,"Negotiation":40,"Prospect":25,"Lost":10}
            health = health_map.get(p["status"], 30)
            if p.get("int_used"): health = min(health + 10, 98)
            if p.get("revshare_active"): health = min(health + 5, 99)
            h_color = "#10B981" if health >= 70 else "#F59E0B" if health >= 40 else "#EF4444"
            h_text = "#065F46" if health >= 70 else "#B45309" if health >= 40 else "#DC2626"

            # Enrichment data lookup
            enrich = _PARTNER_ENRICHMENT.get(p["name"].upper().replace("(","").replace(")","").strip())
            _ext_d = _EXT_CONTACTS.get(p["name"].upper().replace("(","").replace(")","").strip(), {})
            comm_contact = _ext_d.get("commercial_contact","") or (enrich["commercial_contact"] if enrich else "") or "TBD"
            tech_contact = _ext_d.get("technical_contact","") or (enrich["technical_contact"] if enrich else "") or "TBD"
            _ext_verts_d = _ext_d.get("verticals","")
            verticals = _ext_verts_d if _ext_verts_d else (enrich["verticals"] if enrich else "TBD")
            pricing = enrich["avg_pricing"] if enrich else "TBD"
            discount = enrich["discount_rate"] if enrich else "TBD"
            discount_cond = enrich["discount_condition"] if enrich else "TBD"
            revshare_status = enrich["revshare"] if enrich else ("Yes" if p.get("revshare") else "TBD")
            onboarding = enrich["onboarding"] if enrich else "TBD"
            merch_live = enrich["merchants_live"] if enrich else 0
            merch_names = enrich["merchants"] if enrich else "TBD"
            p_health_label = enrich["health"] if enrich else ("Good" if health >= 50 else "New")

            # Payment methods from SOT
            sot_pms = []
            if len(_SOT_DF) > 0:
                pm_rows = _SOT_DF[(_SOT_DF["PROVIDER_NAME"].str.upper() == p["name"].upper()) & (_SOT_DF["Live/NonLive Partner/Contract signed"] == "Live")]
                sot_pms = sorted(pm_rows["PAYMENT_METHOD_TYPE"].dropna().unique().tolist())[:10]
            pm_display = sot_pms if sot_pms else ["Card"]

            # Tick/cross helper
            def _tick(val):
                return '<span style="width:18px;height:18px;border-radius:50%;background:#DCFCE7;display:inline-flex;align-items:center;justify-content:center;font-size:11px;color:#065F46;">&#10003;</span>' if val else '<span style="width:18px;height:18px;border-radius:50%;background:#F1F5F9;display:inline-flex;align-items:center;justify-content:center;font-size:11px;color:#CBD5E1;">&#10007;</span>'

            # Panel container
            st.markdown(f'<div style="background:#FAFBFF;border:1px solid #5B5FDE;border-top:none;border-radius:0 0 12px 12px;padding:20px 22px;margin:0 0 6px;">', unsafe_allow_html=True)

            # Health + Status bar
            h_label_bg = "#DCFCE7" if health >= 70 else "#FEF3C7" if health >= 40 else "#FEE2E2"
            h_label_color = "#065F46" if health >= 70 else "#B45309" if health >= 40 else "#DC2626"
            st.markdown(f'<div style="display:flex;align-items:center;gap:12px;margin-bottom:16px;padding-bottom:14px;border-bottom:1px solid #EEF2FF;">'
                        f'<div style="display:flex;align-items:center;gap:8px;"><span style="font-size:10px;font-weight:700;color:#94a3b8;text-transform:uppercase;">Health</span>'
                        f'<div style="width:120px;height:6px;border-radius:3px;background:#F1F5F9;"><div style="height:100%;border-radius:3px;width:{health}%;background:{h_color};"></div></div>'
                        f'<span style="font-size:14px;font-weight:800;color:{h_text};font-family:monospace;">{health}</span></div>'
                        f'<span style="font-size:10px;font-weight:700;padding:3px 10px;border-radius:6px;background:{h_label_bg};color:{h_label_color};">{p_health_label}</span>'
                        f'<div style="margin-left:auto;display:flex;align-items:center;gap:8px;">'
                        f'<span style="font-size:10px;color:#94a3b8;">Tier</span><span style="font-size:12px;font-weight:700;color:#0f172a;">{tier_short}</span>'
                        f'<span style="font-size:10px;color:#94a3b8;margin-left:8px;">Status</span><span style="font-size:10.5px;font-weight:700;padding:2px 8px;border-radius:6px;background:{s_bg};color:{s_color};">{p["status"]}</span>'
                        f'</div></div>', unsafe_allow_html=True)

            # Row 1: Contacts + Commercial
            col_a, col_b = st.columns(2)
            with col_a:
                st.markdown(f'<div style="background:#fff;border:1px solid #EEF2F7;border-radius:10px;padding:16px 18px;">'
                            f'<div style="font-size:10px;font-weight:700;color:#5B5FDE;text-transform:uppercase;letter-spacing:0.08em;margin-bottom:12px;">Key Contacts</div>'
                            f'<div style="display:flex;flex-direction:column;gap:10px;">'
                            f'<div style="display:flex;align-items:center;gap:10px;padding:8px 10px;background:#F8FAFC;border-radius:8px;">'
                            f'<div style="width:28px;height:28px;border-radius:50%;background:#DBEAFE;display:flex;align-items:center;justify-content:center;font-size:10px;font-weight:700;color:#1D4ED8;flex-shrink:0;">BD</div>'
                            f'<div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;">Commercial Contact</div><div style="font-size:12px;font-weight:600;color:#0f172a;">{comm_contact}</div></div></div>'
                            f'<div style="display:flex;align-items:center;gap:10px;padding:8px 10px;background:#F8FAFC;border-radius:8px;">'
                            f'<div style="width:28px;height:28px;border-radius:50%;background:#CCFBF1;display:flex;align-items:center;justify-content:center;font-size:10px;font-weight:700;color:#065F46;flex-shrink:0;">TC</div>'
                            f'<div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;">Technical Contact</div><div style="font-size:12px;font-weight:600;color:#0f172a;">{tech_contact}</div></div></div>'
                            f'<div style="display:flex;align-items:center;gap:10px;padding:8px 10px;background:#F8FAFC;border-radius:8px;">'
                            f'<div style="width:28px;height:28px;border-radius:50%;background:#FEE2E2;display:flex;align-items:center;justify-content:center;font-size:10px;font-weight:700;color:#DC2626;flex-shrink:0;">ES</div>'
                            f'<div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;">Escalation / Partner Manager</div><div style="font-size:12px;font-weight:600;color:#0f172a;">{manager if manager != "—" else "Partnerships Team"}</div></div></div>'
                            f'</div></div>', unsafe_allow_html=True)
            with col_b:
                st.markdown(f'<div style="background:#fff;border:1px solid #EEF2F7;border-radius:10px;padding:16px 18px;">'
                            f'<div style="font-size:10px;font-weight:700;color:#5B5FDE;text-transform:uppercase;letter-spacing:0.08em;margin-bottom:12px;">Pricing & Commercial</div>'
                            f'<div style="display:grid;grid-template-columns:1fr 1fr;gap:10px;">'
                            f'<div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;">Standard Rate</div><div style="font-size:13px;font-weight:700;color:#0f172a;margin-top:1px;">{pricing}</div></div>'
                            f'<div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;">Discounted Rate</div><div style="font-size:13px;font-weight:700;color:#065F46;margin-top:1px;">{discount}</div></div>'
                            f'<div style="grid-column:span 2;"><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;">Discount Applies When</div><div style="font-size:11px;color:#334155;margin-top:1px;">{discount_cond}</div></div>'
                            f'<div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;">Rev Share</div><div style="font-size:12px;font-weight:600;color:#0f172a;margin-top:1px;">{revshare_status}</div></div>'
                            f'<div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;">Merchants Live</div><div style="font-size:12px;font-weight:600;color:#0f172a;margin-top:1px;">{merch_live}</div></div>'
                            f'</div></div>', unsafe_allow_html=True)

            # Row 2: Verticals + Features + Onboarding
            col_c, col_d = st.columns(2)
            with col_c:
                vert_pills = " ".join(f'<span style="font-size:10px;font-weight:600;padding:3px 10px;border-radius:6px;background:#F1F5F9;color:#334155;">{v.strip()}</span>' for v in verticals.split(",")) if verticals != "—" else '<span style="font-size:11px;color:#94a3b8;">—</span>'
                pm_pills = " ".join(f'<span style="font-size:10px;font-weight:600;padding:3px 10px;border-radius:6px;background:#EFF6FF;color:#1D4ED8;">{pm}</span>' for pm in pm_display[:8])
                st.markdown(f'<div style="background:#fff;border:1px solid #EEF2F7;border-radius:10px;padding:16px 18px;margin-top:10px;">'
                            f'<div style="font-size:10px;font-weight:700;color:#5B5FDE;text-transform:uppercase;letter-spacing:0.08em;margin-bottom:10px;">Verticals & Payment Methods</div>'
                            f'<div style="font-size:9px;color:#94a3b8;text-transform:uppercase;margin-bottom:4px;">Key Verticals</div>'
                            f'<div style="display:flex;flex-wrap:wrap;gap:4px;margin-bottom:10px;">{vert_pills}</div>'
                            f'<div style="font-size:9px;color:#94a3b8;text-transform:uppercase;margin-bottom:4px;">Payment Methods Supported</div>'
                            f'<div style="display:flex;flex-wrap:wrap;gap:4px;">{pm_pills}</div>'
                            f'</div>', unsafe_allow_html=True)
            with col_d:
                st.markdown(f'<div style="background:#fff;border:1px solid #EEF2F7;border-radius:10px;padding:16px 18px;margin-top:10px;">'
                            f'<div style="font-size:10px;font-weight:700;color:#5B5FDE;text-transform:uppercase;letter-spacing:0.08em;margin-bottom:10px;">Integration & Onboarding</div>'
                            f'<div style="display:flex;flex-direction:column;gap:6px;margin-bottom:10px;">'
                            f'<div style="display:flex;align-items:center;gap:8px;">{_tick(p.get("int_ready"))}<span style="font-size:12px;color:#0f172a;">Integration Ready</span></div>'
                            f'<div style="display:flex;align-items:center;gap:8px;">{_tick(p.get("nda"))}<span style="font-size:12px;color:#0f172a;">NDA Signed</span></div>'
                            f'<div style="display:flex;align-items:center;gap:8px;">{_tick(p.get("int_ready"))}<span style="font-size:12px;color:#0f172a;">Latest API Integrated</span></div>'
                            f'<div style="display:flex;align-items:center;gap:8px;">{_tick(p.get("int_used"))}<span style="font-size:12px;color:#0f172a;">Merchant Using Integration</span></div>'
                            f'</div>'
                            f'<div style="font-size:9px;color:#94a3b8;text-transform:uppercase;margin-bottom:3px;">Onboarding Requirements</div>'
                            f'<div style="font-size:11px;color:#334155;line-height:1.5;">{onboarding}</div>'
                            f'</div>', unsafe_allow_html=True)

            # Merchants row (if data available)
            if merch_names and merch_names != "—":
                st.markdown(f'<div style="background:#fff;border:1px solid #EEF2F7;border-radius:10px;padding:12px 18px;margin-top:10px;display:flex;align-items:center;gap:12px;">'
                            f'<span style="font-size:10px;font-weight:700;color:#94a3b8;text-transform:uppercase;">Merchants</span>'
                            f'<span style="font-size:12px;color:#334155;">{merch_names}</span>'
                            f'<button style="margin-left:auto;background:#EFF6FF;border:0.5px solid #BFDBFE;border-radius:6px;padding:5px 12px;font-size:11px;font-weight:600;color:#1D4ED8;cursor:pointer;display:flex;align-items:center;gap:4px;font-family:inherit;">'
                            f'<svg width="11" height="11" viewBox="0 0 11 11" fill="none" stroke="#1D4ED8" stroke-width="1.5" stroke-linecap="round"><path d="M5.5 1v6M3.5 5l2 2 2-2"/><path d="M2 9h7"/></svg>Download RFI/RFP</button>'
                            f'</div>', unsafe_allow_html=True)

            st.markdown('</div>', unsafe_allow_html=True)

    # No extra JS needed — st.expander handles everything
    components.html("""<script>
setTimeout(function(){
</script>""", height=0)

# ── Contacts View ──────────────────────────────────────────────────────────────
YUNO_CONTACTS = [
    {"init":"DR","bg":"rgba(79,70,229,0.2)","color":"#818cf8","name":"Daniela Reyes","role":"Head of Partnerships","dept":"Partnerships · Yuno","badge":"Your Lead","badge_class":"p-purple","last":"Today","rel":"Primary contact","scope":"Strategy & commercial terms"},
    {"init":"MS","bg":"rgba(79,70,229,0.2)","color":"#818cf8","name":"Marco Silva","role":"Partnership Manager · LATAM","dept":"Partnerships · Yuno","badge":"BD","badge_class":"p-blue","last":"1d ago","rel":"Day-to-day","scope":"Pipeline & deal support"},
    {"init":"LG","bg":"rgba(79,70,229,0.2)","color":"#818cf8","name":"Laura Gómez","role":"Partnership Manager · MX & COL","dept":"Partnerships · Yuno","badge":"BD","badge_class":"p-blue","last":"3d ago","rel":"Day-to-day","scope":"MX & COL opportunities"},
    {"init":"JR","bg":"rgba(20,184,166,0.2)","color":"#2dd4bf","name":"Jorge Restrepo","role":"Delivery Manager","dept":"Integrations · Yuno","badge":"Technical","badge_class":"p-teal","last":"2d ago","rel":"Integration lead","scope":"Onboarding & go-live"},
    {"init":"VP","bg":"rgba(20,184,166,0.2)","color":"#2dd4bf","name":"Valentina Perez","role":"Solutions Engineer","dept":"Integrations · Yuno","badge":"Technical","badge_class":"p-teal","last":"5d ago","rel":"Tech support","scope":"API & certification queries"},
    {"init":"AS","bg":"rgba(34,197,94,0.2)","color":"#4ade80","name":"Andrés Suárez","role":"Marketing Manager · Partners","dept":"Marketing · Yuno","badge":"Marketing","badge_class":"p-green","last":"1w ago","rel":"Co-marketing","scope":"Joint campaigns & content"},
    {"init":"CB","bg":"rgba(34,197,94,0.2)","color":"#4ade80","name":"Carolina Blanco","role":"Partner Marketing Specialist","dept":"Marketing · Yuno","badge":"Marketing","badge_class":"p-green","last":"4d ago","rel":"Co-marketing","scope":"Case studies & events"},
    {"init":"RM","bg":"rgba(245,158,11,0.2)","color":"#fbbf24","name":"Rafael Mendoza","role":"Customer Success Manager","dept":"Success · Yuno","badge":"Success","badge_class":"p-amber","last":"Today","rel":"Escalation point","scope":"Performance & health"},
    {"init":"NF","bg":"rgba(245,158,11,0.2)","color":"#fbbf24","name":"Natalia Ferro","role":"Revenue Operations","dept":"RevOps · Yuno","badge":"Ops","badge_class":"p-amber","last":"2d ago","rel":"Reporting","scope":"Invoicing & SLAs"},
]

def show_contacts():
    is_internal = st.session_state.role == "internal"

    BADGE_STYLES = {
        "p-green":  ("background:#1a1a1a;color:#fff;", "#065f46"),
        "p-blue":   ("background:#dbeafe;color:#1e40af;", "#1e40af"),
        "p-amber":  ("background:#fef3c7;color:#92400e;", "#92400e"),
        "p-teal":   ("background:#ccfbf1;color:#115e59;", "#115e59"),
        "p-purple": ("background:#4F46E5;color:#fff;", "#4F46E5"),
    }

    if not is_internal:
        st.markdown('<div style="font-size:28px;font-weight:700;color:#1d1d1f;letter-spacing:-0.8px;margin-bottom:4px;">Your Yuno Team</div>'
                    '<div style="font-size:13px;color:#6b7280;margin-bottom:24px;">Your dedicated Yuno contacts across partnerships, integrations, marketing & success</div>', unsafe_allow_html=True)
        contacts_data = YUNO_CONTACTS
    else:
        st.markdown('<div style="font-size:28px;font-weight:700;color:#1d1d1f;letter-spacing:-0.8px;margin-bottom:4px;">Key Contacts</div>'
                    '<div style="font-size:13px;color:#6b7280;margin-bottom:24px;">BD counterparts, technical leads & executive sponsors across partner organizations</div>', unsafe_allow_html=True)
        contacts_data = CONTACTS

    # ── Summary table ─────────────────────────────────────────────────────────
    header_cols = ["Name", "Role", "Department" if not is_internal else "Company", "Status", "Last Contact"]
    header_html = "".join(f'<th style="padding:10px 14px;font-size:11px;font-weight:600;color:#6b7280;text-transform:uppercase;letter-spacing:0.5px;text-align:left;border-bottom:2px solid #e5e7eb;">{h}</th>' for h in header_cols)
    rows_html = ""
    for c in contacts_data:
        badge_style = BADGE_STYLES.get(c['badge_class'], ("background:#f1f5f9;color:#475569;", "#475569"))[0]
        org = c.get('dept', '') if not is_internal else c.get('company', '')
        rows_html += f"""<tr style="border-bottom:1px solid #f1f5f9;">
  <td style="padding:12px 14px;font-size:13px;font-weight:600;color:#1a1a1a;">{c['name']}</td>
  <td style="padding:12px 14px;font-size:13px;color:#6b7280;">{c['role']}</td>
  <td style="padding:12px 14px;font-size:13px;color:#6b7280;">{org}</td>
  <td style="padding:12px 14px;"><span style="{badge_style}font-size:11px;font-weight:600;padding:3px 12px;border-radius:20px;white-space:nowrap;">{c['badge']}</span></td>
  <td style="padding:12px 14px;font-size:13px;color:#6b7280;font-variant-numeric:tabular-nums;">{c['last']}</td>
</tr>"""
    st.markdown(f"""
<div style="background:#fff;border:1px solid #e5e7eb;border-radius:8px;overflow:hidden;margin-bottom:28px;">
  <table style="width:100%;border-collapse:collapse;">
    <thead><tr>{header_html}</tr></thead>
    <tbody>{rows_html}</tbody>
  </table>
</div>""", unsafe_allow_html=True)

    # ── Detail cards (partner-detail style) ───────────────────────────────────
    cols_per_row = 2
    for i in range(0, len(contacts_data), cols_per_row):
        row = contacts_data[i:i+cols_per_row]
        cols = st.columns(cols_per_row)
        for col, c in zip(cols, row):
            badge_style, badge_border = BADGE_STYLES.get(c['badge_class'], ("background:#f1f5f9;color:#475569;", "#475569"))
            org = c.get('dept', '') if not is_internal else c.get('company', '')

            if not is_internal:
                fields = [
                    ("Department", org),
                    ("Last Contact", c['last']),
                    ("Relationship", c.get('rel', '-')),
                    ("Scope", c.get('scope', '-')),
                ]
            else:
                fields = [
                    ("Company", org),
                    ("Last Contact", c['last']),
                    ("Relationship", c.get('rel', '-')),
                    ("Active Deals", c.get('deals', '-')),
                ]

            fields_html = ""
            for fi, (label, value) in enumerate(fields):
                border_b = "border-bottom:1px solid #f1f5f9;padding-bottom:12px;margin-bottom:12px;" if fi < len(fields) - 1 else ""
                fields_html += f"""<div style="{border_b}">
  <div style="font-size:12px;color:#6b7280;margin-bottom:2px;">{label}</div>
  <div style="font-size:14px;font-weight:600;color:#1a1a1a;">{value}</div>
</div>"""

            subtitle = f"{c['role']} &middot; {org}"

            col.markdown(f"""
<div style="margin-bottom:20px;">
  <div style="display:flex;justify-content:space-between;align-items:flex-start;margin-bottom:12px;">
    <div>
      <div style="font-size:20px;font-weight:800;color:#1a1a1a;letter-spacing:-0.4px;text-transform:uppercase;">{c['name']}</div>
      <div style="font-size:13px;color:#6b7280;margin-top:2px;">{subtitle}</div>
    </div>
    <span style="{badge_style}font-size:11px;font-weight:600;padding:4px 14px;border-radius:20px;white-space:nowrap;">{c['badge']}</span>
  </div>
  <div style="border:1px solid #e5e7eb;border-radius:8px;padding:20px;">
    <div style="font-size:16px;font-weight:700;color:#1a1a1a;margin-bottom:16px;">Contact Details</div>
    <div style="display:grid;grid-template-columns:1fr 1fr;gap:0 40px;">
      {fields_html}
    </div>
    <div style="display:flex;gap:8px;margin-top:16px;border-top:1px solid #f1f5f9;padding-top:14px;">
      <button style="flex:1;padding:7px;border-radius:6px;font-size:11px;font-weight:600;background:#1a1a1a;color:#fff;border:none;cursor:pointer;font-family:inherit;">Email</button>
      <button style="flex:1;padding:7px;border-radius:6px;font-size:11px;font-weight:600;background:#fff;color:#1a1a1a;border:1px solid #e5e7eb;cursor:pointer;font-family:inherit;">Schedule</button>
      <button style="flex:1;padding:7px;border-radius:6px;font-size:11px;font-weight:600;background:#fff;color:#1a1a1a;border:1px solid #e5e7eb;cursor:pointer;font-family:inherit;">{"Slack" if not is_internal else "Note"}</button>
    </div>
  </div>
</div>""", unsafe_allow_html=True)

    if is_internal:
        st.markdown("""
<div style="border:1px solid rgba(245,158,11,.22);background:rgba(245,158,11,.04);border-radius:8px;padding:16px 20px;margin-top:8px;">
  <div style="font-size:10px;font-weight:700;color:#92400e;letter-spacing:.8px;text-transform:uppercase;margin-bottom:12px;">Internal BD Notes — Not Visible to Partners</div>
  <div style="display:flex;justify-content:space-between;align-items:center;padding:8px 0;border-bottom:1px solid rgba(0,0,0,0.06);font-size:13px;"><span style="color:#6b7280;">Tom Kuehn (Adyen) — negotiating exclusivity window on MX corridor, do not share with other acquirers</span><span style="font-size:11px;font-weight:600;color:#92400e;">Confidential</span></div>
  <div style="display:flex;justify-content:space-between;align-items:center;padding:8px 0;border-bottom:1px solid rgba(0,0,0,0.06);font-size:13px;"><span style="color:#6b7280;">Felipe Morales (Getnet) — board approval needed above $500K. Target Jan board meeting.</span><span style="font-size:11px;font-weight:600;color:#92400e;">Action needed</span></div>
  <div style="display:flex;justify-content:space-between;align-items:center;padding:8px 0;font-size:13px;"><span style="color:#6b7280;">Martín Castillo (Pomelo) — exploring white-label orchestration. High strategic value for BaaS play.</span><span style="font-size:11px;font-weight:600;color:#92400e;">Strategic</span></div>
</div>""", unsafe_allow_html=True)

# ── Performance View ───────────────────────────────────────────────────────────
def show_performance():
    is_internal = st.session_state.role == "internal"

    # ── Hero ──────────────────────────────────────────────────────────────────
    st.markdown("""
<div style="background:#0C1220;border-radius:16px;padding:24px 28px;margin-bottom:14px;position:relative;overflow:hidden;">
  <div style="position:absolute;inset:0;background-image:radial-gradient(circle at 1px 1px,rgba(255,255,255,0.03) 1px,transparent 0);background-size:22px 22px;pointer-events:none;"></div>
  <div style="position:relative;z-index:1;">
    <div style="font-size:22px;font-weight:800;color:#F1F5F9;letter-spacing:-0.02em;margin-bottom:4px;">Partner Health</div>
    <div style="font-size:12px;color:#64748b;">Performance analytics across partners and merchants</div>
  </div>
</div>""", unsafe_allow_html=True)

    # ── View toggle ───────────────────────────────────────────────────────────
    perf_view = st.session_state.get("perf_view", "partners")
    vc1, vc2 = st.columns(2)
    with vc1:
        if st.button("Partnership Performance", key="pv_partners", use_container_width=True):
            st.session_state["perf_view"] = "partners"
            st.rerun()
    with vc2:
        if st.button("Merchant Performance", key="pv_merchants", use_container_width=True):
            st.session_state["perf_view"] = "merchants"
            st.rerun()

    # Style toggle buttons
    components.html(f"""<script>
setTimeout(function(){{
  var doc = window.parent.document;
  doc.querySelectorAll('button').forEach(function(btn){{
    var t = btn.textContent.trim();
    if(t==='Partnership Performance'){{
      btn.style.cssText='background:{"#0C1220" if perf_view=="partners" else "#fff"} !important;border:1px solid {"#5B5FDE" if perf_view=="partners" else "#E2E8F0"} !important;border-radius:10px !important;padding:10px !important;font-size:13px !important;color:{"#A5B4FC" if perf_view=="partners" else "#64748b"} !important;font-weight:700 !important;cursor:pointer !important;box-shadow:none !important;';
    }}
    if(t==='Merchant Performance'){{
      btn.style.cssText='background:{"#0C1220" if perf_view=="merchants" else "#fff"} !important;border:1px solid {"#5B5FDE" if perf_view=="merchants" else "#E2E8F0"} !important;border-radius:10px !important;padding:10px !important;font-size:13px !important;color:{"#A5B4FC" if perf_view=="merchants" else "#64748b"} !important;font-weight:700 !important;cursor:pointer !important;box-shadow:none !important;';
    }}
  }});
}},200);
</script>""", height=0)

    # ══════════════════════════════════════════════════════════════════════════
    # MERCHANT PERFORMANCE VIEW
    # ══════════════════════════════════════════════════════════════════════════
    if perf_view == "merchants":
        _MERCH_DATA = [
            {"name":"Uber","tpv":"$5.9M","txns":"323K","ar":93.4,"providers":["Adyen","dLocal","Stripe"],"country":"MX","missing":["Smart retries","PIX in BR"],"upsell":"Add Cielo for local acquiring in Brazil — +3% auth rate"},
            {"name":"Netflix","tpv":"$4.8M","txns":"306K","ar":95.6,"providers":["Stripe","Adyen","dLocal"],"country":"US","missing":["Installments in LATAM"],"upsell":"Enable parcelamento via Cielo/PagBank — unlock +20% AOV in Brazil"},
            {"name":"Spotify","tpv":"$4.2M","txns":"421K","ar":94.8,"providers":["Adyen","Ebanx","dLocal"],"country":"BR","missing":["OXXO in MX","Nequi in CO"],"upsell":"Add Kushki for Colombia APMs — reduce decline rate by 4%"},
            {"name":"Rappi","tpv":"$8.4M","txns":"259K","ar":91.2,"providers":["Adyen","Kushki","MercadoPago"],"country":"CO","missing":["Fraud layer","3DS optimization"],"upsell":"Add SEON fraud screening — reduce chargebacks by 40%"},
            {"name":"MercadoLibre","tpv":"$7.3M","txns":"62K","ar":88.7,"providers":["MercadoPago","Adyen","Cielo"],"country":"AR","missing":["Network tokens","Smart routing"],"upsell":"Enable network tokens via Adyen — +2% auth rate on recurring"},
            {"name":"Falabella","tpv":"$3.1M","txns":"36K","ar":87.1,"providers":["Transbank","Kushki","Adyen"],"country":"CL","missing":["Installments","Webpay Plus"],"upsell":"Upgrade Transbank to Webpay Plus — unlock 6/12 month installments"},
            {"name":"iFood","tpv":"$2.9M","txns":"103K","ar":91.8,"providers":["Cielo","PagBank","Stone"],"country":"BR","missing":["Apple Pay","Google Pay"],"upsell":"Add Apple Pay/Google Pay — +8% mobile conversion"},
            {"name":"Despegar","tpv":"$3.8M","txns":"12K","ar":82.3,"providers":["Cybersource","Adyen","dLocal"],"country":"AR","missing":["Airline data","3DS exemptions"],"upsell":"Enable airline data passthrough — +5% auth rate on high-value bookings"},
        ]

        stat_row([
            {"label":"Live Merchants","value":"12","delta":"↑ 2 this quarter","delta_type":"up"},
            {"label":"Total TPV","value":"$40.4M","delta":"↑ 28% QoQ","delta_type":"up","val_color":"#818cf8"},
            {"label":"Avg Auth Rate","value":"90.6%","delta":"↑ 1.4pts","delta_type":"up"},
            {"label":"Cross-sell Opps","value":"8","delta":"★ actionable now","delta_type":"up","val_color":"#F59E0B"},
        ])

        st.markdown('<div style="font-size:14px;font-weight:700;color:#0f172a;margin:16px 0 8px;">Merchant Portfolio — Partners & Performance</div>', unsafe_allow_html=True)

        for m in _MERCH_DATA:
            ar_color = "#065F46" if m["ar"] >= 90 else "#B45309" if m["ar"] >= 85 else "#DC2626"
            providers_pills = " ".join(f'<span style="font-size:10px;font-weight:600;padding:2px 8px;border-radius:5px;background:#EFF6FF;color:#1D4ED8;">{p}</span>' for p in m["providers"])
            missing_pills = " ".join(f'<span style="font-size:10px;font-weight:600;padding:2px 8px;border-radius:5px;background:#FEF3C7;color:#92400E;">{f}</span>' for f in m["missing"])

            with st.expander(f"{m['name']}  ·  {m['tpv']} TPV  ·  {m['ar']}% auth  ·  {m['country']}", expanded=False):
                col_m1, col_m2 = st.columns(2)
                with col_m1:
                    st.markdown(f'<div style="background:#fff;border:1px solid #EEF2F7;border-radius:10px;padding:16px;">'
                                f'<div style="font-size:10px;font-weight:700;color:#5B5FDE;text-transform:uppercase;letter-spacing:0.08em;margin-bottom:10px;">Performance</div>'
                                f'<div style="display:grid;grid-template-columns:1fr 1fr;gap:8px;">'
                                f'<div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;">TPV/mo</div><div style="font-size:14px;font-weight:700;color:#0f172a;">{m["tpv"]}</div></div>'
                                f'<div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;">Txns/mo</div><div style="font-size:14px;font-weight:700;color:#0f172a;">{m["txns"]}</div></div>'
                                f'<div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;">Auth Rate</div><div style="font-size:14px;font-weight:700;color:{ar_color};">{m["ar"]}%</div></div>'
                                f'<div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;">Country</div><div style="font-size:14px;font-weight:700;color:#0f172a;">{m["country"]}</div></div>'
                                f'</div>'
                                f'<div style="margin-top:10px;"><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;margin-bottom:4px;">Live Partners</div>'
                                f'<div style="display:flex;flex-wrap:wrap;gap:3px;">{providers_pills}</div></div>'
                                f'</div>', unsafe_allow_html=True)
                with col_m2:
                    st.markdown(f'<div style="background:#FFFBEB;border:1px solid #FDE68A;border-radius:10px;padding:16px;">'
                                f'<div style="font-size:10px;font-weight:700;color:#92400E;text-transform:uppercase;letter-spacing:0.08em;margin-bottom:10px;">Optimization Opportunities</div>'
                                f'<div style="font-size:9px;color:#92400E;text-transform:uppercase;margin-bottom:4px;">Missing Features</div>'
                                f'<div style="display:flex;flex-wrap:wrap;gap:3px;margin-bottom:10px;">{missing_pills}</div>'
                                f'<div style="font-size:9px;color:#92400E;text-transform:uppercase;margin-bottom:4px;">Recommended Action</div>'
                                f'<div style="font-size:12px;color:#0f172a;line-height:1.5;font-weight:500;">{m["upsell"]}</div>'
                                f'</div>', unsafe_allow_html=True)
        return

    # ══════════════════════════════════════════════════════════════════════════
    # PARTNERSHIP PERFORMANCE VIEW (original)
    # ══════════════════════════════════════════════════════════════════════════
    stat_row([
        {"label":"Total TPV (Live Partners)","value":"$2.4B","delta":"↑ 34% YoY","delta_type":"up","val_color":"#818cf8"},
        {"label":"Transactions/mo","value":"142M","delta":"↑ 22% MoM","delta_type":"up"},
        {"label":"Auth Rate (avg)","value":"89.4%","delta":"↑ 1.2pts","delta_type":"up"},
        {"label":"Partner Revenue","value":"$18.6M","delta":"↑ 28% QoQ","delta_type":"up"},
        {"label":"NPS (Partners)","value":"72","delta":"↑ 8pts","delta_type":"up"},
    ])

    col_tpv, col_mix = st.columns([2.2, 1])

    with col_tpv:
        months = ["May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]
        fig = go.Figure()
        datasets = [
            ("Acquirers",[820,880,910,950,1020,1100,1180,1240],"rgba(59,130,246,.7)"),
            ("PSPs",[420,460,490,530,580,640,700,780],"rgba(168,85,247,.7)"),
            ("APMs",[180,200,230,260,300,340,390,440],"rgba(20,184,166,.7)"),
            ("Fraud",[60,70,80,90,100,115,130,150],"rgba(239,68,68,.6)"),
            ("BaaS",[0,0,0,0,10,25,45,80],"rgba(245,158,11,.7)"),
        ]
        for name, data, color in datasets:
            fig.add_trace(go.Bar(name=name, x=months, y=data, marker_color=color))
        fig.update_layout(
            barmode="stack", paper_bgcolor="#ffffff", plot_bgcolor="#ffffff",
            font=dict(color="#6e6e73", size=10),
            margin=dict(l=40,r=10,t=30,b=30),
            legend=dict(orientation="h", y=1.1, font=dict(size=10)),
            height=280,
            title=dict(text="Monthly TPV by Partner Category", font=dict(color="#1d1d1f", size=13), x=0),
            yaxis=dict(gridcolor="rgba(0,0,0,0.06)", ticksuffix="M"),
            xaxis=dict(gridcolor="rgba(0,0,0,0.06)"),
        )
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})

    with col_mix:
        fig2 = go.Figure(go.Pie(
            labels=["Acquirers","PSPs","APMs","Fraud","BaaS"],
            values=[38,31,19,8,4],
            hole=0.65,
            marker_colors=["rgba(59,130,246,.8)","rgba(168,85,247,.8)","rgba(20,184,166,.8)","rgba(239,68,68,.75)","rgba(245,158,11,.8)"],
            textinfo="none",
        ))
        fig2.update_layout(
            paper_bgcolor="#ffffff", plot_bgcolor="#ffffff",
            font=dict(color="#6e6e73",size=10),
            margin=dict(l=10,r=10,t=50,b=10),
            legend=dict(font=dict(size=10,color="#6e6e73")),
            height=280,
            title=dict(text="TPV Mix by Category", font=dict(color="#1d1d1f",size=13), x=0),
        )
        st.plotly_chart(fig2, use_container_width=True, config={"displayModeBar": False})

    col_auth, col_health = st.columns(2)

    with col_auth:
        partners_list = ["Adyen","Prosa","Getnet","OpenPix","Kushki","Cielo","SEON","Nuvei"]
        rates = [92.4, 90.1, 89.6, 99.1, 88.2, 86.9, 85.0, 84.1]
        colors = ["rgba(34,197,94,.7)" if r>=90 else "rgba(59,130,246,.7)" if r>=87 else "rgba(239,68,68,.7)" for r in rates]
        fig3 = go.Figure(go.Bar(x=rates, y=partners_list, orientation="h", marker_color=colors, marker=dict(cornerradius=3)))
        fig3.update_layout(
            paper_bgcolor="#ffffff", plot_bgcolor="#ffffff",
            font=dict(color="#6e6e73",size=10),
            margin=dict(l=10,r=20,t=50,b=30),
            height=260,
            title=dict(text="Auth Rate by Partner (Top 8)", font=dict(color="#1d1d1f",size=13), x=0),
            xaxis=dict(range=[80,100], ticksuffix="%", gridcolor="rgba(0,0,0,0.06)"),
            yaxis=dict(gridcolor="rgba(0,0,0,0)"),
        )
        st.plotly_chart(fig3, use_container_width=True, config={"displayModeBar": False})

    with col_health:
        st.markdown("""
<div style="background:#fff;box-shadow:0 1px 3px rgba(0,0,0,0.07),0 0 0 1px rgba(0,0,0,0.04);border:none;border-radius:10px;padding:18px;">
  <div style="font-size:13px;font-weight:700;color:#1d1d1f;margin-bottom:2px;">Partner Health Scorecard</div>
  <div style="font-size:11px;color:#86868b;margin-bottom:14px;">Integration quality · engagement · growth</div>
  <table style="width:100%;border-collapse:collapse;font-size:12px;">
    <tr style="background:#f5f5f7;"><th style="text-align:left;padding:6px 8px;font-size:9.5px;font-weight:700;letter-spacing:.9px;text-transform:uppercase;color:#86868b;">Partner</th><th style="padding:6px 8px;font-size:9.5px;font-weight:700;letter-spacing:.9px;text-transform:uppercase;color:#86868b;">Health</th><th style="padding:6px 8px;font-size:9.5px;font-weight:700;letter-spacing:.9px;text-transform:uppercase;color:#86868b;">TPV Trend</th><th style="padding:6px 8px;font-size:9.5px;font-weight:700;letter-spacing:.9px;text-transform:uppercase;color:#86868b;">Issues</th></tr>
    <tr style="border-bottom:1px solid rgba(0,0,0,0.07);"><td style="padding:7px 8px;font-weight:600;color:#1d1d1f;">Adyen</td><td style="padding:7px 8px;"><span style="background:rgba(34,197,94,.14);color:#4ade80;font-size:10px;font-weight:600;padding:3px 9px;border-radius:20px;font-family:monospace;">95</span></td><td style="padding:7px 8px;color:#22c55e;font-family:monospace;">+18%</td><td style="padding:7px 8px;color:#86868b;">0</td></tr>
    <tr style="border-bottom:1px solid rgba(0,0,0,0.07);"><td style="padding:7px 8px;font-weight:600;color:#1d1d1f;">Getnet</td><td style="padding:7px 8px;"><span style="background:rgba(34,197,94,.14);color:#4ade80;font-size:10px;font-weight:600;padding:3px 9px;border-radius:20px;font-family:monospace;">88</span></td><td style="padding:7px 8px;color:#22c55e;font-family:monospace;">+42%</td><td style="padding:7px 8px;color:#86868b;">1</td></tr>
    <tr style="border-bottom:1px solid rgba(0,0,0,0.07);"><td style="padding:7px 8px;font-weight:600;color:#1d1d1f;">Kushki</td><td style="padding:7px 8px;"><span style="background:rgba(59,130,246,.14);color:#60a5fa;font-size:10px;font-weight:600;padding:3px 9px;border-radius:20px;font-family:monospace;">81</span></td><td style="padding:7px 8px;color:#22c55e;font-family:monospace;">+9%</td><td style="padding:7px 8px;color:#86868b;">2</td></tr>
    <tr style="border-bottom:1px solid rgba(0,0,0,0.07);"><td style="padding:7px 8px;font-weight:600;color:#1d1d1f;">OpenPix</td><td style="padding:7px 8px;"><span style="background:rgba(59,130,246,.14);color:#60a5fa;font-size:10px;font-weight:600;padding:3px 9px;border-radius:20px;font-family:monospace;">79</span></td><td style="padding:7px 8px;color:#22c55e;font-family:monospace;">+31%</td><td style="padding:7px 8px;color:#d97706;">1</td></tr>
    <tr style="border-bottom:1px solid rgba(0,0,0,0.07);"><td style="padding:7px 8px;font-weight:600;color:#1d1d1f;">SEON</td><td style="padding:7px 8px;"><span style="background:rgba(245,158,11,.14);color:#fbbf24;font-size:10px;font-weight:600;padding:3px 9px;border-radius:20px;font-family:monospace;">71</span></td><td style="padding:7px 8px;color:#22c55e;font-family:monospace;">+6%</td><td style="padding:7px 8px;color:#d97706;">3</td></tr>
    <tr><td style="padding:7px 8px;font-weight:600;color:#1d1d1f;">Nuvei</td><td style="padding:7px 8px;"><span style="background:rgba(245,158,11,.14);color:#fbbf24;font-size:10px;font-weight:600;padding:3px 9px;border-radius:20px;font-family:monospace;">68</span></td><td style="padding:7px 8px;color:#ef4444;font-family:monospace;">−4%</td><td style="padding:7px 8px;color:#ef4444;">4</td></tr>
  </table>
</div>""", unsafe_allow_html=True)

    if is_internal:
        st.markdown("""
<div style="border:1px solid rgba(245,158,11,.22);background:rgba(245,158,11,.04);border-radius:10px;padding:13px 16px;margin-top:16px;">
  <div style="font-size:9.5px;font-weight:700;color:#d97706;letter-spacing:.9px;text-transform:uppercase;margin-bottom:10px;">🔒 Internal Revenue Intelligence</div>
  <div style="display:flex;justify-content:space-between;align-items:center;padding:5px 0;border-bottom:1px solid rgba(0,0,0,0.06);font-size:12px;"><span style="color:#6e6e73;">Blended take rate (all live partners)</span><span style="font-family:monospace;color:#d97706;">0.048%</span></div>
  <div style="display:flex;justify-content:space-between;align-items:center;padding:5px 0;border-bottom:1px solid rgba(0,0,0,0.06);font-size:12px;"><span style="color:#6e6e73;">Adyen contract renewal date</span><span style="font-family:monospace;color:#d97706;">Mar 2025 · auto-renew flag on</span></div>
  <div style="display:flex;justify-content:space-between;align-items:center;padding:5px 0;border-bottom:1px solid rgba(0,0,0,0.06);font-size:12px;"><span style="color:#6e6e73;">Nuvei performance clause risk</span><span style="font-family:monospace;color:#d97706;">⚠ Below SLA — escalate</span></div>
  <div style="display:flex;justify-content:space-between;align-items:center;padding:5px 0;font-size:12px;"><span style="color:#6e6e73;">BaaS vertical projected ARR (12mo)</span><span style="font-family:monospace;color:#d97706;">$2.1M if 4 deals close</span></div>
</div>""", unsafe_allow_html=True)

# ── Insights View ──────────────────────────────────────────────────────────────
def show_insights():
    is_internal = st.session_state.role == "internal"
    st.markdown('<div style="font-size:28px;font-weight:700;color:#1d1d1f;letter-spacing:-0.8px;margin-bottom:4px;">Global Market Intel</div>'
                '<div style="font-size:11px;color:#86868b;margin-bottom:16px;">Payment landscape insights across key markets worldwide</div>', unsafe_allow_html=True)

    stat_row([
        {"label":"Markets Covered","value":"45+","delta":"Across LATAM, APAC, MENAT, EU, NA","delta_type":"flat","val_color":"#818cf8"},
        {"label":"Payment Methods Tracked","value":"300+","delta":"Cards, wallets, A2A, BNPL, crypto","delta_type":"flat"},
        {"label":"Avg. Approval Rate (Global)","value":"89.2%","delta":"↑ 1.4pp vs prior year","delta_type":"up"},
        {"label":"Real-Time Insights","value":"Live","delta":"Updated continuously","delta_type":"up","val_color":"#16a34a"},
    ])

    tab_options = ["Country Profiles","Payment Trends","Regulatory Updates","Market News"]
    tab = st.radio("insights_tab", tab_options, horizontal=True, key="insight_tab_radio", label_visibility="collapsed")

    # ── Country Profiles Data ─────────────────────────────────────────────────
    _COUNTRIES = {
        "Brazil": {
            "flag":"🇧🇷","population":"216M","gdp":"$2.17T","ecom_size":"$62B (2025)",
            "payment_methods":[("PIX",40),("Credit Card",30),("Boleto Bancario",15),("Debit Card",10),("Other",5)],
            "top_acquirers":"Cielo, Rede (Itau), Getnet (Santander), Stone, PagSeguro, Adyen",
            "regulatory":"Central Bank regulates PIX. Open Finance mandate in effect. PIX credit launching Q2 2026. Data localization not required but LGPD (data privacy) applies to all payment data.",
            "settlement":"PIX: instant. Cards: D+1 to D+30 (negotiable). Boleto: D+1 after payment confirmation.",
            "fx":"BRL is not freely convertible. Central Bank controls FX. Merchants must use authorized dealers. IOF tax (up to 6.38%) applies to cross-border transactions.",
            "instore_online":"Online 48% / In-store 52%. PIX rapidly gaining in-store share via QR code.",
        },
        "Mexico": {
            "flag":"🇲🇽","population":"130M","gdp":"$1.79T","ecom_size":"$38B (2025)",
            "payment_methods":[("Credit Card",32),("Debit Card",28),("OXXO (cash voucher)",18),("SPEI/CoDi",12),("Digital Wallets",10)],
            "top_acquirers":"BBVA Mexico, Banorte, Conekta, Stripe, Mercado Pago, Kushki",
            "regulatory":"Banxico regulates payment systems. CoDi (instant payments) push. Fintech Law (Ley Fintech) governs digital payments. CNBV oversees PSPs.",
            "settlement":"Cards: D+1 to D+3. SPEI: instant (banking hours). OXXO: D+1 after cash payment.",
            "fx":"MXN is freely traded. No FX restrictions for merchants. Withholding tax on cross-border payments may apply (ISR).",
            "instore_online":"Online 38% / In-store 62%. Cash-based methods (OXXO) remain critical for unbanked population (~45%).",
        },
        "Colombia": {
            "flag":"🇨🇴","population":"52M","gdp":"$363B","ecom_size":"$12B (2025)",
            "payment_methods":[("Credit Card",30),("PSE (bank transfer)",28),("Debit Card",18),("Cash (Efecty/Baloto)",15),("Digital Wallets",9)],
            "top_acquirers":"Redeban, Credibanco, Kushki, Mercado Pago, PayU, Addi (BNPL)",
            "regulatory":"Superintendencia Financiera oversees payments. Transfiya 2.0 (instant payments) launched 2026. GMF tax (4x1000) on financial transactions.",
            "settlement":"Cards: D+1 to D+7. PSE: D+1. Cash vouchers: D+1 after confirmation.",
            "fx":"COP has managed float. Central bank approval needed for large FX transactions. Repatriation of funds possible but requires documentation.",
            "instore_online":"Online 32% / In-store 68%. Cash on delivery still significant in tier-2 cities.",
        },
        "Argentina": {
            "flag":"🇦🇷","population":"46M","gdp":"$641B","ecom_size":"$14B (2025)",
            "payment_methods":[("Credit Card (installments)",38),("Debit Card",25),("Bank Transfer",15),("Digital Wallets",14),("Cash (Rapipago/PagoFacil)",8)],
            "top_acquirers":"Prisma (Visa), First Data (MC), Mercado Pago, Ualá, Modo",
            "regulatory":"BCRA regulates payments. Installment plans (cuotas) are culturally essential. Capital controls (cepo) restrict FX. DEBIN for instant debits.",
            "settlement":"Cards: D+10 to D+18 standard (accelerated available at cost). Bank transfers: D+0 to D+1.",
            "fx":"Strict capital controls. Multiple exchange rates (official, MEP, CCL). Cross-border settlement is complex. Merchants must navigate BCRA restrictions.",
            "instore_online":"Online 40% / In-store 60%. Installment-based credit card payments dominate eCommerce.",
        },
        "Chile": {
            "flag":"🇨🇱","population":"19.5M","gdp":"$335B","ecom_size":"$12B (2025)",
            "payment_methods":[("Debit Card",35),("Credit Card",30),("Bank Transfer (Khipu)",18),("Prepaid Cards",10),("Cash",7)],
            "top_acquirers":"Transbank, Getnet, Kushki, Mercado Pago, Flow",
            "regulatory":"CMF oversees financial regulation. Transbank monopoly broken by four-party model reform. New PSP licensing framework in place since 2022.",
            "settlement":"Cards: D+1 to D+2. Khipu: instant/D+1. Transbank: D+2 standard.",
            "fx":"CLP is freely traded. No significant FX restrictions. Chile has bilateral trade agreements reducing friction for cross-border commerce.",
            "instore_online":"Online 42% / In-store 58%. Debit cards overtook credit cards as primary online method in 2024.",
        },
        "India": {
            "flag":"🇮🇳","population":"1.44B","gdp":"$3.94T","ecom_size":"$120B (2025)",
            "payment_methods":[("UPI",48),("Credit Card",15),("Debit Card",14),("Net Banking",10),("Wallets (Paytm, PhonePe)",8),("Cash on Delivery",5)],
            "top_acquirers":"Razorpay, Cashfree, PayU India, CCAvenue, Juspay, Pine Labs (in-store)",
            "regulatory":"RBI regulates all payments. Zero MDR on UPI (under review). Data localization mandate — all payment data must be stored in India. Recurring payment framework (e-mandate) in effect.",
            "settlement":"UPI: instant. Cards: T+1 to T+3. Net Banking: T+1 to T+2.",
            "fx":"INR is partially convertible. RBI controls capital account. FEMA regulations govern cross-border payments. LRS limit of $250K/year for individuals.",
            "instore_online":"Online 35% / In-store 65%. UPI QR codes dominate in-store. 15B+ monthly UPI transactions as of 2026.",
        },
        "Indonesia": {
            "flag":"🇮🇩","population":"278M","gdp":"$1.42T","ecom_size":"$62B (2025)",
            "payment_methods":[("Digital Wallets (GoPay, OVO, Dana)",35),("Bank Transfer (VA)",28),("Credit Card",12),("QRIS (QR standard)",15),("Convenience Store",10)],
            "top_acquirers":"Midtrans (GoTo), Xendit, DOKU, Nicepay, Ayoconnect",
            "regulatory":"Bank Indonesia oversees payments. QRIS is the national QR standard. PBI regulations on payment systems. Local processing requirement for domestic transactions.",
            "settlement":"Bank transfers: instant (BI-FAST). Cards: D+2 to D+7. E-wallets: D+1.",
            "fx":"IDR is managed float. BI approval needed for large transfers. Repatriation rules apply. Withholding tax on cross-border digital services.",
            "instore_online":"Online 40% / In-store 60%. QRIS unifying fragmented QR code market. E-wallet adoption among highest in Southeast Asia.",
        },
        "UAE": {
            "flag":"🇦🇪","population":"10M","gdp":"$509B","ecom_size":"$10B (2025)",
            "payment_methods":[("Credit Card",40),("Debit Card",22),("Apple Pay / Samsung Pay",18),("AANI (instant)",10),("Cash on Delivery",10)],
            "top_acquirers":"Network International, Checkout.com, Tap Payments, Amazon Payment Services, Telr",
            "regulatory":"CBUAE regulates payments. AANI instant payment system launched nationwide 2026. Strong AML/KYC requirements. UAE PASS digital identity integration for onboarding.",
            "settlement":"Cards: D+1 to D+3. AANI: instant. COD reconciliation: D+1 to D+5.",
            "fx":"AED is pegged to USD (3.6725). No FX restrictions. Free zones offer simplified cross-border payment structures. No income tax; 9% corporate tax introduced 2023.",
            "instore_online":"Online 45% / In-store 55%. Contactless penetration among highest globally (>80% of in-store card transactions).",
        },
        "Saudi Arabia": {
            "flag":"🇸🇦","population":"36M","gdp":"$1.07T","ecom_size":"$16B (2025)",
            "payment_methods":[("Mada (debit)",45),("Credit Card",25),("STC Pay",12),("Apple Pay",10),("SADAD / BNPL (Tamara, Tabby)",8)],
            "top_acquirers":"Neoleap, Geidea, HyperPay, Moyasar, Tap Payments, Checkout.com",
            "regulatory":"SAMA regulates payments. Mada scheme mandatory for all domestic debit. BNPL framework introduced 2023. Open Banking launched. Vision 2030 driving cashless push — targeting 70% non-cash by 2025.",
            "settlement":"Mada: D+1. Credit cards: D+2 to D+5. STC Pay: D+1.",
            "fx":"SAR pegged to USD (3.75). No FX restrictions. Straightforward repatriation. ZATCA manages VAT (15%) compliance.",
            "instore_online":"Online 38% / In-store 62%. Mada contactless payments growing rapidly. NFC penetration > 70%.",
        },
        "Nigeria": {
            "flag":"🇳🇬","population":"224M","gdp":"$477B","ecom_size":"$8B (2025)",
            "payment_methods":[("Bank Transfer",35),("Cards (Verve, Visa, MC)",25),("USSD",18),("Mobile Money",12),("eNaira / Cash",10)],
            "top_acquirers":"Paystack (Stripe), Flutterwave, Interswitch, Kuda, Moniepoint",
            "regulatory":"CBN regulates payments. eNaira (CBDC) in circulation. Cashless policy enforced. Verve is the dominant local card scheme. PSB (Payment Service Bank) licenses expanding mobile money.",
            "settlement":"Bank transfers: instant (NIP). Cards: D+1. USSD: D+1. Mobile money: instant.",
            "fx":"NGN has unified floating rate (post-2023 reform). FX liquidity challenges persist. Repatriation can be delayed. Parallel market premiums have narrowed post-reform.",
            "instore_online":"Online 30% / In-store 70%. USSD payments critical for feature-phone users. Bank transfers growing fastest.",
        },
        "South Africa": {
            "flag":"🇿🇦","population":"62M","gdp":"$399B","ecom_size":"$7.5B (2025)",
            "payment_methods":[("Credit Card",30),("Debit Card",25),("EFT (bank transfer)",22),("Instant EFT (Ozow, Stitch)",15),("Cash / Voucher",8)],
            "top_acquirers":"Nedbank, Standard Bank, Peach Payments, PayFast (Takealot), Ozow, Stitch",
            "regulatory":"SARB and PASA regulate payments. National Payment System Act. Strong consumer protection via NCA. PayShap instant payment system launched 2023.",
            "settlement":"Cards: D+1 to D+3. EFT: D+1 to D+2. Instant EFT: D+0. PayShap: instant.",
            "fx":"ZAR is freely traded. Excon rules apply to cross-border flows. Loop structures available for multinationals. Withholding tax on royalties/services.",
            "instore_online":"Online 32% / In-store 68%. Instant EFT growing rapidly as card alternative. Tap-to-phone pilots expanding.",
        },
        "UK": {
            "flag":"🇬🇧","population":"68M","gdp":"$3.34T","ecom_size":"$170B (2025)",
            "payment_methods":[("Debit Card",42),("Credit Card",25),("Open Banking / FPS",12),("Digital Wallets",15),("BNPL (Klarna, Clearpay)",6)],
            "top_acquirers":"Worldpay, Adyen, Stripe, Barclays Payments, Checkout.com, Trust Payments",
            "regulatory":"FCA and PSR regulate payments. Strong Consumer Duty (2023). Open Banking well-established (9M+ users). SCA enforced under UK-retained PSD2. BNPL regulation pending.",
            "settlement":"Faster Payments: instant. Cards: D+1 to D+3. BACS: D+3.",
            "fx":"GBP is freely traded. No restrictions. London is global FX hub. Cross-border payments straightforward. Interchange capped at 0.2% debit / 0.3% credit.",
            "instore_online":"Online 52% / In-store 48%. Contactless limit raised to £100. Open Banking payments accelerating for eCommerce.",
        },
        "Germany": {
            "flag":"🇩🇪","population":"84M","gdp":"$4.46T","ecom_size":"$130B (2025)",
            "payment_methods":[("PayPal",28),("Invoice/Kauf auf Rechnung",20),("Credit Card",18),("SEPA Direct Debit",15),("Debit (girocard)",12),("SOFORT/Klarna",7)],
            "top_acquirers":"Adyen, Unzer, Computop, Worldline, Stripe, Payone",
            "regulatory":"BaFin and Bundesbank oversee payments. PSD2/SCA enforced. Strong data privacy (GDPR strictest interpretation). E-invoicing mandate from 2025.",
            "settlement":"SEPA: D+1. Cards: D+1 to D+3. Invoice: varies (14-30 days).",
            "fx":"EUR — no FX needed within eurozone. No restrictions for cross-border EUR payments. SEPA zone covers 36 countries.",
            "instore_online":"Online 45% / In-store 55%. Germany historically cash-heavy but post-COVID card/digital adoption surged. girocard dominates in-store.",
        },
        "France": {
            "flag":"🇫🇷","population":"68M","gdp":"$3.05T","ecom_size":"$95B (2025)",
            "payment_methods":[("Credit/Debit Card (CB)",45),("PayPal",15),("Bank Transfer (SEPA)",12),("Digital Wallets",10),("Cartes Bancaires",10),("BNPL / Cheque",8)],
            "top_acquirers":"Worldline, Adyen, Stripe, PayPlug, Dalenys (Natixis), Lyra",
            "regulatory":"ACPR/Banque de France regulate. PSD2/SCA enforced. Cartes Bancaires (CB) is mandatory co-badging on all French cards. DGCCRF consumer protection applies.",
            "settlement":"Cards: D+1 to D+2. SEPA Credit Transfer: D+1. SEPA Instant: seconds.",
            "fx":"EUR — no FX within eurozone. Interchange capped per EU regulation.",
            "instore_online":"Online 42% / In-store 58%. Carte Bancaire network processes majority of domestic transactions. Contactless widely adopted.",
        },
        "USA": {
            "flag":"🇺🇸","population":"335M","gdp":"$28.78T","ecom_size":"$1.14T (2025)",
            "payment_methods":[("Credit Card",38),("Debit Card",28),("Digital Wallets (Apple Pay, Google Pay)",16),("ACH/Bank Transfer",10),("BNPL",5),("Cash/Other",3)],
            "top_acquirers":"Chase Paymentech, Fiserv (First Data), Worldpay, Adyen, Stripe, Square, Braintree",
            "regulatory":"Federal Reserve, OCC, CFPB, and state regulators. Durbin amendment caps debit interchange. No federal data privacy law (state-level patchwork — CCPA, etc.). FedNow instant payment system live since 2023.",
            "settlement":"ACH: D+1 to D+2. FedNow: instant. Cards: D+1 to D+2. Wire: same-day.",
            "fx":"USD is global reserve currency. No FX restrictions. Straightforward cross-border settlement.",
            "instore_online":"Online 42% / In-store 58%. Contactless adoption growing but still below EU levels. Tap-to-pay at ~35% of in-store card transactions.",
        },
        "Canada": {
            "flag":"🇨🇦","population":"40M","gdp":"$2.14T","ecom_size":"$58B (2025)",
            "payment_methods":[("Credit Card",38),("Debit (Interac)",30),("Digital Wallets",14),("Bank Transfer (EFT)",10),("BNPL",5),("Cash",3)],
            "top_acquirers":"Moneris, Global Payments, Chase (Canada), Adyen, Stripe, Nuvei",
            "regulatory":"Bank of Canada and OSFI oversee payments. Payments Canada modernizing with RTR (Real-Time Rail). Interchange regulated. PIPEDA (privacy). Retail Payment Activities Act (2024).",
            "settlement":"Interac e-Transfer: instant. Cards: D+1 to D+2. EFT: D+1 to D+3. RTR (upcoming): instant.",
            "fx":"CAD is freely traded. No FX restrictions. Close integration with US financial system.",
            "instore_online":"Online 40% / In-store 60%. Interac Flash (contactless debit) very popular in-store. Interac Online for eCommerce.",
        },
        "Japan": {
            "flag":"🇯🇵","population":"125M","gdp":"$4.23T","ecom_size":"$180B (2025)",
            "payment_methods":[("Credit Card",40),("Convenience Store Payment (Konbini)",15),("Digital Wallets (PayPay, LinePay)",18),("Bank Transfer (Furikomi)",12),("Carrier Billing",8),("Cash",7)],
            "top_acquirers":"GMO Payment Gateway, SB Payment Service, Adyen, Stripe Japan, Paygent, DG Financial Technology",
            "regulatory":"FSA and Bank of Japan regulate. Cashless Vision targeting 40%+ cashless ratio (achieved ~39% in 2024). JCB is major domestic card scheme alongside Visa/MC. Installments (bunkatsu) culturally important.",
            "settlement":"Cards: D+15 to D+30 (monthly cycle typical). Bank transfer: D+0 to D+1. E-wallets: D+1 to D+15.",
            "fx":"JPY is freely traded. No FX restrictions. Japan has unique long settlement cycles for card payments that require planning.",
            "instore_online":"Online 38% / In-store 62%. IC card payments (Suica, PASMO) dominate transit/convenience. QR code wallets (PayPay) growing rapidly.",
        },
        "Singapore": {
            "flag":"🇸🇬","population":"5.9M","gdp":"$497B","ecom_size":"$9B (2025)",
            "payment_methods":[("Credit Card",38),("Debit Card",20),("PayNow (instant)",18),("Digital Wallets (GrabPay, ShopeePay)",15),("BNPL (Atome, Pace)",5),("NETS",4)],
            "top_acquirers":"Adyen, Stripe, 2C2P, Worldline, Rapyd, Nium",
            "regulatory":"MAS regulates payments. Payment Services Act (PS Act) covers all payment service providers. PayNow corporate expanding. SGQR unified QR standard. Strong AML framework.",
            "settlement":"PayNow: instant. Cards: D+1 to D+3. NETS: D+1. FAST: instant.",
            "fx":"SGD is managed float (MAS band). No FX restrictions. Singapore is Asian payments hub. Cross-border QR payments linked with Thailand (PromptPay) and India (UPI).",
            "instore_online":"Online 48% / In-store 52%. SGQR standardizes QR payments. Contactless card penetration > 85%.",
        },
        "Turkey": {
            "flag":"🇹🇷","population":"85M","gdp":"$1.11T","ecom_size":"$25B (2025)",
            "payment_methods":[("Credit Card (installments)",45),("Debit Card / BKM Express",20),("Bank Transfer (FAST/EFT)",18),("Digital Wallets",10),("Cash on Delivery",7)],
            "top_acquirers":"Iyzico (PayU), PayTR, Param, Craftgate, Stripe Turkey, Garanti BBVA",
            "regulatory":"CBRT and BRSA regulate payments. Installment payments deeply embedded (up to 12 months interest-free). BKM operates domestic card scheme (Troy). FAST instant payment system live.",
            "settlement":"Cards: D+1 to D+30 (installment cycles). FAST: instant. Bank transfer: D+0 to D+1.",
            "fx":"TRY is volatile. Central bank interventions frequent. FX-denominated pricing restricted for domestic transactions. Cross-border settlements can face delays. Withholding taxes apply.",
            "instore_online":"Online 35% / In-store 65%. Installment-based card payments are the norm for eCommerce. Troy card scheme growing domestically.",
        },
        "Egypt": {
            "flag":"🇪🇬","population":"106M","gdp":"$398B","ecom_size":"$8B (2025)",
            "payment_methods":[("Cash on Delivery",35),("Credit/Debit Card",22),("Mobile Wallets (Fawry, Vodafone Cash)",20),("Bank Transfer (InstaPay)",13),("Meeza (local scheme)",10)],
            "top_acquirers":"Fawry, Paymob, Accept (by Paymob), Kashier, ValU (BNPL), Amazon Payment Services",
            "regulatory":"CBE regulates payments. National Payment Council driving cashless. InstaPay instant transfer system. Meeza is the national card scheme. Mobile money regulations enabling wallet interoperability.",
            "settlement":"Cards: D+2 to D+7. InstaPay: instant. Mobile wallets: D+1 to D+2. COD: varies.",
            "fx":"EGP floated in 2022-2024 reforms. FX availability has improved but repatriation can still face delays. Letters of credit requirements eased. Central bank unified the exchange rate.",
            "instore_online":"Online 22% / In-store 78%. Cash on delivery remains dominant for eCommerce. Mobile wallet adoption accelerating among younger demographics.",
        },
    }

    if tab == "Country Profiles":
        country_name = st.selectbox("Select a market", list(_COUNTRIES.keys()), key="country_profile_select")
        c = _COUNTRIES[country_name]

        # Payment method breakdown chart
        fig = go.Figure(go.Bar(
            x=[pm[1] for pm in c["payment_methods"]],
            y=[pm[0] for pm in c["payment_methods"]],
            orientation='h',
            marker_color=['#4F46E5','#818cf8','#a5b4fc','#c7d2fe','#e0e7ff','#eef2ff'][:len(c["payment_methods"])],
            text=[f'{pm[1]}%' for pm in c["payment_methods"]],
            textposition='outside',
        ))
        fig.update_layout(
            height=220, margin=dict(l=0,r=40,t=10,b=10),
            xaxis=dict(title="Share (%)", range=[0,60], showgrid=True, gridcolor='rgba(0,0,0,0.05)'),
            yaxis=dict(autorange="reversed"),
            plot_bgcolor='rgba(0,0,0,0)',
            paper_bgcolor='rgba(0,0,0,0)',
            font=dict(size=11),
        )

        st.markdown(f'''<div style="background:#fff;box-shadow:0 1px 3px rgba(0,0,0,0.07),0 0 0 1px rgba(0,0,0,0.04);border-radius:12px;padding:20px 24px;margin-bottom:14px;">
  <div style="display:flex;align-items:center;gap:10px;margin-bottom:14px;">
    <span style="font-size:32px;">{c["flag"]}</span>
    <div>
      <div style="font-size:20px;font-weight:700;color:#1d1d1f;letter-spacing:-0.5px;">{country_name}</div>
      <div style="font-size:11px;color:#86868b;">Population: {c["population"]} &nbsp;·&nbsp; GDP: {c["gdp"]} &nbsp;·&nbsp; eCommerce: {c["ecom_size"]}</div>
    </div>
  </div>
  <div style="display:flex;gap:8px;flex-wrap:wrap;margin-bottom:6px;">
    <span style="font-size:9px;font-weight:700;padding:3px 8px;border-radius:5px;background:#EEF2FF;color:#4F46E5;">{c["instore_online"]}</span>
  </div>
</div>''', unsafe_allow_html=True)

        col1, col2 = st.columns(2)
        with col1:
            st.markdown('<div style="font-size:9px;font-weight:700;letter-spacing:1px;text-transform:uppercase;color:#86868b;margin-bottom:4px;">Payment Method Mix</div>', unsafe_allow_html=True)
            st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})

            st.markdown(f'''<div style="background:#fff;box-shadow:0 1px 3px rgba(0,0,0,0.07),0 0 0 1px rgba(0,0,0,0.04);border-radius:10px;padding:16px;margin-bottom:14px;">
  <div style="font-size:9px;font-weight:700;letter-spacing:1px;text-transform:uppercase;color:#86868b;margin-bottom:10px;">🏦 Top Acquirers / PSPs</div>
  <div style="font-size:12px;color:#1d1d1f;line-height:1.7;">{c["top_acquirers"]}</div>
</div>''', unsafe_allow_html=True)

        with col2:
            st.markdown(f'''<div style="background:#fff;box-shadow:0 1px 3px rgba(0,0,0,0.07),0 0 0 1px rgba(0,0,0,0.04);border-radius:10px;padding:16px;margin-bottom:14px;">
  <div style="font-size:9px;font-weight:700;letter-spacing:1px;text-transform:uppercase;color:#86868b;margin-bottom:10px;">📋 Regulatory Highlights</div>
  <div style="font-size:12px;color:#1d1d1f;line-height:1.7;">{c["regulatory"]}</div>
</div>''', unsafe_allow_html=True)

            st.markdown(f'''<div style="background:#fff;box-shadow:0 1px 3px rgba(0,0,0,0.07),0 0 0 1px rgba(0,0,0,0.04);border-radius:10px;padding:16px;margin-bottom:14px;">
  <div style="font-size:9px;font-weight:700;letter-spacing:1px;text-transform:uppercase;color:#86868b;margin-bottom:10px;">⏱ Settlement Times</div>
  <div style="font-size:12px;color:#1d1d1f;line-height:1.7;">{c["settlement"]}</div>
</div>''', unsafe_allow_html=True)

            st.markdown(f'''<div style="background:#fff;box-shadow:0 1px 3px rgba(0,0,0,0.07),0 0 0 1px rgba(0,0,0,0.04);border-radius:10px;padding:16px;margin-bottom:14px;">
  <div style="font-size:9px;font-weight:700;letter-spacing:1px;text-transform:uppercase;color:#86868b;margin-bottom:10px;">💱 FX Considerations</div>
  <div style="font-size:12px;color:#1d1d1f;line-height:1.7;">{c["fx"]}</div>
</div>''', unsafe_allow_html=True)

    elif tab == "Payment Trends":
        _TRENDS = [
            {
                "icon":"⚡","title":"Rise of Account-to-Account (A2A) Instant Payments",
                "body":"PIX (Brazil) now processes 4B+ monthly transactions. UPI (India) exceeds 15B/month. FPS (UK) and AANI (UAE) expanding merchant acceptance. A2A payments offer lower costs (0-0.5% vs 1.5-3% for cards), instant settlement, and zero chargeback risk. Projected to capture 25% of global eCommerce payments by 2028.",
                "stats":"PIX: 4B txns/mo · UPI: 15B txns/mo · FPS: 400M txns/mo",
                "tag":"Instant Payments","tag_color":"#059669","tag_bg":"#D1FAE5",
            },
            {
                "icon":"🛒","title":"BNPL Adoption Varies Sharply by Region",
                "body":"BNPL represents ~5% of global eCommerce but adoption varies widely. Nordic countries lead at 12-15%, followed by Australia (10%), UK (8%), and Germany (7% via invoice model). In MENA, Tamara and Tabby are growing 80% YoY. LATAM adoption is 3-4% but growing fast. Regulatory scrutiny increasing globally — UK and EU introducing BNPL-specific frameworks.",
                "stats":"Global BNPL: $350B GMV · MENA growth: 80% YoY · EU regulation pending",
                "tag":"BNPL","tag_color":"#7C3AED","tag_bg":"#EDE9FE",
            },
            {
                "icon":"🌐","title":"Cross-Border Payment Costs Declining",
                "body":"Average cost of sending $200 cross-border has dropped from 6.3% (2020) to 4.1% (2026). Swift GPI now delivers 50% of payments in under 30 minutes. Visa Direct and Mastercard Send enabling real-time cross-border. Wise, Airwallex, and Nium offering 0.5-1.5% FX margins vs 2-4% from traditional banks. G20 target of 3% average cost by 2027 on track.",
                "stats":"Avg cost: 4.1% (down from 6.3%) · Swift GPI: 50% <30min · G20 target: 3%",
                "tag":"Cross-Border","tag_color":"#0369A1","tag_bg":"#E0F2FE",
            },
            {
                "icon":"📱","title":"Mobile Wallet Penetration Accelerating",
                "body":"Global mobile wallet users projected at 5.2B by 2027. In APAC, wallets already account for 60%+ of eCommerce payments (Alipay, WeChat Pay, GCash, GrabPay). In Africa, M-Pesa and MTN Mobile Money process $1.2B daily. Apple Pay and Google Pay growing in developed markets. Super-app wallets (Grab, GoTo, Rappi) blurring lines between payments and commerce.",
                "stats":"5.2B wallet users by 2027 · APAC: 60%+ eCom share · Africa: $1.2B/day",
                "tag":"Mobile Wallets","tag_color":"#EA580C","tag_bg":"#FFF7ED",
            },
            {
                "icon":"₿","title":"Crypto Payment Adoption in Emerging Markets",
                "body":"Stablecoin transaction volume surpassed $10T in 2025. USDT dominates in LATAM, Africa, and Southeast Asia as a remittance and savings tool. Merchants in Argentina, Nigeria, and Turkey increasingly accepting stablecoins to hedge local currency volatility. Stripe and PayPal now support USDC settlements. CBDCs progressing: eNaira (live), Digital Real (pilot), Digital Rupee (pilot).",
                "stats":"Stablecoin vol: $10T+ (2025) · eNaira: live · Digital Real: pilot phase",
                "tag":"Crypto/Stablecoins","tag_color":"#CA8A04","tag_bg":"#FEF9C3",
            },
            {
                "icon":"📶","title":"Contactless & Tap-to-Phone In-Store Growth",
                "body":"Contactless penetration exceeds 75% of in-store card transactions in Europe and MENA. Tap-to-phone (SoftPOS) eliminating need for dedicated terminals — Stripe Terminal, Adyen Tap to Pay, and Square enabling any smartphone as a POS. In-store/online convergence accelerating via unified commerce platforms. Apple Tap to Pay expanding to 15+ markets.",
                "stats":"EU contactless: 75%+ · Tap-to-phone: 15+ markets · SoftPOS reducing terminal costs by 60%",
                "tag":"In-Store Innovation","tag_color":"#BE185D","tag_bg":"#FCE7F3",
            },
        ]

        for t in _TRENDS:
            st.markdown(f'''<div style="background:#fff;box-shadow:0 1px 3px rgba(0,0,0,0.07),0 0 0 1px rgba(0,0,0,0.04);border-radius:10px;padding:16px 18px;margin-bottom:10px;">
  <div style="display:flex;align-items:flex-start;gap:12px;">
    <span style="font-size:24px;flex-shrink:0;margin-top:2px;">{t["icon"]}</span>
    <div style="flex:1;">
      <div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:6px;">
        <div style="font-size:14px;font-weight:700;color:#0f172a;">{t["title"]}</div>
        <span style="font-size:9px;font-weight:700;padding:2px 8px;border-radius:4px;background:{t["tag_bg"]};color:{t["tag_color"]};white-space:nowrap;">{t["tag"]}</span>
      </div>
      <div style="font-size:12px;color:#64748b;line-height:1.65;margin-bottom:8px;">{t["body"]}</div>
      <div style="font-size:10px;font-family:'Menlo',monospace;color:#4F46E5;background:#F8FAFC;padding:6px 10px;border-radius:6px;">{t["stats"]}</div>
    </div>
  </div>
</div>''', unsafe_allow_html=True)

    elif tab == "Regulatory Updates":
        _REGULATIONS = [
            {
                "icon":"🇪🇺","region":"Europe","title":"EU PSD3 & Payment Services Regulation (PSR)",
                "status":"Draft Published — Expected adoption 2026-2027",
                "status_color":"#B45309","status_bg":"#FEF3C7",
                "body":"PSD3 introduces stronger SCA exemptions for low-risk transactions under EUR 500 with merchant Transaction Risk Analysis (TRA). The PSR (a directly applicable regulation replacing PSD2 directives) will harmonize rules across all EU member states. Key changes: expanded open banking access, fraud liability shifts, non-bank PSP access to payment systems, and enhanced consumer rights. Merchants can expect 3-5% auth rate improvement from relaxed SCA and lower friction at checkout.",
            },
            {
                "icon":"🇧🇷","region":"LATAM","title":"Brazil PIX Credit Launch & Open Finance Expansion",
                "status":"Launching Q2 2026",
                "status_color":"#059669","status_bg":"#D1FAE5",
                "body":"The Central Bank of Brazil (BCB) is launching PIX Credit, allowing consumers to make PIX payments using credit lines from their banks. This effectively creates an instant BNPL mechanism on existing PIX rails. Expected to increase eCommerce conversion by 15-20% by combining PIX simplicity with installment flexibility. Additionally, Open Finance Phase 4 expands data sharing to investments and insurance. Merchants should prepare checkout flows to support PIX Credit alongside standard PIX.",
            },
            {
                "icon":"🇮🇳","region":"APAC","title":"India MDR Regulation & Digital Payment Incentives",
                "status":"Under Review — Policy update expected H2 2026",
                "status_color":"#1D4ED8","status_bg":"#DBEAFE",
                "body":"RBI maintains zero MDR on UPI and RuPay debit transactions, subsidized via government reimbursement to banks (~$350M annually). The sustainability of this model is under review. Industry bodies are lobbying for a small MDR (0.15-0.3%) to fund infrastructure investment. Separately, RBI's e-mandate framework now supports recurring UPI payments up to INR 1 lakh. The Digital Personal Data Protection Act (2023) impacts how payment data can be stored and processed, with compliance deadlines in 2026.",
            },
            {
                "icon":"🇸🇦","region":"MENA","title":"Saudi SAMA Mada Requirements & Open Banking",
                "status":"Mandatory — Enforcement active",
                "status_color":"#DC2626","status_bg":"#FEE2E2",
                "body":"SAMA mandates that all domestic debit card transactions route through the Mada scheme. By Q3 2026, merchants without Mada certification will see 100% decline rates on local debit cards (which represent ~45% of all card payments). Open Banking framework launched with 8 banks in phase 1. BNPL regulation formalized — all BNPL providers must be SAMA-licensed. Vision 2030 cashless target of 70% non-cash transactions is driving aggressive fintech licensing and payment modernization.",
            },
            {
                "icon":"🇳🇬","region":"Africa","title":"Nigeria eNaira & CBN Cashless Policy Updates",
                "status":"In Effect — Evolving framework",
                "status_color":"#7C3AED","status_bg":"#EDE9FE",
                "body":"The eNaira (CBDC) is live but adoption remains modest (~$80M in circulation). CBN continues enforcing cashless policy with withdrawal limits (NGN 500K/individual, NGN 5M/corporate per week). Cash processing fees of 5-10% apply above limits. The NIBSS Instant Payment (NIP) system processes 1B+ transactions annually. CBN has issued new Payment Service Bank (PSB) licenses expanding mobile money. New regulations require payment processors to maintain local data storage and obtain CBN licensing.",
            },
            {
                "icon":"🌎","region":"LATAM","title":"Data Localization Trends Across Latin America",
                "status":"Varies by Country — Monitoring recommended",
                "status_color":"#475569","status_bg":"#F1F5F9",
                "body":"Brazil (LGPD) does not mandate data localization but requires consent for cross-border transfers and BCB mandates local payment data processing. Mexico (Ley Fintech) requires financial data to be accessible to Mexican regulators. Colombia requires certain financial records to remain accessible domestically. Argentina's data protection law requires adequate protection for cross-border transfers. Chile is updating its data protection framework (GDPR-aligned). Merchants operating across LATAM should evaluate whether a regional data center strategy is needed.",
            },
        ]

        for r in _REGULATIONS:
            st.markdown(f'''<div style="background:#fff;box-shadow:0 1px 3px rgba(0,0,0,0.07),0 0 0 1px rgba(0,0,0,0.04);border-radius:10px;padding:16px 18px;margin-bottom:10px;">
  <div style="display:flex;align-items:flex-start;gap:12px;">
    <span style="font-size:24px;flex-shrink:0;margin-top:2px;">{r["icon"]}</span>
    <div style="flex:1;">
      <div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:4px;">
        <div style="font-size:14px;font-weight:700;color:#0f172a;">{r["title"]}</div>
        <span style="font-size:9px;font-weight:700;padding:2px 8px;border-radius:4px;background:#EEF2FF;color:#4F46E5;white-space:nowrap;">{r["region"]}</span>
      </div>
      <div style="margin-bottom:8px;"><span style="font-size:10px;font-weight:700;padding:2px 8px;border-radius:4px;background:{r["status_bg"]};color:{r["status_color"]};">{r["status"]}</span></div>
      <div style="font-size:12px;color:#64748b;line-height:1.65;">{r["body"]}</div>
    </div>
  </div>
</div>''', unsafe_allow_html=True)

    elif tab == "Market News":
        _NEWS = [
            {"date":"Mar 22, 2026","region":"LATAM","title":"Brazil Central Bank announces PIX credit launch for Q2 2026","body":"PIX credit will allow merchants to offer instant BNPL at checkout via PIX rails — expected to increase conversion by 15-20% for Brazilian merchants. Cielo, PagBank, and Stone are early adopters.","impact":"High","tag":"Regulation"},
            {"date":"Mar 20, 2026","region":"MENAT","title":"Saudi Arabia mandates Mada scheme for all local transactions","body":"SAMA has confirmed that all domestic card transactions must route through Mada by Q3 2026. Merchants without Mada support will see 100% decline rates on local cards. Checkout.com and Tap Payments are certified.","impact":"Critical","tag":"Compliance"},
            {"date":"Mar 18, 2026","region":"APAC","title":"India UPI crosses 15 billion monthly transactions","body":"UPI volume continues to grow at 40% YoY. RazorPay and Cashfree are leading PSPs for UPI acceptance. Google Pay and PhonePe dominate market share.","impact":"Medium","tag":"Market Growth"},
            {"date":"Mar 15, 2026","region":"Europe","title":"EU PSD3 draft published — stronger SCA exemptions for low-risk transactions","body":"The draft PSD3 regulation proposes wider SCA exemptions for transactions under EUR 500 with merchant TRA. This could improve auth rates by 3-5% for European merchants using Adyen and Stripe.","impact":"High","tag":"Regulation"},
            {"date":"Mar 12, 2026","region":"Global","title":"Visa and Mastercard announce 2026 fee increases for cross-border transactions","body":"Interchange fees for cross-border transactions increasing by 8-12 basis points effective July 2026. Local acquiring becomes even more critical for cost optimization.","impact":"High","tag":"Pricing"},
            {"date":"Mar 10, 2026","region":"LATAM","title":"Colombia launches instant payment system — Transfiya 2.0","body":"Banco de la Republica upgrades Transfiya with merchant acceptance capabilities. PSE alternative expected to reduce costs and increase real-time settlement options.","impact":"Medium","tag":"Infrastructure"},
            {"date":"Mar 8, 2026","region":"MENAT","title":"UAE Central Bank launches AANI instant payments nationwide","body":"AANI is now available across all UAE banks. Merchants can accept instant A2A payments with zero chargeback risk. Integration via Checkout.com and Tap Payments.","impact":"High","tag":"Infrastructure"},
            {"date":"Mar 5, 2026","region":"Global","title":"Apple Pay Later discontinued — shifts focus to issuer installments","body":"Apple confirms BNPL product sunset. Market opportunity for Tamara, Tabby, Klarna, and local BNPL providers to fill the gap.","impact":"Medium","tag":"Competitive"},
        ]
        _impact_style = {"Critical":("#DC2626","#FEE2E2"),"High":("#B45309","#FEF3C7"),"Medium":("#1D4ED8","#DBEAFE")}

        for n in _NEWS:
            ic, ibg = _impact_style.get(n["impact"], ("#475569","#F1F5F9"))
            st.markdown(f'<div style="background:#fff;border:0.5px solid #E2E8F0;border-radius:10px;padding:16px 18px;margin-bottom:8px;">'
                        f'<div style="display:flex;justify-content:space-between;align-items:flex-start;margin-bottom:6px;">'
                        f'<div style="flex:1;"><div style="font-size:14px;font-weight:700;color:#0f172a;line-height:1.4;margin-bottom:4px;">{n["title"]}</div>'
                        f'<div style="display:flex;gap:6px;align-items:center;">'
                        f'<span style="font-size:10px;color:#94a3b8;">{n["date"]}</span>'
                        f'<span style="font-size:9px;font-weight:700;padding:2px 7px;border-radius:4px;background:#EEF2FF;color:#4F46E5;">{n["region"]}</span>'
                        f'<span style="font-size:9px;font-weight:700;padding:2px 7px;border-radius:4px;background:#F1F5F9;color:#475569;">{n["tag"]}</span>'
                        f'<span style="font-size:9px;font-weight:700;padding:2px 7px;border-radius:4px;background:{ibg};color:{ic};">{n["impact"]}</span>'
                        f'</div></div></div>'
                        f'<div style="font-size:12px;color:#64748b;line-height:1.6;">{n["body"]}</div>'
                        f'</div>', unsafe_allow_html=True)

# ── Merchants ──────────────────────────────────────────────────────────────────
def show_merchants():
    # name, country, region, tpv_m, aov, transactions, auth_rate
    MERCHANTS = [
        ("Rappi",        "Colombia",  "LATAM",     8.42,  32.5,  259077, 91.2),
        ("MercadoLibre", "Argentina", "LATAM",     7.31, 118.4,   61739, 88.9),
        ("Uber",         "Mexico",    "LATAM",     5.89,  18.2,  323626, 93.4),
        ("Netflix",      "USA",       "N. America",4.75,  15.5,  306452, 95.6),
        ("Spotify",      "Brazil",    "LATAM",     4.21,  10.0,  421421, 94.8),
        ("Despegar",     "Argentina", "LATAM",     3.84, 312.5,   12288, 87.5),
        ("Falabella",    "Chile",     "LATAM",     3.12,  87.3,   35739, 88.2),
        ("iFood",        "Brazil",    "LATAM",     2.93,  28.4,  103169, 91.8),
        ("Cinépolis",    "Mexico",    "LATAM",     2.14,  24.5,   62857, 92.4),
        ("PedidosYa",    "Argentina", "LATAM",     1.87,  35.8,   62570, 90.1),
        ("Claro",        "Mexico",    "LATAM",     1.52,  22.8,   49123, 89.7),
        ("Linio",        "Colombia",  "LATAM",     1.11,  72.1,   25936, 87.1),
    ]

    st.markdown("""
<div style="display:flex;align-items:baseline;gap:12px;margin-bottom:4px;">
  <span style="font-size:28px;font-weight:700;color:#1d1d1f;">Merchants</span>
</div>
<div style="font-size:11px;color:#86868b;margin-bottom:18px;">Your Merchants</div>
""", unsafe_allow_html=True)

    stat_row([
        {"label":"Merchants Live",    "value":"12",     "delta":"Active & connected",    "delta_type":"flat"},
        {"label":"Total TPV",         "value":"$47.24M","delta":"This month · USD",       "delta_type":"up","val_color":"#818cf8"},
        {"label":"Avg AOV",           "value":"$64.8",  "delta":"Across all merchants",   "delta_type":"flat"},
        {"label":"Avg Approval Rate", "value":"90.2%",  "delta":"↑ +1.2pp vs last mo.", "delta_type":"up"},
    ])

    st.markdown("<div style='height:18px'></div>", unsafe_allow_html=True)

    # Row 1: TPV + Approval Rate
    col_l, col_r = st.columns(2)
    with col_l:
        data = sorted(MERCHANTS, key=lambda x: x[3], reverse=True)
        fig = go.Figure(go.Bar(
            x=[r[3] for r in data], y=[r[0] for r in data], orientation="h",
            marker_color="#6366f1",
            text=[f"${r[3]}M" for r in data], textposition="outside",
            textfont=dict(size=10, color="#6e6e73", family="Menlo"),
        ))
        fig.update_layout(
            title=dict(text="TPV by Merchant", font=dict(size=12, color="#1d1d1f"), x=0),
            paper_bgcolor="#ffffff", plot_bgcolor="#ffffff",
            font=dict(color="#6e6e73", family="Inter"),
            margin=dict(l=10, r=70, t=40, b=10),
            xaxis=dict(showgrid=False, showticklabels=False, zeroline=False),
            yaxis=dict(autorange="reversed", tickfont=dict(size=11)), height=380,
        )
        st.plotly_chart(fig, use_container_width=True)

    with col_r:
        data = sorted(MERCHANTS, key=lambda x: x[6], reverse=True)
        colors = ["#22c55e" if r[6] >= 92 else "#f59e0b" if r[6] >= 89 else "#ef4444" for r in data]
        fig = go.Figure(go.Bar(
            x=[r[6] for r in data], y=[r[0] for r in data], orientation="h",
            marker_color=colors,
            text=[f"{r[6]}%" for r in data], textposition="outside",
            textfont=dict(size=10, color="#6e6e73", family="Menlo"),
        ))
        fig.update_layout(
            title=dict(text="Approval Rate by Merchant", font=dict(size=12, color="#1d1d1f"), x=0),
            paper_bgcolor="#ffffff", plot_bgcolor="#ffffff",
            font=dict(color="#6e6e73", family="Inter"),
            margin=dict(l=10, r=60, t=40, b=10),
            xaxis=dict(showgrid=False, showticklabels=False, zeroline=False, range=[80, 100]),
            yaxis=dict(autorange="reversed", tickfont=dict(size=11)), height=380,
        )
        st.plotly_chart(fig, use_container_width=True)

    # Row 2: AOV + Total Transactions
    col_l2, col_r2 = st.columns(2)
    with col_l2:
        data = sorted(MERCHANTS, key=lambda x: x[4], reverse=True)
        fig = go.Figure(go.Bar(
            x=[r[4] for r in data], y=[r[0] for r in data], orientation="h",
            marker_color="#818cf8",
            text=[f"${r[4]}" for r in data], textposition="outside",
            textfont=dict(size=10, color="#6e6e73", family="Menlo"),
        ))
        fig.update_layout(
            title=dict(text="Average Order Value (AOV)", font=dict(size=12, color="#1d1d1f"), x=0),
            paper_bgcolor="#ffffff", plot_bgcolor="#ffffff",
            font=dict(color="#6e6e73", family="Inter"),
            margin=dict(l=10, r=70, t=40, b=10),
            xaxis=dict(showgrid=False, showticklabels=False, zeroline=False),
            yaxis=dict(autorange="reversed", tickfont=dict(size=11)), height=380,
        )
        st.plotly_chart(fig, use_container_width=True)

    with col_r2:
        data = sorted(MERCHANTS, key=lambda x: x[5], reverse=True)
        fig = go.Figure(go.Bar(
            x=[r[5] for r in data], y=[r[0] for r in data], orientation="h",
            marker_color="#4f46e5",
            text=[f"{r[5]:,}" for r in data], textposition="outside",
            textfont=dict(size=10, color="#6e6e73", family="Menlo"),
        ))
        fig.update_layout(
            title=dict(text="Total Transactions", font=dict(size=12, color="#1d1d1f"), x=0),
            paper_bgcolor="#ffffff", plot_bgcolor="#ffffff",
            font=dict(color="#6e6e73", family="Inter"),
            margin=dict(l=10, r=80, t=40, b=10),
            xaxis=dict(showgrid=False, showticklabels=False, zeroline=False),
            yaxis=dict(autorange="reversed", tickfont=dict(size=11)), height=380,
        )
        st.plotly_chart(fig, use_container_width=True)

    # Merchant Summary table
    def auth_bar(v):
        color = "#22c55e" if v >= 92 else "#f59e0b" if v >= 89 else "#ef4444"
        pct = int((v - 80) / 20 * 100)
        return f'<div style="display:flex;align-items:center;gap:8px;"><div style="flex:1;background:#f0f0f5;border-radius:2px;height:4px;"><div style="width:{pct}%;height:4px;background:{color};border-radius:2px;"></div></div><span style="font-size:11px;font-family:monospace;color:{color};min-width:38px;">{v}%</span></div>'

    rows_html = ""
    for i, r in enumerate(MERCHANTS):
        border = "" if i == len(MERCHANTS)-1 else "border-bottom:1px solid rgba(0,0,0,0.05);"
        rows_html += f"""
  <div style="display:grid;grid-template-columns:1.5fr 1fr 1fr 0.9fr 0.8fr 1fr 1.6fr;padding:9px 16px;{border}align-items:center;">
    <span style="font-size:12px;font-weight:600;color:#1d1d1f;">{r[0]}</span>
    <span style="font-size:11px;color:#6e6e73;">{r[1]}</span>
    <span style="font-size:11px;color:#6e6e73;">{r[2]}</span>
    <span style="font-size:11px;font-family:monospace;color:#1d1d1f;">${r[3]}M</span>
    <span style="font-size:11px;font-family:monospace;color:#6e6e73;">${r[4]}</span>
    <span style="font-size:11px;font-family:monospace;color:#6e6e73;">{r[5]:,}</span>
    {auth_bar(r[6])}
  </div>"""

    st.markdown(f"""
<div style="font-size:13px;font-weight:700;color:#1d1d1f;margin:20px 0 10px;">Merchant Summary</div>
<div style="background:#fff;box-shadow:0 1px 3px rgba(0,0,0,0.07),0 0 0 1px rgba(0,0,0,0.04);border:none;border-radius:10px;overflow:hidden;">
  <div style="display:grid;grid-template-columns:1.5fr 1fr 1fr 0.9fr 0.8fr 1fr 1.6fr;padding:9px 16px;border-bottom:1px solid rgba(0,0,0,0.07);font-size:9.5px;font-weight:700;letter-spacing:.6px;text-transform:uppercase;color:#86868b;">
    <span>Merchant</span><span>Country</span><span>Region</span><span>TPV (M USD)</span><span>AOV (USD)</span><span>Transactions</span><span>Approval Rate</span>
  </div>
  {rows_html}
</div>
""", unsafe_allow_html=True)

# ── Benchmarks ─────────────────────────────────────────────────────────────────
def show_benchmarks():
    # ── Rev Share Dashboard ───────────────────────────────────────────────────
    st.markdown("""
<div style="background:#0C1220;border-radius:16px;padding:24px 28px;margin-bottom:14px;position:relative;overflow:hidden;">
  <div style="position:absolute;inset:0;background-image:radial-gradient(circle at 1px 1px,rgba(255,255,255,0.03) 1px,transparent 0);background-size:22px 22px;pointer-events:none;"></div>
  <div style="position:relative;z-index:1;">
    <div style="font-size:22px;font-weight:800;color:#F1F5F9;letter-spacing:-0.02em;margin-bottom:4px;">Revenue Share</div>
    <div style="font-size:12px;color:#64748b;margin-bottom:14px;">Partner rev share earned by Yuno — Feb 2026</div>
    <div style="display:flex;gap:10px;flex-wrap:wrap;">
      <div style="background:rgba(255,255,255,0.04);border:0.5px solid rgba(255,255,255,0.08);border-radius:8px;padding:10px 14px;text-align:center;"><div style="font-size:18px;font-weight:800;color:#F9A8D4;">$732K</div><div style="font-size:9px;color:#475569;text-transform:uppercase;margin-top:2px;">2024 Total</div></div>
      <div style="background:rgba(255,255,255,0.04);border:0.5px solid rgba(255,255,255,0.08);border-radius:8px;padding:10px 14px;text-align:center;"><div style="font-size:18px;font-weight:800;color:#6EE7B7;">$123K</div><div style="font-size:9px;color:#475569;text-transform:uppercase;margin-top:2px;">Feb 2026</div></div>
      <div style="background:rgba(255,255,255,0.04);border:0.5px solid rgba(255,255,255,0.08);border-radius:8px;padding:10px 14px;text-align:center;"><div style="font-size:18px;font-weight:800;color:#818cf8;">$1.48M</div><div style="font-size:9px;color:#475569;text-transform:uppercase;margin-top:2px;">ARR</div></div>
      <div style="background:rgba(255,255,255,0.04);border:0.5px solid rgba(255,255,255,0.08);border-radius:8px;padding:10px 14px;text-align:center;"><div style="font-size:18px;font-weight:800;color:#FCD34D;">$245K</div><div style="font-size:9px;color:#475569;text-transform:uppercase;margin-top:2px;">Peak (Oct 25)</div></div>
      <div style="background:rgba(255,255,255,0.04);border:0.5px solid rgba(255,255,255,0.08);border-radius:8px;padding:10px 14px;text-align:center;"><div style="font-size:18px;font-weight:800;color:#C7D2FE;">$1.78M</div><div style="font-size:9px;color:#475569;text-transform:uppercase;margin-top:2px;">2025 Total</div></div>
      <div style="background:rgba(255,255,255,0.04);border:0.5px solid rgba(255,255,255,0.08);border-radius:8px;padding:10px 14px;text-align:center;"><div style="font-size:18px;font-weight:800;color:#6EE7B7;">+33%</div><div style="font-size:9px;color:#475569;text-transform:uppercase;margin-top:2px;">vs Feb 25</div></div>
    </div>
  </div>
</div>""", unsafe_allow_html=True)

    # ── Rev Share by Region (bar chart) + Top Partners table ──────────────────
    _rev_regions = [
        ("LATAM", 47600, 42, "#4F46E5"),
        ("Brasil", 32900, 29, "#10B981"),
        ("Global", 28300, 25, "#3B82F6"),
        ("APAC", 3300, 3, "#0EA5E9"),
        ("N. America", 400, 0, "#059669"),
    ]
    _top_partners = [
        ("DLocal", "$25,338"),
        ("Stripe", "$21,705"),
        ("Cielo", "$17,635"),
        ("Bamboo", "$9,491"),
        ("Mercado Pago", "$7,365"),
        ("Unlimint", "$5,949"),
        ("PagBank", "$4,516"),
        ("Nuvei/Paymentez", "$3,772"),
        ("Pagar.me", "$3,441"),
        ("PicPay", "$3,180"),
    ]

    col_chart, col_table = st.columns([1.5, 1])

    with col_chart:
        # Region bars
        max_rev = max(r[1] for r in _rev_regions)
        bars_html = ""
        for region, rev, pct, color in _rev_regions:
            w = int(rev / max_rev * 100)
            bars_html += (f'<div style="display:flex;align-items:center;gap:10px;margin-bottom:10px;">'
                          f'<span style="font-size:12px;color:#64748b;min-width:80px;text-align:right;">{region}</span>'
                          f'<div style="flex:1;height:28px;background:#F1F5F9;border-radius:6px;overflow:hidden;position:relative;">'
                          f'<div style="width:{w}%;height:100%;background:{color};border-radius:6px;display:flex;align-items:center;padding-left:10px;">'
                          f'<span style="font-size:11px;font-weight:700;color:#fff;">${rev/1000:.1f}K</span></div></div>'
                          f'<span style="font-size:11px;font-weight:600;color:#64748b;min-width:30px;">{pct}%</span></div>')
        st.markdown(f'<div style="background:#fff;border:0.5px solid #E2E8F0;border-radius:12px;padding:20px;">'
                    f'<div style="font-size:14px;font-weight:700;color:#0f172a;margin-bottom:14px;">Rev Share by Region</div>'
                    f'{bars_html}</div>', unsafe_allow_html=True)

    with col_table:
        # Top partners table
        rows_html = ""
        for i, (name, amount) in enumerate(_top_partners):
            bg = "#FAFBFF" if i % 2 == 0 else "#fff"
            rows_html += (f'<tr style="background:{bg};">'
                          f'<td style="padding:8px 12px;font-size:12px;color:#94a3b8;font-weight:600;">{i+1}</td>'
                          f'<td style="padding:8px 12px;font-size:13px;font-weight:700;color:#0f172a;">{name}</td>'
                          f'<td style="padding:8px 12px;font-size:13px;font-weight:700;color:#065F46;font-family:monospace;text-align:right;">{amount}</td></tr>')
        st.markdown(f'<div style="background:#fff;border:0.5px solid #E2E8F0;border-radius:12px;overflow:hidden;">'
                    f'<div style="padding:16px 16px 8px;font-size:14px;font-weight:700;color:#0f172a;">Top Partners — Feb 2026</div>'
                    f'<table style="width:100%;border-collapse:collapse;">'
                    f'<thead><tr><th style="padding:6px 12px;font-size:9px;font-weight:700;letter-spacing:0.08em;text-transform:uppercase;color:#94a3b8;text-align:left;border-bottom:1px solid #E2E8F0;">#</th>'
                    f'<th style="padding:6px 12px;font-size:9px;font-weight:700;letter-spacing:0.08em;text-transform:uppercase;color:#94a3b8;text-align:left;border-bottom:1px solid #E2E8F0;">Partner</th>'
                    f'<th style="padding:6px 12px;font-size:9px;font-weight:700;letter-spacing:0.08em;text-transform:uppercase;color:#94a3b8;text-align:right;border-bottom:1px solid #E2E8F0;">Rev Share</th></tr></thead>'
                    f'<tbody>{rows_html}</tbody></table></div>', unsafe_allow_html=True)

    # ── Monthly trend chart ───────────────────────────────────────────────────
    _monthly = [("Nov 24",61),("Dec 24",61),("Jan 25",122),("Feb 25",122),("Mar 25",183),("Apr 25",183),("May 25",144),("Jun 25",145),("Jul 25",166),("Aug 25",219),("Sep 25",245),("Oct 25",245),("Nov 25",178),("Dec 25",161),("Jan 26",123),("Feb 26",123)]
    max_m = max(v for _, v in _monthly)
    trend_bars = ""
    for month, val in _monthly:
        h = int(val / max_m * 100)
        is_peak = val == max_m
        color = "#F59E0B" if is_peak else "#5B5FDE"
        trend_bars += (f'<div style="display:flex;flex-direction:column;align-items:center;flex:1;">'
                       f'<div style="font-size:9px;font-weight:600;color:#64748b;margin-bottom:2px;">${val}K</div>'
                       f'<div style="width:100%;height:120px;display:flex;align-items:flex-end;">'
                       f'<div style="width:100%;height:{h}%;background:{color};border-radius:4px 4px 0 0;min-height:4px;"></div></div>'
                       f'<div style="font-size:8px;color:#94a3b8;margin-top:4px;white-space:nowrap;">{month}</div></div>')
    st.markdown(f'<div style="background:#fff;border:0.5px solid #E2E8F0;border-radius:12px;padding:20px;margin-top:14px;">'
                f'<div style="font-size:14px;font-weight:700;color:#0f172a;margin-bottom:14px;">Monthly Revenue Share Trend</div>'
                f'<div style="display:flex;gap:4px;align-items:flex-end;">{trend_bars}</div></div>', unsafe_allow_html=True)

    return  # Rev Share section complete — benchmarks removed

    months = ["Aug/24","Sep/24","Oct/24","Nov/24","Dec/24","Jan/25"]
    st.markdown("", unsafe_allow_html=True)

    stat_row([
        {"label":"Your Avg Approval Rate", "value":"90.2%",  "delta":"↑ +1.2pp vs last mo.", "delta_type":"up"},
        {"label":"Network Avg",            "value":"87.6%",  "delta":"→ stable",              "delta_type":"flat"},
        {"label":"Top Performer",          "value":"95.6%",  "delta":"↑ Netflix · Streaming", "delta_type":"up", "val_color":"#22c55e"},
        {"label":"Your Recovery Rate",     "value":"4.1%",   "delta":"↑ +0.4pp MoM",          "delta_type":"up"},
        {"label":"Network Recovery Avg",   "value":"3.8%",   "delta":"→ stable",              "delta_type":"flat"},
    ])

    st.markdown("<div style='height:10px'></div>", unsafe_allow_html=True)

    col_l, col_r = st.columns(2)

    with col_l:
        yuno_approval   = [68.2, 67.1, 65.8, 64.3, 62.1, 58.4]
        peer2_approval  = [95.1, 96.2, 96.8, 96.4, 96.9, 97.2]
        peer3_approval  = [82.4, 83.1, 82.7, 82.9, 82.1, 81.8]
        fig = go.Figure()
        fig.add_trace(go.Bar(name="You (Yuno)", x=months, y=yuno_approval,  marker_color="#6366f1", text=[f"{v}%" for v in yuno_approval],  textposition="outside", textfont=dict(size=9, color="#6e6e73")))
        fig.add_trace(go.Bar(name="Peer Avg",   x=months, y=peer2_approval, marker_color="#14b8a6", text=[f"{v}%" for v in peer2_approval], textposition="outside", textfont=dict(size=9, color="#6e6e73")))
        fig.add_trace(go.Bar(name="Top Peer",   x=months, y=peer3_approval, marker_color="#f59e0b", text=[f"{v}%" for v in peer3_approval], textposition="outside", textfont=dict(size=9, color="#6e6e73")))
        fig.update_layout(
            title=dict(text="Approval Rate Benchmark", font=dict(size=12, color="#1d1d1f"), x=0),
            barmode="group", paper_bgcolor="#ffffff", plot_bgcolor="#ffffff",
            font=dict(color="#6e6e73", family="Inter"),
            margin=dict(l=10, r=10, t=40, b=30),
            legend=dict(orientation="h", y=-0.15, font=dict(size=10,color="#6e6e73"), bgcolor="rgba(0,0,0,0)"),
            yaxis=dict(gridcolor="rgba(0,0,0,0.06)", ticksuffix="%", range=[0,105]),
            xaxis=dict(gridcolor="rgba(0,0,0,0.06)"),
            height=320,
        )
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})

    with col_r:
        yuno_recovery  = [0.9, 0.8, 0.9, 0.7, 0.6, 0.5]
        peer2_recovery = [4.8, 5.1, 5.0, 5.3, 2.7, 3.9]
        peer3_recovery = [3.8, 3.9, 4.1, 4.3, 5.2, 2.6]
        fig2 = go.Figure()
        fig2.add_trace(go.Bar(name="You (Yuno)", x=months, y=yuno_recovery,  marker_color="#6366f1", text=[f"{v}%" for v in yuno_recovery],  textposition="outside", textfont=dict(size=9, color="#6e6e73")))
        fig2.add_trace(go.Bar(name="Peer Avg",   x=months, y=peer2_recovery, marker_color="#14b8a6", text=[f"{v}%" for v in peer2_recovery], textposition="outside", textfont=dict(size=9, color="#6e6e73")))
        fig2.add_trace(go.Bar(name="Top Peer",   x=months, y=peer3_recovery, marker_color="#f59e0b", text=[f"{v}%" for v in peer3_recovery], textposition="outside", textfont=dict(size=9, color="#6e6e73")))
        fig2.update_layout(
            title=dict(text="Recovery Rate of Rejected Transactions — Benchmark", font=dict(size=12, color="#1d1d1f"), x=0),
            barmode="group", paper_bgcolor="#ffffff", plot_bgcolor="#ffffff",
            font=dict(color="#6e6e73", family="Inter"),
            margin=dict(l=10, r=10, t=40, b=30),
            legend=dict(orientation="h", y=-0.15, font=dict(size=10,color="#6e6e73"), bgcolor="rgba(0,0,0,0)"),
            yaxis=dict(gridcolor="rgba(0,0,0,0.06)", ticksuffix="%"),
            xaxis=dict(gridcolor="rgba(0,0,0,0.06)"),
            height=320,
        )
        st.plotly_chart(fig2, use_container_width=True, config={"displayModeBar": False})

    # Second row — decline reasons + TPV benchmark
    col_l2, col_r2 = st.columns(2)

    with col_l2:
        reasons = ["Insufficient funds","Do not honor","Invalid card","Expired card","Fraud suspected","Limit exceeded","Tech error"]
        you_pct  = [31.2, 22.4, 14.1, 9.8, 8.3, 7.6, 6.6]
        peer_pct = [28.5, 20.1, 12.8, 11.2, 10.4, 9.5, 7.5]
        fig3 = go.Figure()
        fig3.add_trace(go.Bar(name="You", x=you_pct,  y=reasons, orientation="h", marker_color="#6366f1"))
        fig3.add_trace(go.Bar(name="Peer Avg", x=peer_pct, y=reasons, orientation="h", marker_color="#14b8a6"))
        fig3.update_layout(
            title=dict(text="Decline Reason Breakdown vs Peers", font=dict(size=12, color="#1d1d1f"), x=0),
            barmode="group", paper_bgcolor="#ffffff", plot_bgcolor="#ffffff",
            font=dict(color="#6e6e73", family="Inter"),
            margin=dict(l=10, r=40, t=40, b=10),
            legend=dict(orientation="h", y=1.1, font=dict(size=10,color="#6e6e73"), bgcolor="rgba(0,0,0,0)"),
            xaxis=dict(gridcolor="rgba(0,0,0,0.06)", ticksuffix="%"),
            yaxis=dict(autorange="reversed", tickfont=dict(size=10)),
            height=320,
        )
        st.plotly_chart(fig3, use_container_width=True, config={"displayModeBar": False})

    with col_r2:
        st.markdown("""
<div style="background:#fff;box-shadow:0 1px 3px rgba(0,0,0,0.07),0 0 0 1px rgba(0,0,0,0.04);border:none;border-radius:10px;padding:16px;height:100%;">
  <div style="font-size:12px;font-weight:700;color:#1d1d1f;margin-bottom:14px;">How You Rank</div>
  <div style="display:flex;flex-direction:column;gap:10px;">
    <div>
      <div style="display:flex;justify-content:space-between;margin-bottom:4px;"><span style="font-size:11px;color:#6e6e73;">Approval Rate</span><span style="font-size:11px;font-family:monospace;color:#818cf8;">90.2% · <span style="color:#22c55e;">Top 30%</span></span></div>
      <div style="background:#e8e8ed;border-radius:4px;height:6px;"><div style="width:70%;height:6px;background:linear-gradient(90deg,#4F46E5,#818cf8);border-radius:4px;"></div></div>
    </div>
    <div>
      <div style="display:flex;justify-content:space-between;margin-bottom:4px;"><span style="font-size:11px;color:#6e6e73;">Recovery Rate</span><span style="font-size:11px;font-family:monospace;color:#818cf8;">4.1% · <span style="color:#22c55e;">Top 25%</span></span></div>
      <div style="background:#e8e8ed;border-radius:4px;height:6px;"><div style="width:75%;height:6px;background:linear-gradient(90deg,#4F46E5,#818cf8);border-radius:4px;"></div></div>
    </div>
    <div>
      <div style="display:flex;justify-content:space-between;margin-bottom:4px;"><span style="font-size:11px;color:#6e6e73;">Fraud Rate</span><span style="font-size:11px;font-family:monospace;color:#818cf8;">0.8% · <span style="color:#22c55e;">Top 20%</span></span></div>
      <div style="background:#e8e8ed;border-radius:4px;height:6px;"><div style="width:80%;height:6px;background:linear-gradient(90deg,#4F46E5,#818cf8);border-radius:4px;"></div></div>
    </div>
    <div>
      <div style="display:flex;justify-content:space-between;margin-bottom:4px;"><span style="font-size:11px;color:#6e6e73;">Avg Response Time</span><span style="font-size:11px;font-family:monospace;color:#818cf8;">1.2s · <span style="color:#d97706;">Top 50%</span></span></div>
      <div style="background:#e8e8ed;border-radius:4px;height:6px;"><div style="width:50%;height:6px;background:linear-gradient(90deg,#4F46E5,#818cf8);border-radius:4px;"></div></div>
    </div>
    <div>
      <div style="display:flex;justify-content:space-between;margin-bottom:4px;"><span style="font-size:11px;color:#6e6e73;">Chargeback Rate</span><span style="font-size:11px;font-family:monospace;color:#818cf8;">0.3% · <span style="color:#22c55e;">Top 15%</span></span></div>
      <div style="background:#e8e8ed;border-radius:4px;height:6px;"><div style="width:85%;height:6px;background:linear-gradient(90deg,#4F46E5,#818cf8);border-radius:4px;"></div></div>
    </div>
  </div>
  <div style="margin-top:16px;padding-top:12px;border-top:1px solid rgba(0,0,0,0.07);font-size:10px;color:#86868b;">Peers = anonymised merchants in same category & region · Updated monthly</div>
</div>
""", unsafe_allow_html=True)

# ── Payment Recommendations ──────────────────────────────────────────────────
# Partner enrichment data for expanded detail view
_PARTNER_ENRICHMENT = {
    "ADYEN": {"commercial_contact":"Tom Kuehn","technical_contact":"Jorge Restrepo","verticals":"Retail, Travel, Mobility, Streaming, Marketplaces","revshare":"Yes","avg_pricing":"€0.11 + scheme fees","discount_rate":"—","discount_condition":"—","onboarding":"API keys + KYC docs, 2-4 weeks","health":"Strong","merchants_live":5,"merchants":"Uber, Netflix, Spotify, Rappi, MercadoLibre"},
    "CIELO": {"commercial_contact":"Talita Diaz Gama","technical_contact":"Valentina Perez","verticals":"Retail, eCommerce, Food Delivery, Subscriptions","revshare":"Yes — active","avg_pricing":"From 4.99%","discount_rate":"—","discount_condition":"—","onboarding":"Brazilian CNPJ required, 3-4 weeks","health":"Strong","merchants_live":3,"merchants":"Rappi, iFood, Falabella"},
    "DLOCAL": {"commercial_contact":"Daniela Reyes","technical_contact":"Jorge Restrepo","verticals":"Streaming, Gaming, SaaS, Marketplaces, Ride-hailing","revshare":"Yes — active","avg_pricing":"—","discount_rate":"—","discount_condition":"—","onboarding":"Single API, no local entity needed, 1-2 weeks","health":"Strong","merchants_live":4,"merchants":"Spotify, Netflix, Uber, MercadoLibre"},
    "STRIPE": {"commercial_contact":"Daniela Reyes","technical_contact":"Valentina Perez","verticals":"SaaS, Marketplaces, eCommerce, On-demand","revshare":"Yes — active","avg_pricing":"2.9% + $0.30","discount_rate":"—","discount_condition":"—","onboarding":"Self-serve API, 1 week","health":"Strong","merchants_live":2,"merchants":"Netflix, Despegar"},
    "CHECKOUT": {"commercial_contact":"Francisco Quintana","technical_contact":"Jorge Restrepo","verticals":"Travel, Gaming, Crypto, Forex, Marketplaces","revshare":"Yes","avg_pricing":"—","discount_rate":"—","discount_condition":"—","onboarding":"API integration, 2-3 weeks","health":"Strong","merchants_live":0,"merchants":"No merchants live yet — in onboarding"},
    "CHECKOUT MENA": {"commercial_contact":"Francisco Quintana","technical_contact":"Jorge Restrepo","verticals":"Travel, Gaming, Crypto, Forex, Marketplaces","revshare":"Yes","avg_pricing":"—","discount_rate":"—","discount_condition":"—","onboarding":"API integration, 2-3 weeks","health":"Strong","merchants_live":0,"merchants":"No merchants live yet — in onboarding"},
    "GETNET BR": {"commercial_contact":"Johanderson Guevara","technical_contact":"Valentina Perez","verticals":"Retail, Food Delivery, Subscriptions","revshare":"Yes","avg_pricing":"From 2%","discount_rate":"—","discount_condition":"—","onboarding":"Brazilian CNPJ, Santander relationship helps, 3 weeks","health":"Good","merchants_live":2,"merchants":"Rappi, Falabella"},
    "MERCADOPAGO BR": {"commercial_contact":"Johanderson Guevara","technical_contact":"Jorge Restrepo","verticals":"eCommerce, Marketplaces, Retail, Delivery","revshare":"Yes — active","avg_pricing":"4.99%","discount_rate":"—","discount_condition":"—","onboarding":"API + MercadoLibre ecosystem, 2 weeks","health":"Strong","merchants_live":3,"merchants":"MercadoLibre, Rappi, PedidosYa"},
    "EBANX": {"commercial_contact":"Talita Diaz Gama","technical_contact":"Valentina Perez","verticals":"SaaS, Gaming, Streaming, Digital Goods","revshare":"Yes — active","avg_pricing":"—","discount_rate":"—","discount_condition":"—","onboarding":"No local entity needed, 1-2 weeks","health":"Good","merchants_live":2,"merchants":"Spotify, Uber"},
    "BAMBOO PAYMENTS": {"commercial_contact":"Talita Diaz Gama","technical_contact":"Jorge Restrepo","verticals":"eCommerce, Retail, Delivery, Fintech","revshare":"Yes — active","avg_pricing":"—","discount_rate":"—","discount_condition":"—","onboarding":"Single API for LATAM, 2 weeks","health":"Good","merchants_live":2,"merchants":"Rappi, PedidosYa"},
    "NUVEI": {"commercial_contact":"Alessandra Rospigliosi","technical_contact":"Valentina Perez","verticals":"Gaming, Crypto, Forex, Travel, iGaming","revshare":"Yes","avg_pricing":"—","discount_rate":"—","discount_condition":"—","onboarding":"API + compliance review, 3-4 weeks","health":"Good","merchants_live":0,"merchants":"No merchants live yet — in pipeline"},
    "TAP PAYMENTS": {"commercial_contact":"Francisco Quintana","technical_contact":"Jorge Restrepo","verticals":"eCommerce, Retail, Food Delivery, SaaS","revshare":"Yes","avg_pricing":"2.75% + $0.30","discount_rate":"—","discount_condition":"—","onboarding":"GCC entity required, 2-3 weeks","health":"Good","merchants_live":0,"merchants":"No merchants live yet — in pipeline"},
    "PAYMOB": {"commercial_contact":"Francisco Quintana","technical_contact":"Valentina Perez","verticals":"eCommerce, Retail, SME, Food Delivery","revshare":"Pending","avg_pricing":"—","discount_rate":"—","discount_condition":"—","onboarding":"Local partnership, 2-3 weeks","health":"New","merchants_live":0,"merchants":"No merchants live yet — prospecting"},
    "PAGBANK BRASIL": {"commercial_contact":"Johanderson Guevara","technical_contact":"Jorge Restrepo","verticals":"eCommerce, Retail, SME, Subscriptions","revshare":"Yes — active","avg_pricing":"From 4.99%","discount_rate":"—","discount_condition":"—","onboarding":"Brazilian CNPJ required, 2 weeks","health":"Good","merchants_live":2,"merchants":"iFood, Claro"},
    "STONE": {"commercial_contact":"Johanderson Guevara","technical_contact":"Valentina Perez","verticals":"Retail, eCommerce, SME","revshare":"Yes","avg_pricing":"—","discount_rate":"—","discount_condition":"—","onboarding":"Brazilian CNPJ required, 3 weeks","health":"Good","merchants_live":1,"merchants":"Falabella"},
    "PAGAR.ME": {"commercial_contact":"Johanderson Guevara","technical_contact":"Jorge Restrepo","verticals":"eCommerce, SaaS, Subscriptions, Marketplaces","revshare":"Yes — active","avg_pricing":"From 3.49% + R$0.39","discount_rate":"—","discount_condition":"—","onboarding":"Brazilian CNPJ, 2-3 weeks","health":"Strong","merchants_live":2,"merchants":"Spotify, Netflix"},
    "KUSHKI": {"commercial_contact":"Talita Diaz Gama","technical_contact":"Valentina Perez","verticals":"eCommerce, Ride-hailing, Delivery, SaaS","revshare":"Yes","avg_pricing":"—","discount_rate":"—","discount_condition":"—","onboarding":"Single API, 2 weeks","health":"Good","merchants_live":2,"merchants":"Uber, Rappi"},
    "CYBERSOURCE": {"commercial_contact":"Francisco Quintana","technical_contact":"Jorge Restrepo","verticals":"Enterprise, Travel, Retail, eCommerce","revshare":"Yes","avg_pricing":"—","discount_rate":"—","discount_condition":"—","onboarding":"Visa subsidiary, 3-4 weeks","health":"Strong","merchants_live":1,"merchants":"Despegar"},
    "RAZORPAY": {"commercial_contact":"Francisco Quintana","technical_contact":"Valentina Perez","verticals":"eCommerce, SaaS, Fintech, Subscriptions","revshare":"Yes","avg_pricing":"2% per txn","discount_rate":"—","discount_condition":"—","onboarding":"Indian entity required, 2 weeks","health":"Good","merchants_live":0,"merchants":"No merchants live yet — in pipeline"},
    "XENDIT": {"commercial_contact":"Francisco Quintana","technical_contact":"Jorge Restrepo","verticals":"eCommerce, Marketplaces, Fintech, Gaming","revshare":"Yes","avg_pricing":"~2.9% + fixed fee","discount_rate":"—","discount_condition":"—","onboarding":"SEA entity required, 2 weeks","health":"Good","merchants_live":0,"merchants":"No merchants live yet — in pipeline"},
}

def show_payment_recs():
    # ── Hero ──────────────────────────────────────────────────────────────────
    st.markdown("""
<div style="background:linear-gradient(160deg,#0a0e1a 0%,#111827 40%,#1a1f3a 100%);border-radius:20px;padding:48px 44px 40px;margin-bottom:0;position:relative;overflow:hidden;">
  <div style="position:absolute;top:-80px;right:-80px;width:400px;height:400px;border-radius:50%;background:radial-gradient(circle,rgba(99,102,241,0.12),transparent 65%);"></div>
  <div style="position:absolute;bottom:-60px;left:30%;width:300px;height:300px;border-radius:50%;background:radial-gradient(circle,rgba(20,184,166,0.06),transparent 65%);"></div>
  <div style="position:absolute;top:20px;right:40px;font-size:120px;font-weight:900;color:rgba(255,255,255,0.02);letter-spacing:-6px;line-height:1;">PIE</div>
  <div style="position:relative;z-index:1;">
    <div style="display:inline-flex;align-items:center;gap:8px;padding:4px 14px;border-radius:20px;background:rgba(99,102,241,0.12);border:1px solid rgba(99,102,241,0.2);margin-bottom:16px;">
      <span style="font-size:14px;">✦</span>
      <span style="font-size:10px;font-weight:700;letter-spacing:1.5px;text-transform:uppercase;color:#a5b4fc;">Payment Intelligence Engine</span>
    </div>
    <div style="font-size:30px;font-weight:800;color:#fff;letter-spacing:-1px;line-height:1.15;margin-bottom:10px;">Find the right payment<br>partner for any market</div>
    <div style="font-size:13px;color:#94a3b8;line-height:1.7;max-width:480px;">Describe where your merchant is expanding. We match them to the best-fit payment partners — with live integration status, pricing, onboarding timeline, and merchant traction across 189 countries.</div>
  </div>
</div>""", unsafe_allow_html=True)

    # ── Search bar ────────────────────────────────────────────────────────────
    st.markdown("""
<div style="background:#fff;border:1px solid #e2e8f0;border-radius:0 0 20px 20px;padding:20px 28px 10px;margin-bottom:4px;box-shadow:0 8px 32px rgba(0,0,0,0.04);">
  <div style="font-size:10px;font-weight:700;color:#6366f1;text-transform:uppercase;letter-spacing:1.2px;margin-bottom:6px;">Describe your merchant's needs</div>
</div>""", unsafe_allow_html=True)

    # Free-text search
    free_search = st.text_input("", placeholder="e.g. Uber wants to expand into Brazil with local processing...", key="payrec_free_search", label_visibility="collapsed")

    # Parse free text into filters
    _COUNTRY_KEYWORDS = {name.lower(): iso for iso, name in _ISO_TO_COUNTRY.items()}
    # Common aliases and abbreviations
    _COUNTRY_ALIASES = {
        "u.s.": "US", "u.s": "US", "usa": "US", "united states": "US", "the us": "US", "the u.s.": "US",
        "estados unidos": "US", "america": "US",
        "u.k.": "GB", "uk": "GB", "united kingdom": "GB", "england": "GB", "britain": "GB",
        "uae": "AE", "u.a.e.": "AE", "u.a.e": "AE", "emirates": "AE", "dubai": "AE", "abu dhabi": "AE",
        "ksa": "SA", "saudi": "SA", "saudi arabia": "SA", "arabia": "SA",
        "brasil": "BR", "brazil": "BR",
        "mexico": "MX", "méxico": "MX",
        "colombia": "CO",
        "argentina": "AR",
        "chile": "CL",
        "peru": "PE", "perú": "PE",
        "turkey": "TR", "türkiye": "TR", "turkiye": "TR",
        "singapore": "SG",
        "india": "IN",
        "indonesia": "ID",
        "japan": "JP",
        "south korea": "KR", "korea": "KR",
        "germany": "DE",
        "france": "FR",
        "spain": "ES",
        "italy": "IT",
        "netherlands": "NL", "holland": "NL",
        "nigeria": "NG",
        "south africa": "ZA",
        "egypt": "EG",
        "qatar": "QA",
        "kuwait": "KW",
        "bahrain": "BH",
        "oman": "OM",
        "jordan": "JO",
        "iraq": "IQ",
        "morocco": "MA",
        "canada": "CA",
        "australia": "AU",
        "philippines": "PH",
        "thailand": "TH",
        "vietnam": "VN",
        "malaysia": "MY",
        "pakistan": "PK",
        "bangladesh": "BD",
        "kenya": "KE",
        "ghana": "GH",
        "uruguay": "UY",
        "ecuador": "EC",
        "panama": "PA",
        "costa rica": "CR",
        "guatemala": "GT",
        "dominican republic": "DO",
        "puerto rico": "PR",
        "el salvador": "SV",
        "honduras": "HN",
        "bolivia": "BO",
        "paraguay": "PY",
    }
    _auto_country = ""
    _auto_vertical = "All Verticals"
    _auto_proc = "Any"
    if free_search:
        fs_lower = free_search.lower()
        # Try aliases first (longest match)
        for alias, iso in sorted(_COUNTRY_ALIASES.items(), key=lambda x: -len(x[0])):
            if alias in fs_lower:
                _auto_country = _ISO_TO_COUNTRY.get(iso, "")
                break
        # Fallback to full country names from pycountry
        if not _auto_country:
            for cname, ciso in sorted(_COUNTRY_KEYWORDS.items(), key=lambda x: -len(x[0])):
                if cname in fs_lower:
                    _auto_country = _ISO_TO_COUNTRY.get(ciso, "")
                    break
        # Verticals — expanded keywords
        vertical_kw = {
            "gambling":"Gambling", "igaming":"Gambling", "betting":"Gambling", "casino":"Gambling",
            "gaming":"Gaming", "game":"Gaming", "esports":"Gaming",
            "crypto":"Crypto", "cryptocurrency":"Crypto", "bitcoin":"Crypto", "web3":"Crypto",
            "forex":"Forex", "trading":"Forex", "fx":"Forex",
            "airline":"Airlines", "airlines":"Airlines", "aviation":"Airlines", "flight":"Airlines", "travel":"Airlines",
            "high risk":"High Risk", "high-risk":"High Risk",
            "adult":"Adult",
            "mlm":"MLM", "multi-level":"MLM",
        }
        for kw, vl in sorted(vertical_kw.items(), key=lambda x: -len(x[0])):
            if kw in fs_lower:
                _auto_vertical = vl
                break
        if "local" in fs_lower and "processing" in fs_lower:
            _auto_proc = "LOCAL"
        elif "local" in fs_lower:
            _auto_proc = "LOCAL"
        elif "cross-border" in fs_lower or "cross border" in fs_lower or "crossborder" in fs_lower:
            _auto_proc = "CROSS_BORDER"

    st.markdown('<div style="font-size:10px;color:#94a3b8;margin:-8px 0 10px;">or refine with filters</div>', unsafe_allow_html=True)

    col1, col2, col3 = st.columns([2, 1.5, 1.5])
    with col1:
        country_options = [""] + _SOT_COUNTRIES
        default_idx = country_options.index(_auto_country) if _auto_country in country_options else 0
        selected_country = st.selectbox("Target Country", country_options, key="payrec_country",
                                        index=default_idx if not st.session_state.get("payrec_country") else 0,
                                        format_func=lambda x: "Select target market..." if x == "" else x)
        # If free search found a country but dropdown wasn't set, use it
        if not selected_country and _auto_country:
            selected_country = _auto_country
    with col2:
        vertical_options = ["All Verticals"] + list(_VERTICAL_COLS.keys())
        selected_vertical = st.selectbox("Merchant Vertical", vertical_options, key="payrec_vertical")
        if selected_vertical == "All Verticals" and _auto_vertical != "All Verticals":
            selected_vertical = _auto_vertical
    with col3:
        proc_options = ["Any", "LOCAL", "CROSS_BORDER", "BOTH"]
        selected_proc = st.selectbox("Processing Type", proc_options, key="payrec_proc",
                                     format_func=lambda x: {"Any":"Any processing","LOCAL":"Local processing","CROSS_BORDER":"Cross-border","BOTH":"Both"}[x])
        if selected_proc == "Any" and _auto_proc != "Any":
            selected_proc = _auto_proc

    col_live, col_region = st.columns([1, 2])
    with col_live:
        live_only = st.checkbox("Show live partners only", value=True, key="payrec_live")
    with col_region:
        region_filter = st.selectbox("Filter by region", ["All Regions","LATAM","Brasil","MENAT","APAC","Europe","North America","Africa"], key="payrec_region_filter", label_visibility="collapsed",
                                     format_func=lambda x: "All Regions" if x == "All Regions" else x)

    # ── Landing state ─────────────────────────────────────────────────────────
    if not selected_country:
        total_p = len(PARTNERS_DATA)
        live_p = sum(1 for p in PARTNERS_DATA if p["status"] == "Live")
        st.markdown(f"""
<div style="display:grid;grid-template-columns:repeat(4,1fr);gap:10px;margin:8px 0 16px;">
  <div style="background:#fff;border:1px solid #e2e8f0;border-radius:12px;padding:14px 16px;text-align:center;">
    <div style="font-size:24px;font-weight:800;color:#0f172a;">{total_p}</div>
    <div style="font-size:10px;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-top:2px;">Partners</div>
  </div>
  <div style="background:#fff;border:1px solid #e2e8f0;border-radius:12px;padding:14px 16px;text-align:center;">
    <div style="font-size:24px;font-weight:800;color:#22c55e;">{live_p}</div>
    <div style="font-size:10px;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-top:2px;">Live</div>
  </div>
  <div style="background:#fff;border:1px solid #e2e8f0;border-radius:12px;padding:14px 16px;text-align:center;">
    <div style="font-size:24px;font-weight:800;color:#6366f1;">189</div>
    <div style="font-size:10px;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-top:2px;">Countries</div>
  </div>
  <div style="background:#fff;border:1px solid #e2e8f0;border-radius:12px;padding:14px 16px;text-align:center;">
    <div style="font-size:24px;font-weight:800;color:#f59e0b;">7</div>
    <div style="font-size:10px;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-top:2px;">Regions</div>
  </div>
</div>""", unsafe_allow_html=True)

        # Clickable sample prompts
        _PROMPTS = [
            ("Copa Airlines expanding into the U.S.", "United States", "Airlines", "Any"),
            ("Uber needs local processing in Brazil", "Brazil", "All Verticals", "LOCAL"),
            ("iGaming operator entering Turkey", "Turkey", "Gambling", "Any"),
            ("Crypto exchange launching in UAE", "United Arab Emirates", "Crypto", "Any"),
            ("SaaS recurring payments in Saudi Arabia", "Saudi Arabia", "All Verticals", "Any"),
        ]
        st.markdown('<div style="font-size:10px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:1px;margin-bottom:6px;">Quick search</div>', unsafe_allow_html=True)
        prompt_cols = st.columns(len(_PROMPTS))
        for pi, (plabel, pcountry, pvert, pproc) in enumerate(_PROMPTS):
            if prompt_cols[pi].button(plabel, key=f"prompt_{pi}", use_container_width=True):
                st.session_state["payrec_country"] = pcountry
                st.session_state["payrec_vertical"] = pvert
                st.session_state["payrec_proc"] = pproc
                st.session_state["payrec_free_search"] = plabel
                st.rerun()

        # Style prompt buttons
        components.html("""<script>
setTimeout(function(){
  var doc = window.parent.document;
  doc.querySelectorAll('button').forEach(function(btn){
    var t = btn.textContent.trim();
    var prompts = ['Copa Airlines','Uber needs','iGaming','Crypto exchange','SaaS recurring'];
    if(prompts.some(function(p){ return t.startsWith(p); })){
      btn.style.cssText='background:#eef2ff !important;border:1px solid #e0e7ff !important;border-radius:20px !important;padding:8px 14px !important;font-size:11px !important;color:#6366f1 !important;font-weight:600 !important;cursor:pointer !important;box-shadow:none !important;white-space:normal !important;line-height:1.3 !important;text-align:left !important;min-height:0 !important;';
    }
  });
},250);
</script>""", height=0)

        # Region data with insights from PPTX
        _REGION_HEATMAP = {
            "LATAM":         {"count":"105","color":"#4F46E5","countries":["Brazil","Mexico","Colombia","Argentina","Chile","Peru"],
                              "live":15,"strategic":2,"tier1":31,"revshare":"$47.6K/mo",
                              "offering_mix":{"Acquirer":48,"PSP":27,"APM":15,"Other":10},
                              "pipeline":{"Prospect":24,"Negotiation":12,"Review":8,"Signed":32,"Live":15},
                              "trends":["Fastest-growing eCommerce region globally","PIX in Brazil driving APM adoption","BNPL and wallet penetration accelerating"]},
            "Brasil":        {"count":"73","color":"#4F46E5","countries":["Brazil"],
                              "live":32,"strategic":9,"tier1":4,"revshare":"$32.9K/mo",
                              "offering_mix":{"Acquirer":26,"PSP":24,"APM":12,"Other":5},
                              "pipeline":{"Prospect":7,"Negotiation":2,"Review":2,"Signed":14,"Live":32},
                              "trends":["PIX represents 40%+ of online payments","CNPJ required for local acquiring","Installments (parcelamento) expected at checkout"]},
            "MENAT":         {"count":"81","color":"#7c3aed","countries":["Saudi Arabia","United Arab Emirates","Turkey","Qatar","Kuwait","Egypt"],
                              "live":2,"strategic":3,"tier1":21,"revshare":"—",
                              "offering_mix":{"Acquirer":31,"APM":22,"PSP":21,"Other":6},
                              "pipeline":{"Prospect":29,"Negotiation":9,"Review":19,"Signed":15,"Live":2},
                              "trends":["GCC Vision 2030 driving cashless economies","BNPL (Tamara/Tabby) becoming standard","Local card schemes (MADA, TROY, KNET) are mandatory"]},
            "APAC":          {"count":"53","color":"#0891b2","countries":["Singapore","Thailand","Philippines","Japan","India","Indonesia"],
                              "live":6,"strategic":4,"tier1":10,"revshare":"$32.9K/mo",
                              "offering_mix":{"PSP":41,"Acquirer":7,"Fraud":2,"Other":2},
                              "pipeline":{"Prospect":8,"Negotiation":8,"Review":17,"Signed":10,"Live":6},
                              "trends":["Wallet-first markets (GrabPay, GCash, PayPay)","Real-time payments rails maturing","Data localization requirements in India, Indonesia"]},
            "Europe":        {"count":"82","color":"#2563eb","countries":["United Kingdom","Germany","France","Spain","Netherlands","Italy"],
                              "live":2,"strategic":3,"tier1":21,"revshare":"—",
                              "offering_mix":{"Acquirer":31,"APM":22,"PSP":21,"Other":6},
                              "pipeline":{"Prospect":29,"Negotiation":9,"Review":19,"Signed":15,"Live":2},
                              "trends":["iDEAL (NL), Bancontact (BE), SEPA dominate locally","Strong Data Privacy (GDPR) requirements","Open Banking / PSD2 enabling A2A payments"]},
            "North America": {"count":"37","color":"#059669","countries":["United States","Canada","Mexico"],
                              "live":3,"strategic":5,"tier1":4,"revshare":"$0.4K/mo",
                              "offering_mix":{"Acquirer":35,"APM":1,"PSP":1},
                              "pipeline":{"Prospect":6,"Negotiation":6,"Review":4,"Signed":14,"Live":3},
                              "trends":["Cards dominate (Visa/MC/Amex/Discover)","ACH and real-time payments (FedNow) emerging","Highly regulated, PCI DSS compliance critical"]},
            "Africa":        {"count":"5","color":"#d97706","countries":["South Africa","Nigeria","Kenya"],
                              "live":0,"strategic":0,"tier1":2,"revshare":"—",
                              "offering_mix":{"PSP":3,"Acquirer":2},
                              "pipeline":{"Prospect":0,"Negotiation":0,"Review":1,"Signed":4,"Live":0},
                              "trends":["Mobile money dominant (M-Pesa, MTN)","Low card penetration, high cash usage","Frontier market with long-term upside"]},
        }

        st.markdown("""
<div style="display:flex;align-items:center;gap:10px;margin:4px 0 14px;">
  <div style="height:1px;flex:1;background:linear-gradient(90deg,transparent,#e2e8f0,transparent);"></div>
  <span style="font-size:10px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:1.5px;">Global Coverage</span>
  <div style="height:1px;flex:1;background:linear-gradient(90deg,transparent,#e2e8f0,transparent);"></div>
</div>""", unsafe_allow_html=True)

        region_list = list(_REGION_HEATMAP.keys())
        for ci, region in enumerate(region_list):
            rd = _REGION_HEATMAP[region]
            exp_key = f"reg_exp_{region}"
            is_open = st.session_state.get(exp_key, False)

            if st.button(region, key=f"region_{ci}", use_container_width=True):
                for rr in region_list:
                    st.session_state[f"reg_exp_{rr}"] = False
                st.session_state[exp_key] = not is_open
                st.rerun()

            live_pct = int(int(rd["live"]) / int(rd["count"]) * 100) if int(rd["count"]) else 0
            border_c = rd["color"] if is_open else "#eef2f7"
            bg = "#fafbff" if is_open else "#fff"
            radius = "14px 14px 0 0" if is_open else "14px"
            st.markdown(f'<div style="background:{bg};border:1px solid {border_c};border-radius:{radius};padding:16px 22px;margin:-14px 0 {"0" if is_open else "6"}px;">'
                        f'<div style="display:flex;justify-content:space-between;align-items:center;">'
                        f'<div style="display:flex;align-items:center;gap:16px;">'
                        f'<div style="width:40px;height:40px;border-radius:10px;background:{rd["color"]}10;display:flex;align-items:center;justify-content:center;border:1px solid {rd["color"]}20;">'
                        f'<div style="width:10px;height:10px;border-radius:50%;background:{rd["color"]};"></div></div>'
                        f'<div><div style="font-size:15px;font-weight:700;color:#0f172a;letter-spacing:-0.3px;">{region}</div>'
                        f'<div style="font-size:11px;color:#94a3b8;margin-top:1px;">{rd["count"]} partners &middot; {rd["live"]} live &middot; {rd["revshare"]}</div></div></div>'
                        f'<div style="display:flex;align-items:center;gap:14px;">'
                        f'<div style="text-align:right;"><div style="font-size:14px;font-weight:800;color:{rd["color"]};">{live_pct}%</div>'
                        f'<div style="font-size:9px;color:#94a3b8;">live</div></div>'
                        f'<div style="width:60px;height:5px;background:#f1f5f9;border-radius:3px;overflow:hidden;"><div style="width:{live_pct}%;height:100%;background:{rd["color"]};border-radius:3px;"></div></div>'
                        f'<span style="font-size:20px;color:{rd["color"] if is_open else "#d1d5db"};font-weight:200;">{"−" if is_open else "+"}</span>'
                        f'</div></div></div>', unsafe_allow_html=True)

            if is_open:
                st.markdown(f'<div style="background:{bg};border:1px solid {rd["color"]};border-top:none;border-radius:0 0 14px 14px;padding:4px 22px 22px;margin:0 0 6px;">', unsafe_allow_html=True)

                st.markdown(f'<div style="display:grid;grid-template-columns:repeat(5,1fr);gap:8px;padding:12px 0 14px;border-bottom:1px solid #eef2ff;">'
                            f'<div style="text-align:center;padding:10px;background:#fff;border-radius:10px;border:1px solid #eef2f7;"><div style="font-size:20px;font-weight:800;color:#0f172a;">{rd["count"]}</div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-top:3px;">Total</div></div>'
                            f'<div style="text-align:center;padding:10px;background:#fff;border-radius:10px;border:1px solid #eef2f7;"><div style="font-size:20px;font-weight:800;color:#22c55e;">{rd["live"]}</div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-top:3px;">Live</div></div>'
                            f'<div style="text-align:center;padding:10px;background:#fff;border-radius:10px;border:1px solid #eef2f7;"><div style="font-size:20px;font-weight:800;color:#818cf8;">{rd["strategic"]}</div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-top:3px;">Strategic</div></div>'
                            f'<div style="text-align:center;padding:10px;background:#fff;border-radius:10px;border:1px solid #eef2f7;"><div style="font-size:20px;font-weight:800;color:#f59e0b;">{rd["tier1"]}</div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-top:3px;">Tier 1</div></div>'
                            f'<div style="text-align:center;padding:10px;background:#fff;border-radius:10px;border:1px solid #eef2f7;"><div style="font-size:20px;font-weight:800;color:{rd["color"]};">{rd["revshare"]}</div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;letter-spacing:0.5px;margin-top:3px;">Rev Share</div></div>'
                            f'</div>', unsafe_allow_html=True)

                # Offering mix + Pipeline
                col_o, col_p = st.columns(2)
                with col_o:
                    mix_bars = ""
                    total_mix = sum(rd["offering_mix"].values())
                    for cat, val in sorted(rd["offering_mix"].items(), key=lambda x: -x[1]):
                        pct = int(val / total_mix * 100)
                        mix_bars += (f'<div style="display:flex;align-items:center;gap:8px;margin-bottom:8px;">'
                                     f'<span style="font-size:11px;color:#64748b;min-width:65px;">{cat}</span>'
                                     f'<div style="flex:1;height:8px;background:#eef2ff;border-radius:4px;overflow:hidden;"><div style="width:{pct}%;height:100%;background:linear-gradient(90deg,{rd["color"]},{rd["color"]}88);border-radius:4px;"></div></div>'
                                     f'<span style="font-size:11px;font-weight:700;color:#0f172a;min-width:24px;text-align:right;">{val}</span></div>')
                    st.markdown(f'<div style="background:#fff;border:1px solid #eef2f7;border-radius:12px;padding:16px 18px;margin-top:14px;">'
                                f'<div style="font-size:10px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:0.8px;margin-bottom:12px;">Offering Mix</div>'
                                f'{mix_bars}</div>', unsafe_allow_html=True)
                with col_p:
                    pipe_stages = [("Prospect","#94a3b8"),("Negotiation","#f59e0b"),("Review","#8b5cf6"),("Signed","#3b82f6"),("Live","#22c55e")]
                    pipe_bars = ""
                    for stage, scolor in pipe_stages:
                        val = rd["pipeline"].get(stage, 0)
                        max_v = max(rd["pipeline"].values()) if rd["pipeline"] else 1
                        pct = int(val / max_v * 100) if max_v else 0
                        pipe_bars += (f'<div style="display:flex;align-items:center;gap:8px;margin-bottom:8px;">'
                                      f'<span style="font-size:11px;color:#64748b;min-width:65px;">{stage}</span>'
                                      f'<div style="flex:1;height:8px;background:#f1f5f9;border-radius:4px;overflow:hidden;"><div style="width:{pct}%;height:100%;background:{scolor};border-radius:4px;"></div></div>'
                                      f'<span style="font-size:11px;font-weight:700;color:#0f172a;min-width:24px;text-align:right;">{val}</span></div>')
                    st.markdown(f'<div style="background:#fff;border:1px solid #eef2f7;border-radius:12px;padding:16px 18px;margin-top:14px;">'
                                f'<div style="font-size:10px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:0.8px;margin-bottom:12px;">Pipeline Stage</div>'
                                f'{pipe_bars}</div>', unsafe_allow_html=True)

                trends_html = "".join(f'<div style="display:flex;gap:10px;padding:8px 0;border-bottom:1px solid #f8fafc;"><span style="width:6px;height:6px;border-radius:50%;background:{rd["color"]};flex-shrink:0;margin-top:6px;"></span><span style="font-size:12px;color:#334155;line-height:1.6;">{t}</span></div>' for t in rd["trends"])
                country_pills = " ".join(f'<span style="font-size:11px;font-weight:600;padding:6px 14px;border-radius:8px;background:#fff;color:#0f172a;border:1px solid #eef2f7;">{c}</span>' for c in rd["countries"])
                st.markdown(f'<div style="display:grid;grid-template-columns:1fr 1fr;gap:10px;margin-top:12px;">'
                            f'<div style="background:#fff;border:1px solid #eef2f7;border-radius:12px;padding:16px 18px;">'
                            f'<div style="font-size:10px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:0.8px;margin-bottom:10px;">Market Trends</div>'
                            f'{trends_html}</div>'
                            f'<div style="background:#fff;border:1px solid #eef2f7;border-radius:12px;padding:16px 18px;">'
                            f'<div style="font-size:10px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:0.8px;margin-bottom:10px;">Key Markets</div>'
                            f'<div style="display:flex;flex-wrap:wrap;gap:6px;margin-bottom:8px;">{country_pills}</div></div>'
                            f'</div>', unsafe_allow_html=True)

                st.markdown('</div>', unsafe_allow_html=True)

        # Style region toggle buttons to be invisible (card above handles visuals)
        components.html("""<script>
setTimeout(function(){
  var doc = window.parent.document;
  var main = doc.querySelector('[data-testid="stMainBlockContainer"]');
  if(!main) return;
  main.querySelectorAll('button').forEach(function(btn){
    var t = btn.textContent.trim();
    var regions = ['LATAM','Brasil','MENAT','APAC','Europe','North America','Africa'];
    if(regions.indexOf(t) !== -1){
      btn.style.cssText = 'width:100% !important;background:none !important;border:none !important;box-shadow:none !important;padding:0 !important;margin:0 !important;height:1px !important;min-height:0 !important;opacity:0 !important;cursor:pointer !important;position:absolute !important;z-index:10 !important;';
    }
  });
},200);
</script>""", height=0)

        # ── Global Insights Report ────────────────────────────────────────────
        st.markdown("""
<div style="display:flex;align-items:center;gap:10px;margin:28px 0 16px;">
  <div style="height:1px;flex:1;background:linear-gradient(90deg,transparent,#e2e8f0,transparent);"></div>
  <span style="font-size:10px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:1.5px;">Global Payments Landscape</span>
  <div style="height:1px;flex:1;background:linear-gradient(90deg,transparent,#e2e8f0,transparent);"></div>
</div>""", unsafe_allow_html=True)

        # Industry snapshot
        st.markdown("""
<div style="background:linear-gradient(160deg,#0a0e1a,#111827,#1a1f3a);border-radius:14px;padding:28px 30px;margin-bottom:14px;">
  <div style="font-size:16px;font-weight:800;color:#fff;margin-bottom:16px;">Industry Snapshot — 2026</div>
  <div style="display:grid;grid-template-columns:repeat(4,1fr);gap:12px;">
    <div style="text-align:center;padding:14px;background:rgba(255,255,255,0.04);border-radius:10px;border:1px solid rgba(255,255,255,0.06);">
      <div style="font-size:22px;font-weight:800;color:#a5b4fc;">$8.5T</div>
      <div style="font-size:9px;color:#64748b;text-transform:uppercase;letter-spacing:0.5px;margin-top:4px;">Global eComm Volume</div>
    </div>
    <div style="text-align:center;padding:14px;background:rgba(255,255,255,0.04);border-radius:10px;border:1px solid rgba(255,255,255,0.06);">
      <div style="font-size:22px;font-weight:800;color:#2dd4bf;">62%</div>
      <div style="font-size:9px;color:#64748b;text-transform:uppercase;letter-spacing:0.5px;margin-top:4px;">Digital Wallet Share</div>
    </div>
    <div style="text-align:center;padding:14px;background:rgba(255,255,255,0.04);border-radius:10px;border:1px solid rgba(255,255,255,0.06);">
      <div style="font-size:22px;font-weight:800;color:#fbbf24;">$1.2T</div>
      <div style="font-size:9px;color:#64748b;text-transform:uppercase;letter-spacing:0.5px;margin-top:4px;">BNPL Market</div>
    </div>
    <div style="text-align:center;padding:14px;background:rgba(255,255,255,0.04);border-radius:10px;border:1px solid rgba(255,255,255,0.06);">
      <div style="font-size:22px;font-weight:800;color:#f87171;">14%</div>
      <div style="font-size:9px;color:#64748b;text-transform:uppercase;letter-spacing:0.5px;margin-top:4px;">YoY Growth</div>
    </div>
  </div>
</div>""", unsafe_allow_html=True)

        # Two columns: eComm trends + POS trends
        col_ec, col_pos = st.columns(2)
        with col_ec:
            st.markdown("""
<div style="background:#fff;border:1px solid #eef2f7;border-radius:12px;padding:20px;">
  <div style="font-size:12px;font-weight:700;color:#0f172a;margin-bottom:14px;">eCommerce Trends</div>
  <div style="display:flex;gap:8px;padding:8px 0;border-bottom:1px solid #f8fafc;"><span style="width:6px;height:6px;border-radius:50%;background:#6366f1;flex-shrink:0;margin-top:6px;"></span><span style="font-size:12px;color:#334155;line-height:1.6;">Digital wallets overtook cards as #1 global eComm payment method in 2025</span></div>
  <div style="display:flex;gap:8px;padding:8px 0;border-bottom:1px solid #f8fafc;"><span style="width:6px;height:6px;border-radius:50%;background:#6366f1;flex-shrink:0;margin-top:6px;"></span><span style="font-size:12px;color:#334155;line-height:1.6;">A2A / real-time payments (PIX, UPI, SEPA Instant) are the fastest-growing rail globally</span></div>
  <div style="display:flex;gap:8px;padding:8px 0;border-bottom:1px solid #f8fafc;"><span style="width:6px;height:6px;border-radius:50%;background:#6366f1;flex-shrink:0;margin-top:6px;"></span><span style="font-size:12px;color:#334155;line-height:1.6;">BNPL adoption accelerating in MENAT (+18% CAGR) and LATAM (+22% CAGR)</span></div>
  <div style="display:flex;gap:8px;padding:8px 0;"><span style="width:6px;height:6px;border-radius:50%;background:#6366f1;flex-shrink:0;margin-top:6px;"></span><span style="font-size:12px;color:#334155;line-height:1.6;">Cross-border eCommerce growing 2x faster than domestic — orchestration is key</span></div>
</div>""", unsafe_allow_html=True)
        with col_pos:
            st.markdown("""
<div style="background:#fff;border:1px solid #eef2f7;border-radius:12px;padding:20px;">
  <div style="font-size:12px;font-weight:700;color:#0f172a;margin-bottom:14px;">POS & In-Store Trends</div>
  <div style="display:flex;gap:8px;padding:8px 0;border-bottom:1px solid #f8fafc;"><span style="width:6px;height:6px;border-radius:50%;background:#22c55e;flex-shrink:0;margin-top:6px;"></span><span style="font-size:12px;color:#334155;line-height:1.6;">Contactless payments represent 75%+ of in-store transactions in Europe and APAC</span></div>
  <div style="display:flex;gap:8px;padding:8px 0;border-bottom:1px solid #f8fafc;"><span style="width:6px;height:6px;border-radius:50%;background:#22c55e;flex-shrink:0;margin-top:6px;"></span><span style="font-size:12px;color:#334155;line-height:1.6;">QR code payments dominating LATAM and SEA — PIX QR, PromptPay, GCash</span></div>
  <div style="display:flex;gap:8px;padding:8px 0;border-bottom:1px solid #f8fafc;"><span style="width:6px;height:6px;border-radius:50%;background:#22c55e;flex-shrink:0;margin-top:6px;"></span><span style="font-size:12px;color:#334155;line-height:1.6;">Tap-to-Pay on phone (SoftPOS) eliminating hardware dependency for SMEs</span></div>
  <div style="display:flex;gap:8px;padding:8px 0;"><span style="width:6px;height:6px;border-radius:50%;background:#22c55e;flex-shrink:0;margin-top:6px;"></span><span style="font-size:12px;color:#334155;line-height:1.6;">Unified commerce (online + POS) is becoming table stakes for enterprise merchants</span></div>
</div>""", unsafe_allow_html=True)

        # Yuno coverage summary
        acq_ct = sum(1 for p in PARTNERS_DATA if p["cat"] == "Acquirer" and p["status"] == "Live")
        psp_ct = sum(1 for p in PARTNERS_DATA if p["cat"] == "PSP" and p["status"] == "Live")
        apm_ct = sum(1 for p in PARTNERS_DATA if p["cat"] == "APM" and p["status"] == "Live")
        fraud_ct = sum(1 for p in PARTNERS_DATA if p["cat"] == "Fraud" and p["status"] == "Live")
        sot_live_providers = len(_SOT_DF[_SOT_DF["Live/NonLive Partner/Contract signed"] == "Live"]["PROVIDER_NAME"].unique()) if len(_SOT_DF) > 0 else 0
        sot_live_countries = len(_SOT_DF[_SOT_DF["Live/NonLive Partner/Contract signed"] == "Live"]["COUNTRY_ISO"].unique()) if len(_SOT_DF) > 0 else 0
        sot_pm_count = len(_SOT_DF[_SOT_DF["Live/NonLive Partner/Contract signed"] == "Live"]["PAYMENT_METHOD_TYPE"].unique()) if len(_SOT_DF) > 0 else 0

        st.markdown(f"""
<div style="background:#fff;border:1px solid #eef2f7;border-radius:12px;padding:22px;margin-top:12px;">
  <div style="font-size:12px;font-weight:700;color:#0f172a;margin-bottom:16px;">What Yuno Has Integrated</div>
  <div style="display:grid;grid-template-columns:repeat(6,1fr);gap:10px;margin-bottom:16px;">
    <div style="text-align:center;padding:12px;background:#f8fafc;border-radius:10px;"><div style="font-size:18px;font-weight:800;color:#3b82f6;">{sot_live_providers}</div><div style="font-size:9px;color:#94a3b8;margin-top:2px;">LIVE PROVIDERS</div></div>
    <div style="text-align:center;padding:12px;background:#f8fafc;border-radius:10px;"><div style="font-size:18px;font-weight:800;color:#22c55e;">{sot_live_countries}</div><div style="font-size:9px;color:#94a3b8;margin-top:2px;">COUNTRIES</div></div>
    <div style="text-align:center;padding:12px;background:#f8fafc;border-radius:10px;"><div style="font-size:18px;font-weight:800;color:#6366f1;">{sot_pm_count}</div><div style="font-size:9px;color:#94a3b8;margin-top:2px;">PAYMENT METHODS</div></div>
    <div style="text-align:center;padding:12px;background:#f8fafc;border-radius:10px;"><div style="font-size:18px;font-weight:800;color:#3b82f6;">{acq_ct}</div><div style="font-size:9px;color:#94a3b8;margin-top:2px;">ACQUIRERS</div></div>
    <div style="text-align:center;padding:12px;background:#f8fafc;border-radius:10px;"><div style="font-size:18px;font-weight:800;color:#8b5cf6;">{psp_ct}</div><div style="font-size:9px;color:#94a3b8;margin-top:2px;">PSPs</div></div>
    <div style="text-align:center;padding:12px;background:#f8fafc;border-radius:10px;"><div style="font-size:18px;font-weight:800;color:#14b8a6;">{apm_ct}</div><div style="font-size:9px;color:#94a3b8;margin-top:2px;">APMs</div></div>
  </div>
  <div style="display:grid;grid-template-columns:1fr 1fr;gap:10px;">
    <div style="padding:12px 16px;background:#eef2ff;border-radius:10px;border:1px solid #e0e7ff;">
      <div style="font-size:11px;font-weight:700;color:#4F46E5;margin-bottom:4px;">Orchestration Layer</div>
      <div style="font-size:12px;color:#334155;line-height:1.5;">Yuno routes transactions across multiple providers per country — enabling smart retries, cost optimization, and approval rate maximization through a single API.</div>
    </div>
    <div style="padding:12px 16px;background:#ecfdf5;border-radius:10px;border:1px solid #a7f3d0;">
      <div style="font-size:11px;font-weight:700;color:#065f46;margin-bottom:4px;">Integration Speed</div>
      <div style="font-size:12px;color:#334155;line-height:1.5;">New payment providers can be integrated in under 30 days. Merchants connect once to Yuno and gain access to the entire partner ecosystem — no additional development required.</div>
    </div>
  </div>
</div>""", unsafe_allow_html=True)

        return

    # ── Find ISO + run search ─────────────────────────────────────────────────
    country_iso = None
    for iso, name in _ISO_TO_COUNTRY.items():
        if name == selected_country:
            country_iso = iso
            break

    verticals_filter = [selected_vertical] if selected_vertical != "All Verticals" else None
    proc_filter = selected_proc if selected_proc != "Any" else None
    results = find_partners(country_iso=country_iso, verticals=verticals_filter, live_only=live_only, processing_type=proc_filter)
    v_label = selected_vertical if selected_vertical != "All Verticals" else "all verticals"

    # ── Results header ────────────────────────────────────────────────────────
    st.markdown(f'<div style="background:#fff;border:1px solid #e2e8f0;border-radius:12px;padding:18px 22px;margin:8px 0 12px;">'
                f'<div style="display:flex;justify-content:space-between;align-items:center;">'
                f'<div><div style="font-size:11px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:1px;">Results for {selected_country}</div>'
                f'<div style="font-size:22px;font-weight:800;color:#0f172a;margin-top:2px;">{len(results)} partner{"s" if len(results)!=1 else ""} found</div>'
                f'<div style="font-size:12px;color:#64748b;margin-top:2px;">{v_label}</div></div>'
                f'<div style="font-size:48px;font-weight:800;color:#4F46E5;opacity:0.1;">{len(results)}</div></div></div>', unsafe_allow_html=True)

    if not results:
        st.markdown('<div style="border:1px solid #fde68a;border-radius:10px;padding:24px;background:#fffbeb;text-align:center;">'
                    '<div style="font-size:15px;font-weight:600;color:#92400e;">No partners match this combination</div>'
                    '<div style="font-size:13px;color:#b45309;margin-top:4px;">Try broadening your filters.</div></div>', unsafe_allow_html=True)
        return

    # ── Results with inline expandable profiles ───────────────────────────────
    CAT_LABELS = {"ACQUIRER":"Acquirer","GATEWAY":"PSP/Gateway","AGREGATOR":"Aggregator","AGREGATOR / GATEWAY":"Aggregator/Gateway","PAYMENT_METHOD":"APM"}

    for idx, r in enumerate(results):
        cat_label = CAT_LABELS.get(r["categories"][0], r["categories"][0]) if r["categories"] else ""
        pm_display = ", ".join(r["payment_methods"][:5])
        if len(r["payment_methods"]) > 5:
            pm_display += f" +{len(r['payment_methods'])-5}"
        proc_display = " / ".join(r["processing_types"])
        rank_color = "#4F46E5" if idx < 3 else "#94a3b8"
        s_bg = "#ecfdf5" if r["status"] == "Live" else "#f8fafc"
        s_color = "#065f46" if r["status"] == "Live" else "#64748b"

        # Card header
        st.markdown(f'<div style="background:#fff;border:1px solid #e2e8f0;border-radius:10px;padding:14px 20px;margin-top:6px;">'
                    f'<div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:8px;">'
                    f'<div style="display:flex;align-items:center;gap:10px;">'
                    f'<span style="font-size:16px;font-weight:800;color:{rank_color};">#{idx+1}</span>'
                    f'<span style="font-size:15px;font-weight:700;color:#0f172a;">{r["name"]}</span>'
                    f'<span style="font-size:10px;font-weight:600;padding:2px 8px;border-radius:4px;background:#eff6ff;color:#1e40af;">{cat_label}</span></div>'
                    f'<span style="font-size:10px;font-weight:600;padding:3px 10px;border-radius:20px;background:{s_bg};color:{s_color};">{r["status"]}</span></div>'
                    f'<div style="display:grid;grid-template-columns:1fr 2fr 1fr;gap:8px;">'
                    f'<div><div style="font-size:10px;color:#94a3b8;text-transform:uppercase;">Processing</div><div style="font-size:12px;font-weight:600;color:#0f172a;">{proc_display}</div></div>'
                    f'<div><div style="font-size:10px;color:#94a3b8;text-transform:uppercase;">Payment Methods</div><div style="font-size:12px;color:#64748b;">{pm_display}</div></div>'
                    f'<div><div style="font-size:10px;color:#94a3b8;text-transform:uppercase;">Countries</div><div style="font-size:12px;font-weight:600;color:#0f172a;">{r["countries_count"]}</div></div>'
                    f'</div></div>', unsafe_allow_html=True)

        # Expand button (Streamlit native — clean toggle)
        enrich_key = r["name"].upper().replace("(","").replace(")","").strip()
        enrich = _PARTNER_ENRICHMENT.get(enrich_key)

        exp_key = f'exp_{r["name"]}'
        is_expanded = st.session_state.get(exp_key, False)
        arrow = "▾" if is_expanded else "▸"
        btn_label = f"{arrow} {r['name']} — Details"
        if st.button(btn_label, key=f"btn_exp_{idx}", use_container_width=True):
            st.session_state[exp_key] = not is_expanded
            st.rerun()

        if is_expanded:
            feat_map = {"SUPPORTS_TOKENIZATION":"Tokenization","SUPPORTS_RECURRING_PAYMENTS":"Recurring","SUPPORTS_PAYOUTS":"Payouts","SUPPORTS_INSTALLMENTS":"Installments","SUPPORTS_PAYFAC":"PayFac","SUPPORTS_SPLIT_PAYMENTS":"Split Pay","3DS":"3D Secure"}
            on_pills = " ".join(f'<span style="font-size:10px;padding:3px 8px;border-radius:4px;background:#ecfdf5;color:#065f46;">{v}</span>' for k,v in feat_map.items() if r["supports"].get(k))
            off_pills = " ".join(f'<span style="font-size:10px;padding:3px 8px;border-radius:4px;background:#f8fafc;color:#cbd5e1;text-decoration:line-through;">{v}</span>' for k,v in feat_map.items() if not r["supports"].get(k))

            if enrich:
                h_color = "#065f46" if enrich["health"] == "Strong" else "#d97706" if enrich["health"] == "Good" else "#94a3b8"
                h_bg = "#ecfdf5" if enrich["health"] == "Strong" else "#fffbeb" if enrich["health"] == "Good" else "#f8fafc"

                # Integrated features
                st.markdown(f'<div style="background:#fafafa;border:1px solid #e2e8f0;border-radius:10px;padding:18px 20px;margin:-4px 0 4px;">'
                            f'<div style="font-size:10px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:0.8px;margin-bottom:8px;">Yuno Integrated Capabilities</div>'
                            f'<div style="display:flex;flex-wrap:wrap;gap:4px;">{on_pills} {off_pills}</div></div>', unsafe_allow_html=True)

                # Detail grid — contacts + commercial
                dc1, dc2 = st.columns(2)
                with dc1:
                    _pkey = r["name"].upper().replace("(","").replace(")","").strip()
                    _ec = _EXT_CONTACTS.get(_pkey, {})
                    _cl = _ec.get("commercial_contact","") or enrich["commercial_contact"]
                    _tl = _ec.get("technical_contact","") or enrich["technical_contact"]
                    _ce = _ec.get("commercial_email","")
                    _te = _ec.get("technical_email","")
                    _ev = _ec.get("verticals","") or enrich["verticals"]
                    _cl_line = f'{_cl}<br><span style="font-size:10px;color:#64748b;">{_ce}</span>' if _ce else _cl
                    _tl_line = f'{_tl}<br><span style="font-size:10px;color:#64748b;">{_te}</span>' if _te else _tl
                    st.markdown(f'<div style="background:#fff;border:1px solid #e2e8f0;border-radius:10px;padding:18px 20px;">'
                                f'<div style="font-size:10px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:0.8px;margin-bottom:12px;">Contacts & Verticals</div>'
                                f'<div style="margin-bottom:10px;"><div style="font-size:10px;color:#94a3b8;">Commercial Lead</div><div style="font-size:13px;font-weight:600;color:#0f172a;">{_cl_line}</div></div>'
                                f'<div style="margin-bottom:10px;"><div style="font-size:10px;color:#94a3b8;">Technical Lead</div><div style="font-size:13px;font-weight:600;color:#0f172a;">{_tl_line}</div></div>'
                                f'<div style="margin-bottom:10px;"><div style="font-size:10px;color:#94a3b8;">Focus Verticals</div><div style="font-size:12px;color:#334155;">{_ev}</div></div>'
                                f'<div><div style="font-size:10px;color:#94a3b8;">Rev Share</div><div style="font-size:13px;font-weight:600;color:#0f172a;">{enrich["revshare"]}</div></div>'
                                f'</div>', unsafe_allow_html=True)
                with dc2:
                    st.markdown(f'<div style="background:#fff;border:1px solid #e2e8f0;border-radius:10px;padding:18px 20px;">'
                                f'<div style="font-size:10px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:0.8px;margin-bottom:12px;">Pricing & Onboarding</div>'
                                f'<div style="margin-bottom:10px;"><div style="font-size:10px;color:#94a3b8;">Standard Rate</div><div style="font-size:13px;font-weight:600;color:#0f172a;">{enrich["avg_pricing"]}</div></div>'
                                f'<div style="margin-bottom:10px;"><div style="font-size:10px;color:#94a3b8;">Discounted Rate</div><div style="font-size:13px;font-weight:700;color:#065f46;">{enrich["discount_rate"]}</div></div>'
                                f'<div style="margin-bottom:10px;"><div style="font-size:10px;color:#94a3b8;">Discount Condition</div><div style="font-size:12px;color:#334155;">{enrich["discount_condition"]}</div></div>'
                                f'<div><div style="font-size:10px;color:#94a3b8;">Onboarding</div><div style="font-size:12px;color:#334155;">{enrich["onboarding"]}</div></div>'
                                f'</div>', unsafe_allow_html=True)

                # Health + merchants bar
                st.markdown(f'<div style="background:#fff;border:1px solid #e2e8f0;border-radius:10px;padding:14px 20px;margin-bottom:4px;">'
                            f'<div style="display:flex;gap:20px;align-items:center;">'
                            f'<div style="display:flex;align-items:center;gap:6px;"><span style="font-size:10px;color:#94a3b8;">Health</span><span style="font-size:11px;font-weight:700;padding:2px 10px;border-radius:4px;background:{h_bg};color:{h_color};">{enrich["health"]}</span></div>'
                            f'<div style="display:flex;align-items:center;gap:6px;"><span style="font-size:10px;color:#94a3b8;">Live Merchants</span><span style="font-size:13px;font-weight:700;color:#0f172a;">{enrich["merchants_live"]}</span></div>'
                            f'<div style="flex:1;display:flex;align-items:center;gap:6px;"><span style="font-size:10px;color:#94a3b8;">Merchants</span><span style="font-size:12px;color:#334155;">{enrich["merchants"]}</span></div>'
                            f'</div></div>', unsafe_allow_html=True)
            else:
                st.markdown(f'<div style="background:#fafafa;border:1px solid #e2e8f0;border-radius:10px;padding:18px 20px;margin:-4px 0 4px;">'
                            f'<div style="font-size:10px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:0.8px;margin-bottom:8px;">Yuno Integrated Capabilities</div>'
                            f'<div style="display:flex;flex-wrap:wrap;gap:4px;margin-bottom:10px;">{on_pills} {off_pills}</div>'
                            f'<div style="font-size:12px;color:#94a3b8;">Full commercial profile pending — contact partnerships team.</div></div>', unsafe_allow_html=True)

    # Style detail toggle buttons
    components.html("""<script>
setTimeout(function(){
  var doc = window.parent.document;
  doc.querySelectorAll('button').forEach(function(btn){
    var t = btn.textContent.trim();
    if(t.includes('— Details')){
      btn.style.cssText='background:none !important;border:1px solid #e2e8f0 !important;border-radius:8px !important;padding:6px 16px !important;margin:-2px 0 2px !important;font-size:11px !important;color:#4F46E5 !important;font-weight:600 !important;cursor:pointer !important;box-shadow:none !important;text-align:left !important;';
    }
  });
},300);
</script>""", height=0)

    # Summary insight
    top3 = [r["name"] for r in results[:3]]
    has_local = any("LOCAL" in r.get("processing_types",[]) for r in results[:3])
    st.markdown(f'<div style="background:linear-gradient(135deg,#f8fafc,#eff6ff);border:1px solid #e2e8f0;border-radius:10px;padding:18px 22px;margin-top:16px;">'
                f'<div style="font-size:11px;font-weight:700;color:#4F46E5;text-transform:uppercase;letter-spacing:1px;margin-bottom:6px;">Recommendation</div>'
                f'<div style="font-size:14px;color:#0f172a;line-height:1.6;">For <strong>{selected_country}</strong> ({v_label}), your top options are <strong>{", ".join(top3)}</strong>. '
                f'{"These partners offer local processing with established merchant volume." if has_local else "Consider cross-border routing for initial testing before local entity setup."}'
                f'</div></div>', unsafe_allow_html=True)

    # ── Market Intelligence (MENAT insights from PPTX) ────────────────────────
    _MARKET_INSIGHTS = {
        "Saudi Arabia": {
            "tier": "Tier 1", "tier_color": "#065f46", "tier_bg": "#ecfdf5",
            "population": "36M", "ecomm_size": "$13.3B", "ecomm_growth": "12.5%",
            "smartphone": "97%", "internet": "99%",
            "digital_trends": [
                "Vision 2030 targets cashless economy",
                "Surge in online fashion, electronics, and food delivery post-pandemic",
                "Saudis spend ~$7B/year buying goods from cross-border e-commerce merchants",
            ],
            "pro_tips": [
                "MADA is the national debit scheme — mandatory for local card acceptance",
                "SARIE (A2A) delivers cost savings and eliminates chargeback exposure",
                "BNPL via Tamara/Tabby is becoming a standard checkout expectation",
            ],
            "key_verticals": ["Electronics", "Fashion", "Health & Beauty", "Groceries & Deliveries"],
            "local_methods": ["MADA", "SARIE", "Tamara", "Tabby", "Apple Pay", "STC Pay"],
        },
        "United Arab Emirates": {
            "tier": "Tier 1", "tier_color": "#065f46", "tier_bg": "#ecfdf5",
            "population": "10M", "ecomm_size": "$8B", "ecomm_growth": "11%",
            "smartphone": "96%", "internet": "99%",
            "digital_trends": [
                "Credit cards dominant in eCommerce, wallets rising fast",
                "New local card rail Jaywan emerging",
                "BNPL becoming standard checkout expectation",
            ],
            "pro_tips": [
                "Local entity or free zone setup is required to operate",
                "~90% of population are expats — extremely diverse payment preferences",
                "Arabic RTL checkout support is expected, not optional",
            ],
            "key_verticals": ["Electronics", "Fashion", "Health & Beauty", "Groceries", "Baby Products"],
            "local_methods": ["Jaywan", "AANI", "Tamara", "Tabby", "Apple Pay"],
        },
        "Turkey": {
            "tier": "Tier 1", "tier_color": "#065f46", "tier_bg": "#ecfdf5",
            "population": "86M", "ecomm_size": "$28B", "ecomm_growth": "6.5% of GDP",
            "smartphone": "92%", "internet": "88%",
            "digital_trends": [
                "eCommerce scaling to ~$45B by 2030, 6.5% of GDP (highest in region)",
                "Consumers expect 3/6/12-month installment splits — primary driver of card dominance",
                "With ~39% inflation, BNPL and installments are survival tools, not convenience",
            ],
            "pro_tips": [
                "Turkey mandates data localization for payment data, BDDK/CBT licensing required",
                "TROY must be enabled alongside Visa & Mastercard — skipping it increases costs and declines",
                "FAST (A2A) delivers meaningful savings while removing dispute exposure",
            ],
            "key_verticals": ["Electronics", "Beauty", "Fashion", "Food Delivery", "Gaming"],
            "local_methods": ["TROY", "FAST", "BKM Express", "Papara"],
        },
        "Qatar": {
            "tier": "Tier 2", "tier_color": "#1e40af", "tier_bg": "#eff6ff",
            "population": "3M", "ecomm_size": "$3.8B", "ecomm_growth": "10.3%",
            "smartphone": "95%", "internet": "99%",
            "digital_trends": [
                "eCommerce growing at 9.31% CAGR, transactions up 15% YoY",
                "Digital economy pillar of Qatar National Vision 2030",
                "69% of eCommerce revenue comes via smartphones, 5G covers 99% of population",
            ],
            "pro_tips": [
                "Qatar is small — Doha covers majority of merchant customer base",
                "BNPL growing at 18%+ CAGR, directly lifts AOV",
                "Winning checkout must cover cards, Fawran and digital wallets",
            ],
            "key_verticals": ["Electronics", "Beauty", "Fashion", "Food Delivery", "Travel"],
            "local_methods": ["NAPS", "Fawran", "QMP", "Dynamic QR"],
        },
        "Kuwait": {
            "tier": "Tier 2", "tier_color": "#1e40af", "tier_bg": "#eff6ff",
            "population": "4.3M", "ecomm_size": "$1.85B", "ecomm_growth": "5.5%",
            "smartphone": "98%", "internet": "99%",
            "digital_trends": [
                "KNET national debit network used in ~80% of all online transactions",
                "WAMD instant payments processed $1.5B across 5.4M txns in launch month",
                "eCommerce projected to reach $2.42B by 2030",
            ],
            "pro_tips": [
                "KNET is non-negotiable — local debit network supported by all 11 national banks",
                "Add WAMD for instant A2A payments — brand new and growing rapidly",
                "Younger demographics actively seek BNPL installment options",
            ],
            "key_verticals": ["Electronics", "Beauty", "Fashion", "Home & Furniture", "Groceries"],
            "local_methods": ["KNET", "WAMD", "Tamara", "Tabby"],
        },
        "Bahrain": {
            "tier": "Tier 2", "tier_color": "#1e40af", "tier_bg": "#eff6ff",
            "population": "1.5M", "ecomm_size": "$1.2B", "ecomm_growth": "8%",
            "smartphone": "92%", "internet": "99%",
            "digital_trends": [
                "BenefitPay holds 46.8% of eCommerce payment share — highest wallet dominance in GCC",
                "Cash on Delivery fallen below 10% — lowest in GCC",
                "Two-thirds of all transactions happen via mobile",
            ],
            "pro_tips": [
                "Integrate BenefitPay — it's the market's payment spine (debit, QR, A2A, direct debit)",
                "BNPL reduces abandonment and increases AOV simultaneously",
                "Don't assume GCC incorporation elsewhere automatically covers Bahrain operations",
            ],
            "key_verticals": ["Electronics", "Beauty", "Fashion", "Home & Furniture", "Food Delivery"],
            "local_methods": ["BENEFIT", "Fawri+", "Tamara", "Tabby"],
        },
        "Egypt": {
            "tier": "Tier 3", "tier_color": "#92400e", "tier_bg": "#fffbeb",
            "population": "106M", "ecomm_size": "$10B", "ecomm_growth": "25%",
            "smartphone": "78%", "internet": "72%",
            "digital_trends": [
                "eCommerce scaling to ~$19B by 2030",
                "Cash on Delivery remains dominant at ~45-50%",
                "Digital wallets surging: 50-55M active e-wallets, 80% YoY volume growth",
            ],
            "pro_tips": [
                "Multi-wallet integration is mandatory, not optional",
                "InstaPay (A2A) is maturing — low cost and high adoption",
                "Currency devaluation impacts consumer purchasing power and merchant economics",
                "Incentivize digital payments at checkout to convert CoD users",
            ],
            "key_verticals": ["Electronics", "Food Delivery", "Grocery", "Pharmacy", "Home & Furniture"],
            "local_methods": ["MEEZA", "InstaPay", "Fawry", "Vodafone Cash"],
        },
        "Iraq": {
            "tier": "High Risk", "tier_color": "#dc2626", "tier_bg": "#fef2f2",
            "population": "47M", "ecomm_size": "$780M", "ecomm_growth": "11.2%",
            "smartphone": "87%", "internet": "83%",
            "digital_trends": [
                "Over $15B+ in electronic transactions processed last year",
                "CBI launched National Financial Inclusion Strategy 2025-2029",
                "Fewer than 20% hold a formal bank account, 85% of eCommerce is cash",
            ],
            "pro_tips": [
                "Millions receive salaries via QiCard — leverage this as a competitive advantage",
                "E-wallets are essential given low banked population, not just a convenience",
                "Offer intelligent retries across PSPs and wallets given infrastructure instability",
            ],
            "key_verticals": ["Electronics", "Beauty", "Fashion", "Home & Appliances", "Daily Essentials"],
            "local_methods": ["QiCard", "Zain Cash", "IPS"],
        },
    }

    market = _MARKET_INSIGHTS.get(selected_country)
    if market:
        # Market overview card
        st.markdown(f'<div style="margin-top:24px;margin-bottom:8px;">'
                    f'<div style="font-size:11px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:1.5px;">Market Intelligence</div></div>', unsafe_allow_html=True)

        # Hero stats bar
        st.markdown(f'<div style="background:linear-gradient(135deg,#0f172a,#1e293b);border-radius:12px;padding:28px 32px;margin-bottom:12px;">'
                    f'<div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:20px;">'
                    f'<div style="font-size:22px;font-weight:800;color:#fff;letter-spacing:-0.5px;">{selected_country}</div>'
                    f'<span style="font-size:10px;font-weight:700;padding:4px 14px;border-radius:20px;background:{market["tier_bg"]};color:{market["tier_color"]};">{market["tier"]}</span>'
                    f'</div>'
                    f'<div style="display:grid;grid-template-columns:repeat(5,1fr);gap:12px;">'
                    f'<div style="text-align:center;"><div style="font-size:22px;font-weight:800;color:#fff;">{market["population"]}</div><div style="font-size:10px;color:#64748b;margin-top:2px;">Population</div></div>'
                    f'<div style="text-align:center;"><div style="font-size:22px;font-weight:800;color:#818cf8;">{market["ecomm_size"]}</div><div style="font-size:10px;color:#64748b;margin-top:2px;">eCommerce Size</div></div>'
                    f'<div style="text-align:center;"><div style="font-size:22px;font-weight:800;color:#2dd4bf;">{market["ecomm_growth"]}</div><div style="font-size:10px;color:#64748b;margin-top:2px;">Growth Rate</div></div>'
                    f'<div style="text-align:center;"><div style="font-size:22px;font-weight:800;color:#fff;">{market["smartphone"]}</div><div style="font-size:10px;color:#64748b;margin-top:2px;">Smartphone</div></div>'
                    f'<div style="text-align:center;"><div style="font-size:22px;font-weight:800;color:#fff;">{market["internet"]}</div><div style="font-size:10px;color:#64748b;margin-top:2px;">Internet</div></div>'
                    f'</div></div>', unsafe_allow_html=True)

        # Two-column: Digital Trends + Pro Tips
        col_t, col_p = st.columns(2)
        with col_t:
            trends_html = "".join(f'<div style="display:flex;gap:10px;padding:10px 0;border-bottom:1px solid #f1f5f9;"><span style="color:#818cf8;font-size:14px;flex-shrink:0;">&#9672;</span><span style="font-size:13px;color:#334155;line-height:1.5;">{t}</span></div>' for t in market["digital_trends"])
            st.markdown(f'<div style="background:#fff;border:1px solid #e2e8f0;border-radius:10px;padding:20px;height:100%;">'
                        f'<div style="font-size:13px;font-weight:700;color:#0f172a;margin-bottom:12px;">Digital Trends</div>'
                        f'{trends_html}</div>', unsafe_allow_html=True)
        with col_p:
            tips_html = "".join(f'<div style="display:flex;gap:10px;padding:10px 0;border-bottom:1px solid #f1f5f9;"><span style="color:#f59e0b;font-size:12px;flex-shrink:0;font-weight:800;">TIP</span><span style="font-size:13px;color:#334155;line-height:1.5;">{t}</span></div>' for t in market["pro_tips"])
            st.markdown(f'<div style="background:#fff;border:1px solid #e2e8f0;border-radius:10px;padding:20px;height:100%;">'
                        f'<div style="font-size:13px;font-weight:700;color:#0f172a;margin-bottom:12px;">Pro Tips</div>'
                        f'{tips_html}</div>', unsafe_allow_html=True)

        # Bottom row: Verticals + Local Payment Methods
        col_v, col_m = st.columns(2)
        with col_v:
            vert_pills = " ".join(f'<span style="font-size:11px;font-weight:600;padding:6px 14px;border-radius:8px;background:#f8fafc;color:#334155;border:1px solid #e2e8f0;">{v}</span>' for v in market["key_verticals"])
            st.markdown(f'<div style="background:#fff;border:1px solid #e2e8f0;border-radius:10px;padding:20px;margin-top:10px;">'
                        f'<div style="font-size:13px;font-weight:700;color:#0f172a;margin-bottom:12px;">Key eCommerce Verticals</div>'
                        f'<div style="display:flex;flex-wrap:wrap;gap:6px;">{vert_pills}</div></div>', unsafe_allow_html=True)
        with col_m:
            # Check which methods Yuno has available via which partner
            method_items = []
            for m in market["local_methods"]:
                # Search SOT for this payment method in this country
                m_upper = m.upper().replace(" ","_")
                matches = _SOT_DF[(_SOT_DF["COUNTRY_ISO"] == country_iso) & (_SOT_DF["PAYMENT_METHOD_TYPE"].str.upper().str.contains(m_upper, na=False)) & (_SOT_DF["Live/NonLive Partner/Contract signed"] == "Live")]
                if len(matches) > 0:
                    via_partner = matches["PROVIDER_NAME"].iloc[0]
                    method_items.append(f'<div style="display:flex;justify-content:space-between;align-items:center;padding:8px 0;border-bottom:1px solid #f1f5f9;">'
                                        f'<div style="display:flex;align-items:center;gap:8px;"><span style="width:8px;height:8px;border-radius:50%;background:#22c55e;flex-shrink:0;"></span>'
                                        f'<span style="font-size:13px;font-weight:600;color:#0f172a;">{m}</span></div>'
                                        f'<span style="font-size:11px;color:#065f46;font-weight:600;">via {via_partner}</span></div>')
                else:
                    # Check non-live
                    matches_nl = _SOT_DF[(_SOT_DF["COUNTRY_ISO"] == country_iso) & (_SOT_DF["PAYMENT_METHOD_TYPE"].str.upper().str.contains(m_upper, na=False))]
                    if len(matches_nl) > 0:
                        via_partner = matches_nl["PROVIDER_NAME"].iloc[0]
                        method_items.append(f'<div style="display:flex;justify-content:space-between;align-items:center;padding:8px 0;border-bottom:1px solid #f1f5f9;">'
                                            f'<div style="display:flex;align-items:center;gap:8px;"><span style="width:8px;height:8px;border-radius:50%;background:#f59e0b;flex-shrink:0;"></span>'
                                            f'<span style="font-size:13px;font-weight:600;color:#0f172a;">{m}</span></div>'
                                            f'<span style="font-size:11px;color:#92400e;font-weight:600;">pipeline — {via_partner}</span></div>')
                    else:
                        method_items.append(f'<div style="display:flex;justify-content:space-between;align-items:center;padding:8px 0;border-bottom:1px solid #f1f5f9;">'
                                            f'<div style="display:flex;align-items:center;gap:8px;"><span style="width:8px;height:8px;border-radius:50%;background:#ef4444;flex-shrink:0;"></span>'
                                            f'<span style="font-size:13px;font-weight:600;color:#0f172a;">{m}</span></div>'
                                            f'<span style="font-size:11px;color:#dc2626;font-weight:600;">not available</span></div>')
            methods_html = "".join(method_items)
            available_ct = sum(1 for m in market["local_methods"] if any(m.upper().replace(" ","_") in str(row).upper() for row in _SOT_DF[(_SOT_DF["COUNTRY_ISO"]==country_iso)&(_SOT_DF["Live/NonLive Partner/Contract signed"]=="Live")]["PAYMENT_METHOD_TYPE"].dropna().values))
            coverage_pct = int(available_ct / len(market["local_methods"]) * 100) if market["local_methods"] else 0
            cov_color = "#065f46" if coverage_pct >= 70 else "#d97706" if coverage_pct >= 40 else "#dc2626"
            st.markdown(f'<div style="background:#fff;border:1px solid #e2e8f0;border-radius:10px;padding:20px;margin-top:10px;">'
                        f'<div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:14px;">'
                        f'<div style="font-size:13px;font-weight:700;color:#0f172a;">Must-Have Payment Methods</div>'
                        f'<span style="font-size:11px;font-weight:700;color:{cov_color};">{coverage_pct}% covered</span></div>'
                        f'{methods_html}</div>', unsafe_allow_html=True)


# ── Partners In Flight ─────────────────────────────────────────────────────────
_MC_STAGES = [
    {"id":"prospect","label":"Prospect","short":"PROS","owner":"Partnership BD","owner_color":"#94a3b8","desc":"Market assessment, initial contact, opportunity identification"},
    {"id":"negotiation","label":"Negotiation","short":"NEG","owner":"Partnership BD","owner_color":"#6366f1","desc":"Commercial terms discussion, pricing alignment, scope definition"},
    {"id":"agreement","label":"Agreement Review","short":"AGR","owner":"Partnerships + Legal","owner_color":"#8b5cf6","desc":"NDA review, contract drafting, compliance checks, legal sign-off"},
    {"id":"agreement_signed","label":"Agreement Signed","short":"SIGN","owner":"Partnerships + Legal","owner_color":"#7c3aed","desc":"Contract executed, ready for prioritization"},
    {"id":"prioritized","label":"Prioritized","short":"PRI","owner":"Product + Partnerships","owner_color":"#0891b2","desc":"Roadmap slot confirmed, resources allocated, kickoff planned"},
    {"id":"in_development","label":"In Development","short":"DEV","owner":"Integrations","owner_color":"#2563eb","desc":"API integration, sandbox testing, certification in progress"},
    {"id":"integration_complete","label":"Development Complete","short":"QA","owner":"Integrations","owner_color":"#059669","desc":"Testing passed, sandbox certified, ready for production"},
    {"id":"go_live","label":"Go Live","short":"LIVE","owner":"Merchant BD/KAM + Partnerships + Product","owner_color":"#16a34a","desc":"Production deployment, merchant activation, initial transaction flow"},
    {"id":"pilot","label":"Pilot Phase","short":"PILOT","owner":"Merchant BD/KAM + Partnerships","owner_color":"#d97706","desc":"First merchant onboarding, volume ramp, performance validation"},
    {"id":"g2m","label":"Go to Market","short":"G2M","owner":"Marketing + Partnerships","owner_color":"#e11d48","desc":"Joint GTM campaigns, co-marketing, case studies, press releases"},
    {"id":"live_processing","label":"Live & Processing","short":"LIVE+","owner":"Partnership Account Mgmt","owner_color":"#065f46","desc":"Active volume, rev share collecting, ongoing optimization"},
]

def _map_deal_stage_to_mc(row):
    """Map Excel deal stage + integration columns to mission control stage."""
    stage = str(row.get("Deal Stage","")).strip()
    int_ready = bool(row.get("Integration Ready by Yuno", False))
    int_used = bool(row.get("Integration Used by Merchants", False))
    revshare_active = bool(row.get("Revshare active", False))
    nda = bool(row.get("NDA Signed and in drive", False))
    revshare_contract = bool(row.get("Revshare Contract", False))
    mgmt = str(row.get("Type of Management","")).strip()

    if stage == "Live Partner" and int_used and revshare_active:
        return 10  # live_processing
    elif stage == "Live Partner" and int_used:
        return 9  # g2m
    elif stage == "Live Partner" and int_ready:
        return 8  # pilot
    elif stage == "Live Partner":
        return 7  # go_live
    elif stage == "Agreement Signed" and int_ready:
        return 6  # integration_complete
    elif stage == "Agreement Signed" and nda:
        return 5  # in_development
    elif stage == "Agreement Signed":
        return 4  # prioritized
    elif stage == "Agreement Review":
        return 2  # agreement review
    elif stage == "Initial Negotiation":
        return 1  # negotiation
    elif stage == "Opportunity Identification":
        return 0  # prospect
    else:
        return 0

def _get_pending_task(row, stage_idx):
    """Determine who currently holds the ball and what's pending."""
    nda = bool(row.get("NDA Signed and in drive", False))
    revshare = bool(row.get("Revshare Contract", False))
    int_ready = bool(row.get("Integration Ready by Yuno", False))

    if stage_idx == 0:
        return ("Partnership BD", "Qualify opportunity & initial outreach")
    elif stage_idx == 1:
        return ("Partnership BD", "Align commercial terms & pricing")
    elif stage_idx == 2:
        if not nda:
            return ("Legal", "NDA pending signature")
        return ("Legal", "Contract under review")
    elif stage_idx == 3:
        return ("Partnerships + Legal", "Awaiting final contract sign-off")
    elif stage_idx == 4:
        return ("Product", "Awaiting roadmap slot & resource allocation")
    elif stage_idx == 5:
        return ("Integrations", "API integration in progress")
    elif stage_idx == 6:
        return ("Integrations", "QA certification & sandbox testing")
    elif stage_idx == 7:
        return ("Merchant BD/KAM", "Production go-live & merchant activation")
    elif stage_idx == 8:
        return ("Merchant BD/KAM", "First merchant onboarding & volume ramp")
    elif stage_idx == 9:
        return ("Marketing + Partnerships", "GTM campaigns, co-marketing & case studies")
    elif stage_idx == 10:
        if not revshare:
            return ("Partnership BD", "Rev share contract pending")
        return ("Partnership Acct Mgmt", "Ongoing management & optimization")
    return ("—", "—")

@st.cache_data
def load_mission_control_data():
    try:
        df = pd.read_excel(_os.path.join(_os.path.dirname(_os.path.abspath(__file__)), "data", "strategic_accounts.xlsx"), sheet_name="All partners")
        strat = df[df["Strategic?"] == True].drop_duplicates(subset="Provider")
        partners = []
        for _, row in strat.iterrows():
            name = str(row.get("Provider","")).strip()
            if not name or name == "nan":
                continue
            mc_stage = _map_deal_stage_to_mc(row)
            pending_owner, pending_task = _get_pending_task(row, mc_stage)
            partners.append({
                "name": name,
                "region": str(row.get("Region","")).strip(),
                "country": str(row.get("Country","")).strip(),
                "manager": str(row.get("Partner Manager","")).strip(),
                "offering": str(row.get("Partner offering (acquirer, APM, Fraud etc)","")).strip(),
                "mgmt_type": str(row.get("Type of Management","")).strip(),
                "stage_idx": mc_stage,
                "deal_stage": str(row.get("Deal Stage","")).strip(),
                "nda": bool(row.get("NDA Signed and in drive", False)),
                "int_ready": bool(row.get("Integration Ready by Yuno", False)),
                "int_used": bool(row.get("Integration Used by Merchants", False)),
                "revshare": bool(row.get("Revshare Contract", False)),
                "revshare_active": bool(row.get("Revshare active", False)),
                "pending_owner": pending_owner,
                "pending_task": pending_task,
            })
        return sorted(partners, key=lambda x: (-x["stage_idx"], x["name"]))
    except Exception:
        return []

def show_mission_control():
    mc_data = load_mission_control_data()

    # Build JSON for HTML component
    import json as _json
    _mc_json = _json.dumps(mc_data).replace("'", "\\'").replace("</", "<\\/")

    _html_file = _os.path.join(_BASE, "partners_in_flight.html")
    with open(_html_file, "r") as _f:
        _html_content = _f.read()
    _html_content = _html_content.replace("var DATA = [];", f"var DATA = {_mc_json};")
    components.html(_html_content, height=1800, scrolling=True)
    return

    # ── OLD CODE (kept for reference) ─────────────────────────────────────────
    stage_counts = [0] * len(_MC_STAGES)
    for p in mc_data:
        stage_counts[p["stage_idx"]] += 1
    total_mc = len(mc_data)
    live_mc = stage_counts[10] if len(stage_counts) > 10 else 0
    in_dev = stage_counts[5] + stage_counts[6] if len(stage_counts) > 6 else 0

    st.markdown(f"""
<div style="background:#0C1220;border-radius:16px;padding:24px 28px;margin-bottom:14px;position:relative;overflow:hidden;">
  <div style="position:absolute;inset:0;background-image:radial-gradient(circle at 1px 1px,rgba(255,255,255,0.03) 1px,transparent 0);background-size:22px 22px;pointer-events:none;"></div>
  <div style="position:relative;z-index:1;">
    <div style="display:flex;align-items:center;gap:8px;margin-bottom:10px;">
      <span style="width:6px;height:6px;border-radius:50%;background:#22c55e;"></span>
      <span style="font-size:10px;font-weight:700;letter-spacing:1.2px;text-transform:uppercase;color:#4ade80;">Live Tracker</span>
    </div>
    <div style="font-size:22px;font-weight:800;color:#F1F5F9;letter-spacing:-0.02em;margin-bottom:4px;">Partners In Flight</div>
    <div style="font-size:12px;color:#64748b;margin-bottom:14px;">End-to-end lifecycle of every strategic partner — from prospect to live processing.</div>
    <div style="display:flex;gap:10px;flex-wrap:wrap;">
      <div style="background:rgba(255,255,255,0.04);border:0.5px solid rgba(255,255,255,0.08);border-radius:8px;padding:10px 14px;text-align:center;"><div style="font-size:18px;font-weight:800;color:#C7D2FE;">{total_mc}</div><div style="font-size:9px;color:#475569;text-transform:uppercase;margin-top:2px;">Strategic</div></div>
      <div style="background:rgba(255,255,255,0.04);border:0.5px solid rgba(255,255,255,0.08);border-radius:8px;padding:10px 14px;text-align:center;"><div style="font-size:18px;font-weight:800;color:#6EE7B7;">{live_mc}</div><div style="font-size:9px;color:#475569;text-transform:uppercase;margin-top:2px;">Live</div></div>
      <div style="background:rgba(255,255,255,0.04);border:0.5px solid rgba(255,255,255,0.08);border-radius:8px;padding:10px 14px;text-align:center;"><div style="font-size:18px;font-weight:800;color:#818cf8;">{in_dev}</div><div style="font-size:9px;color:#475569;text-transform:uppercase;margin-top:2px;">In Dev</div></div>
    </div>
  </div>
</div>""", unsafe_allow_html=True)

    # ── Region toggles ────────────────────────────────────────────────────────
    mc_regions = sorted(set(p["region"] for p in mc_data if p["region"] and p["region"] != "nan"))
    _region_colors = {"Global":"#5B5FDE","LATAM":"#4F46E5","Brasil":"#10B981","EMEA":"#8B5CF6","APAC":"#0EA5E9","North America":"#059669","Africa":"#D97706"}

    sel_region = st.session_state.get("mc_region_sel", "All")
    region_opts = ["All"] + mc_regions
    reg_cols = st.columns(len(region_opts))
    for ri, rname in enumerate(region_opts):
        rcount = len([p for p in mc_data if p["region"] == rname]) if rname != "All" else total_mc
        if reg_cols[ri].button(f"{rname} ({rcount})", key=f"mc_reg_{ri}", use_container_width=True):
            st.session_state["mc_region_sel"] = rname
            st.rerun()

    # Style region buttons
    components.html(f"""<script>
setTimeout(function(){{
  var doc = window.parent.document;
  var active = '{sel_region}';
  doc.querySelectorAll('button').forEach(function(btn){{
    var t = btn.textContent.trim();
    if(t.match(/^(All|Global|LATAM|Brasil|EMEA|APAC|North America|Africa)\s*\(/)){{
      var isOn = t.startsWith(active+' (') || (active==='All' && t.startsWith('All ('));
      btn.style.cssText='background:'+(isOn?'#0C1220':'#fff')+' !important;border:1px solid '+(isOn?'#5B5FDE':'#E2E8F0')+' !important;border-radius:8px !important;padding:8px 10px !important;font-size:11px !important;color:'+(isOn?'#A5B4FC':'#64748b')+' !important;font-weight:'+(isOn?'700':'500')+' !important;cursor:pointer !important;box-shadow:none !important;';
    }}
  }});
}},200);
</script>""", height=0)

    # ── Filters ───────────────────────────────────────────────────────────────
    col_f1, col_f2 = st.columns(2)
    with col_f1:
        mc_owners = sorted(set(p["manager"] for p in mc_data if p["manager"] and p["manager"] != "nan"))
        mc_owner_filter = st.selectbox("Manager", ["Manager"] + mc_owners, key="mc_owner_f", label_visibility="collapsed")
    with col_f2:
        active_stage_filter = st.selectbox("Stage", ["Stage"] + [s["label"] for s in _MC_STAGES], key="mc_stage_sel", label_visibility="collapsed")

    # ── Filter logic ──────────────────────────────────────────────────────────
    filtered = list(mc_data)
    if sel_region != "All":
        filtered = [p for p in filtered if p["region"] == sel_region]
    if mc_owner_filter != "Manager":
        filtered = [p for p in filtered if p["manager"] == mc_owner_filter]
    if active_stage_filter != "Stage":
        stage_idx = next((i for i, s in enumerate(_MC_STAGES) if s["label"] == active_stage_filter), -1)
        if stage_idx >= 0:
            filtered = [p for p in filtered if p["stage_idx"] == stage_idx]

    # ── Pipeline progress bar ─────────────────────────────────────────────────
    filt_counts = [0] * len(_MC_STAGES)
    for p in filtered:
        filt_counts[p["stage_idx"]] += 1
    vis_total = sum(filt_counts) or 1
    bar_html = '<div style="display:flex;gap:1px;height:6px;margin:10px 0 14px;border-radius:3px;overflow:hidden;">'
    for si, s in enumerate(_MC_STAGES):
        w = max(int(filt_counts[si] / vis_total * 100), 1) if filt_counts[si] > 0 else 0
        if w > 0:
            bar_html += f'<div style="width:{w}%;background:{s["owner_color"]};"></div>'
    bar_html += '</div>'
    st.markdown(bar_html, unsafe_allow_html=True)

    st.markdown(f'<div style="font-size:12px;color:#94a3b8;margin-bottom:8px;"><b style="color:#0f172a;">{len(filtered)}</b> partners</div>', unsafe_allow_html=True)

    # ── Partner list by stage (expanders) ─────────────────────────────────────
    owner_colors = {"Legal":"#8b5cf6","Partnership BD":"#6366f1","Partnerships + Legal":"#7c3aed","Product":"#0891b2","Integrations":"#2563eb","Merchant BD/KAM":"#d97706","Marketing + Partnerships":"#e11d48","Partnership Acct Mgmt":"#065f46"}

    for si, s in enumerate(_MC_STAGES):
        partners_in_stage = [p for p in filtered if p["stage_idx"] == si]
        if not partners_in_stage:
            continue

        st.markdown(f'<div style="display:flex;align-items:center;gap:8px;margin:12px 0 4px;padding-bottom:4px;border-bottom:2px solid {s["owner_color"]}20;">'
                    f'<span style="font-size:13px;font-weight:700;color:#0f172a;">{s["label"]}</span>'
                    f'<span style="font-size:10px;font-weight:600;padding:2px 8px;border-radius:5px;background:{s["owner_color"]}10;color:{s["owner_color"]};">{s["owner"]}</span>'
                    f'<span style="margin-left:auto;font-size:12px;font-weight:800;color:#0f172a;">{len(partners_in_stage)}</span></div>', unsafe_allow_html=True)

        for p in partners_in_stage:
            po = p.get("pending_owner", "")
            pt = p.get("pending_task", "")
            po_color = owner_colors.get(po, "#94a3b8")
            checks = [("NDA", p["nda"]), ("Int", p["int_ready"]), ("Merch", p["int_used"]), ("Rev$", p["revshare"])]
            ch = " ".join(f'<span style="font-size:8px;padding:1px 5px;border-radius:3px;background:{"#ecfdf5" if v else "#f8fafc"};color:{"#065f46" if v else "#d1d5db"};">{l}</span>' for l, v in checks)

            with st.expander(f"{p['name']}  ·  {p['offering']}  ·  {p['region']}  ·  {p['manager']}", expanded=False):
                st.markdown(f'<div style="display:flex;align-items:center;gap:8px;margin-bottom:10px;">'
                            f'<span style="font-size:9px;font-weight:700;padding:3px 8px;border-radius:4px;background:{po_color}10;color:{po_color};border:0.5px solid {po_color}20;">{po}</span>'
                            f'<span style="font-size:11px;color:#64748b;">{pt}</span></div>', unsafe_allow_html=True)
                st.markdown(f'<div style="display:flex;gap:3px;margin-bottom:8px;">{ch}</div>', unsafe_allow_html=True)
                st.markdown(f'<div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:8px;">'
                            f'<div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;">Country</div><div style="font-size:12px;font-weight:600;color:#0f172a;">{p["country"]}</div></div>'
                            f'<div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;">Deal Stage</div><div style="font-size:12px;font-weight:600;color:#0f172a;">{p["deal_stage"]}</div></div>'
                            f'<div><div style="font-size:9px;color:#94a3b8;text-transform:uppercase;">Management</div><div style="font-size:12px;font-weight:600;color:#0f172a;">{p.get("mgmt_type","—")}</div></div>'
                            f'</div>', unsafe_allow_html=True)


# ── Merchant Simulation ──────────────────────────────────────────────────────
_MERCHANT_PROFILES = {
    "uber":      {"name":"Uber","sector":"Mobility & Delivery","countries":["MX","BR","CO","AR","CL"],"aov":18.2,"txn_mo":323000,"ar":93.4,"methods":["Card Visa","Card Mastercard","PIX","OXXO","Nequi"],"providers":["Adyen","dLocal","Stripe"],"color":"#1a1a1a"},
    "rappi":     {"name":"Rappi","sector":"Super App / Delivery","countries":["CO","MX","BR","CL","PE"],"aov":32.5,"txn_mo":259000,"ar":91.2,"methods":["Card Visa","Card Mastercard","PSE","Nequi","PIX","OXXO"],"providers":["Adyen","Kushki","MercadoPago"],"color":"#FF5722"},
    "netflix":   {"name":"Netflix","sector":"Streaming / Subscriptions","countries":["US","MX","BR","AR","CO"],"aov":15.5,"txn_mo":306000,"ar":95.6,"methods":["Card Visa","Card Mastercard","Amex"],"providers":["Stripe","Adyen","dLocal"],"color":"#E50914"},
    "spotify":   {"name":"Spotify","sector":"Streaming / Subscriptions","countries":["BR","MX","AR","CO","CL"],"aov":9.99,"txn_mo":421000,"ar":94.8,"methods":["Card Visa","Card Mastercard","PIX","Boleto"],"providers":["Adyen","Ebanx","dLocal"],"color":"#1DB954"},
    "mercadolibre":{"name":"MercadoLibre","sector":"Marketplace / eCommerce","countries":["AR","BR","MX","CO","CL"],"aov":118.4,"txn_mo":61000,"ar":88.7,"methods":["Card Visa","Card Mastercard","PIX","Boleto","MercadoPago"],"providers":["MercadoPago","Adyen","Cielo"],"color":"#FFE600"},
    "falabella": {"name":"Falabella","sector":"Retail / eCommerce","countries":["CL","CO","PE","AR"],"aov":87.3,"txn_mo":35000,"ar":87.1,"methods":["Card Visa","Card Mastercard","Webpay","PSE"],"providers":["Transbank","Kushki","Adyen"],"color":"#009900"},
    "ifood":     {"name":"iFood","sector":"Food Delivery","countries":["BR"],"aov":28.4,"txn_mo":103000,"ar":91.8,"methods":["Card Visa","Card Mastercard","PIX","Boleto"],"providers":["Cielo","PagBank","Stone"],"color":"#EA1D2C"},
    "despegar":  {"name":"Despegar","sector":"Travel / OTA","countries":["AR","BR","MX","CO","CL"],"aov":312.5,"txn_mo":12000,"ar":82.3,"methods":["Card Visa","Card Mastercard","Amex","PIX","Installments"],"providers":["Cybersource","Adyen","dLocal"],"color":"#0066CC"},
}

_OPTIMIZATION_TIPS = {
    "low_ar": {"title":"Approval Rate Below 90%","desc":"Enable smart retries across multiple acquirers. Route declines to a secondary provider — this alone can recover 3-5% of failed transactions.","impact":"+3-5% approval rate"},
    "single_provider": {"title":"Single Provider Risk","desc":"Relying on one provider creates a single point of failure. Add a backup acquirer for failover routing and cost negotiation leverage.","impact":"+99.9% uptime, -15% processing cost"},
    "no_apm": {"title":"Missing Local Payment Methods","desc":"Offer local APMs (PIX in Brazil, OXXO in Mexico, PSE in Colombia). Merchants see 15-25% conversion lift when local methods are available.","impact":"+15-25% conversion"},
    "no_installments": {"title":"No Installment Options","desc":"In LATAM, 60%+ of card transactions use installments. Enable parcelamento/MSI to unlock higher AOV and reduce cart abandonment.","impact":"+20% AOV"},
    "high_aov": {"title":"High-AOV Without Fraud Layer","desc":"Transactions above $100 need a fraud prevention layer. Add 3DS and a fraud partner (SEON, Riskified) to reduce chargebacks.","impact":"-40% chargebacks"},
    "cross_border": {"title":"Cross-Border Processing","desc":"Local acquiring delivers 8-12% higher approval rates vs cross-border. Establish local entity or use a provider with local acquiring license.","impact":"+8-12% approval rate"},
}

def show_merchant_sim():
    # ── Hero ──────────────────────────────────────────────────────────────────
    st.markdown("""
<div style="background:linear-gradient(160deg,#0a0e1a 0%,#0f1629 40%,#1a1040 100%);border-radius:20px;padding:44px 44px 36px;margin-bottom:4px;position:relative;overflow:hidden;">
  <div style="position:absolute;top:-60px;right:-40px;width:350px;height:350px;border-radius:50%;background:radial-gradient(circle,rgba(168,85,247,0.1),transparent 65%);"></div>
  <div style="position:relative;z-index:1;">
    <div style="display:inline-flex;align-items:center;gap:8px;padding:4px 14px;border-radius:20px;background:rgba(168,85,247,0.12);border:1px solid rgba(168,85,247,0.2);margin-bottom:16px;">
      <span style="font-size:14px;">◇</span>
      <span style="font-size:10px;font-weight:700;letter-spacing:1.5px;text-transform:uppercase;color:#c4b5fd;">Merchant Simulation</span>
    </div>
    <div style="font-size:28px;font-weight:800;color:#fff;letter-spacing:-0.8px;line-height:1.2;margin-bottom:8px;">Simulate. Optimize. Close.</div>
    <div style="font-size:13px;color:#94a3b8;line-height:1.7;max-width:520px;">Enter a merchant name to generate a checkout flow mockup, identify payment optimization opportunities, and model the revenue impact of adding new Yuno partners to their stack.</div>
  </div>
</div>""", unsafe_allow_html=True)

    # ── Search bar ────────────────────────────────────────────────────────────
    st.markdown("""
<div style="background:#fff;border:1px solid #e2e8f0;border-radius:0 0 20px 20px;padding:18px 28px 10px;margin-bottom:16px;box-shadow:0 8px 32px rgba(0,0,0,0.04);">
  <div style="font-size:10px;font-weight:700;color:#8b5cf6;text-transform:uppercase;letter-spacing:1.2px;margin-bottom:4px;">Enter merchant name</div>
</div>""", unsafe_allow_html=True)

    merchant_input = st.text_input("", placeholder="e.g. Uber, Rappi, Netflix, Spotify, Falabella...", key="merchsim_input", label_visibility="collapsed")

    # Use quick-select override if set
    active_merchant = merchant_input or st.session_state.get("merchsim_selected", "")

    if not active_merchant:
        # Landing state
        st.markdown("""
<div style="border:1px solid #e2e8f0;border-radius:14px;padding:48px 40px;text-align:center;background:#fff;margin-top:4px;">
  <div style="font-size:36px;margin-bottom:14px;opacity:0.2;">◇</div>
  <div style="font-size:18px;font-weight:700;color:#0f172a;margin-bottom:8px;">Enter a merchant to begin</div>
  <div style="font-size:13px;color:#94a3b8;max-width:440px;margin:0 auto;line-height:1.6;">Type a merchant name above. We'll generate their checkout flow, flag optimization gaps, and let you model the impact of adding new payment partners.</div>
</div>""", unsafe_allow_html=True)

        # Quick select
        st.markdown('<div style="font-size:10px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:1px;margin:20px 0 8px;">Quick select</div>', unsafe_allow_html=True)
        qcols = st.columns(4)
        for qi, (mkey, mdata) in enumerate(list(_MERCHANT_PROFILES.items())[:8]):
            if qcols[qi % 4].button(mdata["name"], key=f"qs_{mkey}", use_container_width=True):
                st.session_state["merchsim_selected"] = mdata["name"]
                st.rerun()

        components.html("""<script>
setTimeout(function(){
  var doc = window.parent.document;
  var names = ['Uber','Rappi','Netflix','Spotify','MercadoLibre','Falabella','iFood','Despegar'];
  doc.querySelectorAll('button').forEach(function(btn){
    if(names.indexOf(btn.textContent.trim()) !== -1){
      btn.style.cssText='background:#faf5ff !important;border:1px solid #e9d5ff !important;border-radius:10px !important;padding:10px !important;font-size:13px !important;color:#7c3aed !important;font-weight:700 !important;cursor:pointer !important;box-shadow:none !important;';
    }
  });
},250);
</script>""", height=0)
        return
    else:
        # Clear quick-select if user typed something
        if merchant_input:
            st.session_state["merchsim_selected"] = ""

    # ── Find merchant profile ─────────────────────────────────────────────────
    # Back button if loaded via quick-select
    if st.session_state.get("merchsim_selected"):
        if st.button("← Back to merchant select", key="merchsim_back"):
            st.session_state["merchsim_selected"] = ""
            st.rerun()

    m_key = active_merchant.lower().replace(" ", "")
    profile = _MERCHANT_PROFILES.get(m_key)

    if not profile:
        # Generate a generic profile
        profile = {
            "name": active_merchant, "sector": "eCommerce", "countries": ["BR","MX","CO"],
            "aov": 45.0, "txn_mo": 50000, "ar": 88.0,
            "methods": ["Card Visa","Card Mastercard"],
            "providers": ["Adyen"], "color": "#4F46E5",
        }

    m = profile
    tpv_mo = m["aov"] * m["txn_mo"]
    tpv_display = f"${tpv_mo/1e6:.1f}M" if tpv_mo >= 1e6 else f"${tpv_mo/1e3:.0f}K"
    country_names = [_ISO_TO_COUNTRY.get(c, c) for c in m["countries"]]

    # ── Merchant header ───────────────────────────────────────────────────────
    st.markdown(f'<div style="background:#fff;border:1px solid #e2e8f0;border-radius:14px;padding:22px 26px;margin-bottom:14px;">'
                f'<div style="display:flex;justify-content:space-between;align-items:center;">'
                f'<div style="display:flex;align-items:center;gap:16px;">'
                f'<div style="width:48px;height:48px;border-radius:12px;background:{m["color"]};display:flex;align-items:center;justify-content:center;font-size:18px;font-weight:800;color:#fff;">{m["name"][0]}</div>'
                f'<div><div style="font-size:20px;font-weight:800;color:#0f172a;letter-spacing:-0.5px;">{m["name"]}</div>'
                f'<div style="font-size:12px;color:#94a3b8;margin-top:1px;">{m["sector"]} &middot; {", ".join(country_names)}</div></div></div>'
                f'<div style="display:flex;gap:12px;">'
                f'<div style="text-align:center;"><div style="font-size:18px;font-weight:800;color:#0f172a;">{tpv_display}</div><div style="font-size:9px;color:#94a3b8;">TPV/mo</div></div>'
                f'<div style="text-align:center;"><div style="font-size:18px;font-weight:800;color:#0f172a;">${m["aov"]:.0f}</div><div style="font-size:9px;color:#94a3b8;">AOV</div></div>'
                f'<div style="text-align:center;"><div style="font-size:18px;font-weight:800;color:{"#22c55e" if m["ar"]>=90 else "#f59e0b" if m["ar"]>=85 else "#ef4444"};">{m["ar"]}%</div><div style="font-size:9px;color:#94a3b8;">Auth Rate</div></div>'
                f'<div style="text-align:center;"><div style="font-size:18px;font-weight:800;color:#0f172a;">{m["txn_mo"]:,}</div><div style="font-size:9px;color:#94a3b8;">Txn/mo</div></div>'
                f'</div></div></div>', unsafe_allow_html=True)

    # ── Checkout Flow Mockup ──────────────────────────────────────────────────
    st.markdown(f'<div style="font-size:11px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:1.5px;margin:8px 0 10px;">Checkout Flow Simulation</div>', unsafe_allow_html=True)

    steps = [
        ("Cart","Customer adds items","#94a3b8"),
        ("Checkout","Selects payment method","#6366f1"),
        ("Auth","Transaction processed","#f59e0b"),
        ("3DS","Identity verified","#8b5cf6"),
        ("Approve","Payment confirmed","#22c55e"),
    ]
    steps_html = ""
    for si, (step, desc, scolor) in enumerate(steps):
        arrow = '<div style="font-size:16px;color:#d1d5db;margin:0 4px;">&#8594;</div>' if si < len(steps)-1 else ""
        steps_html += f'<div style="display:flex;align-items:center;gap:0;"><div style="text-align:center;padding:14px 16px;background:#fff;border:1px solid #e2e8f0;border-radius:12px;min-width:100px;"><div style="font-size:13px;font-weight:700;color:{scolor};">{step}</div><div style="font-size:10px;color:#94a3b8;margin-top:2px;">{desc}</div></div>{arrow}</div>'
    st.markdown(f'<div style="display:flex;align-items:center;justify-content:center;gap:0;padding:16px;background:#fafbff;border:1px solid #e2e8f0;border-radius:14px;margin-bottom:14px;">{steps_html}</div>', unsafe_allow_html=True)

    # Payment methods mockup
    methods_html = "".join(f'<div style="display:flex;align-items:center;gap:8px;padding:10px 14px;background:#fff;border:1px solid #e2e8f0;border-radius:10px;"><span style="width:8px;height:8px;border-radius:50%;background:#22c55e;"></span><span style="font-size:12px;font-weight:600;color:#0f172a;">{pm}</span></div>' for pm in m["methods"])
    providers_html = "".join(f'<span style="font-size:11px;font-weight:600;padding:5px 14px;border-radius:8px;background:#eef2ff;color:#4F46E5;border:1px solid #e0e7ff;">{pv}</span>' for pv in m["providers"])

    col_pm, col_pv = st.columns(2)
    with col_pm:
        st.markdown(f'<div style="background:#fff;border:1px solid #e2e8f0;border-radius:12px;padding:18px;">'
                    f'<div style="font-size:11px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:0.8px;margin-bottom:10px;">Active Payment Methods</div>'
                    f'<div style="display:flex;flex-wrap:wrap;gap:6px;">{methods_html}</div></div>', unsafe_allow_html=True)
    with col_pv:
        st.markdown(f'<div style="background:#fff;border:1px solid #e2e8f0;border-radius:12px;padding:18px;">'
                    f'<div style="font-size:11px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:0.8px;margin-bottom:10px;">Current Providers</div>'
                    f'<div style="display:flex;flex-wrap:wrap;gap:6px;">{providers_html}</div></div>', unsafe_allow_html=True)

    # ── Optimization Tips ─────────────────────────────────────────────────────
    st.markdown(f'<div style="font-size:11px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:1.5px;margin:20px 0 10px;">Optimization Opportunities</div>', unsafe_allow_html=True)

    # Determine which tips apply
    tips = []
    if m["ar"] < 90:
        tips.append(_OPTIMIZATION_TIPS["low_ar"])
    if len(m["providers"]) <= 1:
        tips.append(_OPTIMIZATION_TIPS["single_provider"])
    if not any(pm in ["PIX","OXXO","PSE","Nequi","Boleto"] for pm in m["methods"]):
        tips.append(_OPTIMIZATION_TIPS["no_apm"])
    if "Installments" not in " ".join(m["methods"]) and m["aov"] > 30:
        tips.append(_OPTIMIZATION_TIPS["no_installments"])
    if m["aov"] > 100:
        tips.append(_OPTIMIZATION_TIPS["high_aov"])
    if len(m["countries"]) > 2 and len(m["providers"]) < 3:
        tips.append(_OPTIMIZATION_TIPS["cross_border"])
    if not tips:
        tips.append({"title":"Well Optimized","desc":"This merchant has a solid payment setup. Consider fine-tuning routing rules and exploring new markets for growth.","impact":"Maintain & grow"})

    for tip in tips:
        st.markdown(f'<div style="background:#fff;border:1px solid #e2e8f0;border-radius:12px;padding:16px 20px;margin-bottom:8px;">'
                    f'<div style="display:flex;justify-content:space-between;align-items:flex-start;">'
                    f'<div style="flex:1;"><div style="font-size:14px;font-weight:700;color:#0f172a;margin-bottom:4px;">{tip["title"]}</div>'
                    f'<div style="font-size:12px;color:#64748b;line-height:1.6;">{tip["desc"]}</div></div>'
                    f'<span style="font-size:11px;font-weight:700;padding:4px 12px;border-radius:8px;background:#ecfdf5;color:#065f46;white-space:nowrap;margin-left:16px;">{tip["impact"]}</span>'
                    f'</div></div>', unsafe_allow_html=True)

    # ── Partner Impact Simulator ──────────────────────────────────────────────
    st.markdown(f'<div style="font-size:11px;font-weight:700;color:#94a3b8;text-transform:uppercase;letter-spacing:1.5px;margin:24px 0 10px;">Partner Impact Simulator</div>', unsafe_allow_html=True)

    st.markdown('<div style="font-size:12px;color:#64748b;margin-bottom:10px;">Add a payment partner and see the projected impact on this merchant.</div>', unsafe_allow_html=True)

    col_add, col_country = st.columns(2)
    with col_add:
        add_options = ["Select a partner..."] + _SOT_PROVIDERS[:50]
        add_partner = st.selectbox("Add Partner", add_options, key="sim_add_partner", label_visibility="collapsed")
    with col_country:
        sim_country = st.selectbox("Target Country", [""] + _SOT_COUNTRIES, key="sim_country",
                                   format_func=lambda x: "Select country..." if x == "" else x, label_visibility="collapsed")

    if add_partner != "Select a partner..." and sim_country:
        # Find partner capabilities
        sim_iso = None
        for iso, name in _ISO_TO_COUNTRY.items():
            if name == sim_country:
                sim_iso = iso
                break

        sim_results = find_partners(country_iso=sim_iso, live_only=False) if sim_iso else []
        partner_match = next((r for r in sim_results if r["name"].upper() == add_partner.upper()), None)

        # Projected impact
        ar_boost = 3.5 if len(m["providers"]) <= 1 else 1.8
        new_ar = min(m["ar"] + ar_boost, 99.5)
        recovered_txn = int(m["txn_mo"] * (ar_boost / 100))
        revenue_gain = recovered_txn * m["aov"]
        rev_display = f"${revenue_gain/1e3:.0f}K" if revenue_gain >= 1000 else f"${revenue_gain:.0f}"

        st.markdown(f'<div style="background:linear-gradient(160deg,#0a0e1a,#111827);border-radius:14px;padding:24px 28px;margin-top:10px;">'
                    f'<div style="font-size:14px;font-weight:700;color:#fff;margin-bottom:16px;">Projected Impact: Adding {add_partner} in {sim_country}</div>'
                    f'<div style="display:grid;grid-template-columns:repeat(4,1fr);gap:12px;">'
                    f'<div style="text-align:center;padding:14px;background:rgba(255,255,255,0.04);border-radius:10px;border:1px solid rgba(255,255,255,0.06);"><div style="font-size:22px;font-weight:800;color:#22c55e;">{new_ar:.1f}%</div><div style="font-size:9px;color:#64748b;margin-top:3px;">New Auth Rate</div><div style="font-size:10px;color:#22c55e;margin-top:2px;">+{ar_boost:.1f}%</div></div>'
                    f'<div style="text-align:center;padding:14px;background:rgba(255,255,255,0.04);border-radius:10px;border:1px solid rgba(255,255,255,0.06);"><div style="font-size:22px;font-weight:800;color:#a5b4fc;">{recovered_txn:,}</div><div style="font-size:9px;color:#64748b;margin-top:3px;">Recovered Txns/mo</div><div style="font-size:10px;color:#a5b4fc;margin-top:2px;">previously declined</div></div>'
                    f'<div style="text-align:center;padding:14px;background:rgba(255,255,255,0.04);border-radius:10px;border:1px solid rgba(255,255,255,0.06);"><div style="font-size:22px;font-weight:800;color:#fbbf24;">{rev_display}</div><div style="font-size:9px;color:#64748b;margin-top:3px;">Revenue Recovered/mo</div><div style="font-size:10px;color:#fbbf24;margin-top:2px;">from retries</div></div>'
                    f'<div style="text-align:center;padding:14px;background:rgba(255,255,255,0.04);border-radius:10px;border:1px solid rgba(255,255,255,0.06);"><div style="font-size:22px;font-weight:800;color:#fff;">{len(m["providers"])+1}</div><div style="font-size:9px;color:#64748b;margin-top:3px;">Total Providers</div><div style="font-size:10px;color:#22c55e;margin-top:2px;">redundancy added</div></div>'
                    f'</div></div>', unsafe_allow_html=True)

        # Partner details if found
        if partner_match:
            pm_list = ", ".join(partner_match["payment_methods"][:5])
            st.markdown(f'<div style="background:#fff;border:1px solid #e2e8f0;border-radius:12px;padding:16px 20px;margin-top:8px;">'
                        f'<div style="font-size:12px;font-weight:700;color:#0f172a;margin-bottom:8px;">{add_partner} in {sim_country}</div>'
                        f'<div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:12px;">'
                        f'<div><div style="font-size:10px;color:#94a3b8;">Status</div><div style="font-size:13px;font-weight:600;color:#0f172a;">{partner_match["status"]}</div></div>'
                        f'<div><div style="font-size:10px;color:#94a3b8;">Payment Methods</div><div style="font-size:12px;color:#334155;">{pm_list}</div></div>'
                        f'<div><div style="font-size:10px;color:#94a3b8;">Processing</div><div style="font-size:13px;font-weight:600;color:#0f172a;">{" / ".join(partner_match["processing_types"])}</div></div>'
                        f'</div></div>', unsafe_allow_html=True)

        # Before/After comparison
        st.markdown(f'<div style="display:grid;grid-template-columns:1fr 1fr;gap:10px;margin-top:10px;">'
                    f'<div style="background:#fff;border:1px solid #e2e8f0;border-radius:12px;padding:16px 20px;">'
                    f'<div style="font-size:11px;font-weight:700;color:#ef4444;text-transform:uppercase;letter-spacing:0.8px;margin-bottom:10px;">Before</div>'
                    f'<div style="margin-bottom:6px;"><span style="font-size:12px;color:#94a3b8;">Auth Rate</span><span style="font-size:14px;font-weight:700;color:#0f172a;float:right;">{m["ar"]}%</span></div>'
                    f'<div style="margin-bottom:6px;"><span style="font-size:12px;color:#94a3b8;">Providers</span><span style="font-size:14px;font-weight:700;color:#0f172a;float:right;">{len(m["providers"])}</span></div>'
                    f'<div><span style="font-size:12px;color:#94a3b8;">Failed Txns/mo</span><span style="font-size:14px;font-weight:700;color:#ef4444;float:right;">{int(m["txn_mo"]*(1-m["ar"]/100)):,}</span></div>'
                    f'</div>'
                    f'<div style="background:#fff;border:1px solid #a7f3d0;border-radius:12px;padding:16px 20px;">'
                    f'<div style="font-size:11px;font-weight:700;color:#22c55e;text-transform:uppercase;letter-spacing:0.8px;margin-bottom:10px;">After (Projected)</div>'
                    f'<div style="margin-bottom:6px;"><span style="font-size:12px;color:#94a3b8;">Auth Rate</span><span style="font-size:14px;font-weight:700;color:#22c55e;float:right;">{new_ar:.1f}%</span></div>'
                    f'<div style="margin-bottom:6px;"><span style="font-size:12px;color:#94a3b8;">Providers</span><span style="font-size:14px;font-weight:700;color:#0f172a;float:right;">{len(m["providers"])+1}</span></div>'
                    f'<div><span style="font-size:12px;color:#94a3b8;">Failed Txns/mo</span><span style="font-size:14px;font-weight:700;color:#22c55e;float:right;">{int(m["txn_mo"]*(1-new_ar/100)):,}</span></div>'
                    f'</div></div>', unsafe_allow_html=True)


# ── Home / Welcome Panel ────────────────────────────────────────────────────────
def show_home():
    is_internal = st.session_state.role == "internal"
    badge = "INTERNAL" if is_internal else "PARTNER"
    headline = "Welcome back, Daniela" if is_internal else "Welcome, Partner"
    sub = "Your partnership command center — everything in one place." if is_internal else "Your shared workspace with Yuno."

    st.markdown(f"""
<div style="padding:24px 8px 0;">
  <div style="display:flex;align-items:center;gap:10px;margin-bottom:16px;">
    <span style="font-size:9px;font-weight:700;letter-spacing:0.1em;background:rgba(91,95,222,0.12);color:#5B5FDE;border:0.5px solid rgba(91,95,222,0.25);padding:3px 10px;border-radius:5px;text-transform:uppercase;">{badge}</span>
  </div>
  <div style="font-size:26px;font-weight:800;color:#0f172a;letter-spacing:-0.8px;line-height:1.2;margin-bottom:6px;">{headline}</div>
  <p style="font-size:13px;color:#64748b;margin-bottom:24px;">{sub}</p>
</div>""", unsafe_allow_html=True)

    if is_internal:
        # 4 action cards — grouped by purpose
        cards = [
            ("Partners & Connectors", "#5B5FDE", "◈",
             "Yuno Connectors — search, filter and explore 460+ payment partners. Partners In Flight — track every deal from prospect to live.",
             [("Partner Portfolio","Partners"),("Partners In Flight","MissionCtrl")]),
            ("Intelligence", "#0EA5E9", "✦",
             "Market insights across LATAM, MENAT, APAC and Europe. Coverage analysis and strategic context.",
             [("Market Intel","Insights")]),
            ("Pipeline & Merchants", "#10B981", "⟳",
             "24 active deals in pipeline. Live merchant analytics — TPV, AOV, approval rates across all merchants.",
             [("Partner Leads","Pipeline"),("Merchants","Merchants")]),
            ("Performance & Tools", "#F59E0B", "▲",
             "Partner health dashboards, rev share tracking ($123K/mo), and merchant checkout simulator.",
             [("Partner Health","Performance"),("Rev Share","Benchmarks"),("Merchant Sim","MerchSim")]),
        ]
    else:
        cards = [
            ("Partners", "#5B5FDE", "◈",
             "Browse partner details, contacts, contracts, SLAs and integration status.",
             [("Partner Portfolio","Partners")]),
            ("Pipeline & Merchants", "#10B981", "⟳",
             "Shared merchant pipeline and live performance data — TPV, AOV, approval rates.",
             [("Partner Leads","Pipeline"),("Merchants","Merchants")]),
            ("Performance", "#F59E0B", "▲",
             "Monitor TPV, auth rates, and transaction volumes for your Yuno integration.",
             [("Partner Health","Performance"),("Benchmarks","Benchmarks")]),
            ("Market Intel", "#0EA5E9", "◎",
             "Coverage analysis, payment trends, and strategic context for your region.",
             [("Market Intel","Insights")]),
        ]

    cols = st.columns(2)
    for ci, (title, color, icon, desc, links) in enumerate(cards):
        with cols[ci % 2]:
            link_pills = " ".join(f'<span style="font-size:10px;font-weight:600;padding:3px 10px;border-radius:6px;background:{color}10;color:{color};border:0.5px solid {color}25;cursor:pointer;">{lbl}</span>' for lbl, _ in links)
            st.markdown(f'<div style="background:#fff;border:1px solid #E2E8F0;border-radius:14px;padding:22px;margin-bottom:10px;border-top:3px solid {color};min-height:180px;display:flex;flex-direction:column;">'
                        f'<div style="display:flex;align-items:center;gap:10px;margin-bottom:10px;">'
                        f'<span style="font-size:20px;">{icon}</span>'
                        f'<span style="font-size:15px;font-weight:700;color:#0f172a;">{title}</span></div>'
                        f'<div style="font-size:12px;color:#64748b;line-height:1.6;flex:1;">{desc}</div>'
                        f'<div style="display:flex;flex-wrap:wrap;gap:5px;margin-top:14px;">{link_pills}</div>'
                        f'</div>', unsafe_allow_html=True)



# ── Main App ───────────────────────────────────────────────────────────────────
inject_css(st.session_state.role)

if st.session_state.role is None:
    show_landing()
else:
    show_sidebar()
    # Top-right sign out button
    _top_left, _top_right = st.columns([9, 1])
    with _top_right:
        if st.button("Sign out", key="goto_landing", use_container_width=True):
            st.session_state.role = None
            st.session_state.page = "Home"
            st.query_params.clear()
            st.rerun()
    # Style the sign-out button
    st.markdown("""<style>
    button[data-testid="baseButton-secondary"][kind="secondary"]:has(~ div) {},
    [data-testid="stMain"] [data-testid="stColumns"]:first-of-type [data-testid="stColumn"]:last-child button {
        background: transparent !important;
        border: 1px solid #E2E8F0 !important;
        border-radius: 8px !important;
        color: #64748B !important;
        font-size: 12px !important;
        font-weight: 600 !important;
        padding: 6px 14px !important;
        transition: all 0.15s !important;
        letter-spacing: 0.01em !important;
    }
    [data-testid="stMain"] [data-testid="stColumns"]:first-of-type [data-testid="stColumn"]:last-child button:hover {
        background: #F8FAFC !important;
        border-color: #CBD5E1 !important;
        color: #334155 !important;
    }
    </style>""", unsafe_allow_html=True)
    page = st.session_state.page
    if page == "Home":
        show_home()
    elif page == "Pipeline":
        show_pipeline()
    elif page == "Partners":
        show_partners()
    elif page == "Merchants":
        show_merchants()
    elif page == "Contacts":
        st.session_state.page = "Partners"
        st.rerun()
    elif page == "Performance":
        show_performance()
    elif page == "Benchmarks":
        show_benchmarks()
    elif page == "Insights":
        show_insights()
    elif page == "PayRec":
        show_payment_recs()
    elif page == "MerchSim":
        show_merchant_sim()
    elif page == "MissionCtrl":
        show_mission_control()
